from fastapi import FastAPI, HTTPException, UploadFile, status, File, Depends, Form, Header, Query, Response
from pymongo import MongoClient
from pymongo.errors import PyMongoError
from bson import ObjectId
import boto3
from typing import Any, List
from Models.model import CSVNodeQueryRequest, CSVNodeQueryResponse, Flow
from Models.model import PDFNodeQueryRequest
from Models.model import PDFNodeQueryResponse
from Models.model import TXTNodeQueryRequest
from Models.model import TXTNodeQueryResponse
from Models.model import MDNodeQueryRequest
from Models.model import MDNodeQueryResponse
from Models.model import HTMLNodeQueryRequest
from Models.model import HTMLNodeQueryResponse
from Models.model import DOCXNodeQueryRequest
from Models.model import DOCXNodeQueryResponse
from Models.model import PPTXNodeQueryRequest
from Models.model import PPTXNodeQueryResponse
from Models.model import ImgNodeQueryRequest
from Models.model import ImgNodeQueryResponse
from Models.model import AudioNodeQueryRequest
from Models.model import AudioNodeQueryResponse
from Models.model import YoutubeNodeQueryRequest
from Models.model import YoutubeNodeQueryResponse
from Models.model import VideoNodeQueryRequest
from Models.model import VideoNodeQueryResponse
from Models.model import WebNodeQueryRequest
from Models.model import WebNodeQueryResponse
from Models.model import SQLComponentRequest
from Models.model import SQLComponentResponse
from Models.model import SQLNodeQueryRequest
from Models.model import SQLNodeQueryResponse
from Models.model import ComponentFollowUpQueryRequest
from Models.model import ComponentFollowUpQueryResponse
from Models.model import MultipleQuestionAnswerQueryRequest
from Models.model import MultipleQuestionAnswerQueryResponse
from Models.model import FlowSummarizeRequest
from Models.model import FlowSummarizeResponse
from pymongo.mongo_client import MongoClient
from botocore.exceptions import ClientError
from hashlib import sha256
from io import BytesIO
from trp import Document
import time
import pandas as pd
from langchain_experimental.text_splitter import SemanticChunker
from langchain_openai.embeddings import OpenAIEmbeddings
from langchain_chroma import Chroma
from langchain_community.document_loaders import UnstructuredMarkdownLoader
import chromadb
from uuid import uuid4
from langchain_core.documents import Document as LangDocument
from langchain_openai import ChatOpenAI
from langchain_core.prompts import PromptTemplate
import datetime
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel, Field
import pandas as pd
import sqlite3
from vanna.openai import OpenAI_Chat
from vanna.chromadb import ChromaDB_VectorStore
from typing import List
from copy import deepcopy
from crawl4ai import *
import os
import base64
import openai
import pypdfium2 as pdfium
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain.chains.combine_documents.stuff import StuffDocumentsChain
from langchain.chains.combine_documents.reduce import ReduceDocumentsChain
from langchain.chains import MapReduceDocumentsChain
import json
from unstructured.partition.pdf import partition_pdf
import camelot
import re
from unstructured.documents.elements import (
    Text,
    Title,
    NarrativeText,
    ListItem,
    Header as UnstructuredHeader,
)
import google.generativeai as genai
import vertexai
from vertexai.generative_models import GenerativeModel, Part
from google.oauth2 import service_account
import io
import markdown
from bs4 import BeautifulSoup
from docx import Document
from pathlib import Path
from pptx import Presentation
import tiktoken
import traceback
from dotenv import load_dotenv
from config import (
    MissingConfigurationError,
    configuration_http_error,
    get_setting,
    require_settings,
    reset_request_settings,
    set_request_settings,
)
from code_intelligence import (
    CodeIntelligenceCapabilityError,
    GitHubRepoScanError,
    build_code_intelligence_artifacts,
    code_intelligence_capability_contract,
    code_intelligence_to_markdown,
    require_code_intelligence_enabled,
    resolve_allowed_local_repo_root,
    scan_github_repo,
    scan_local_repo,
)
from integrations.github import GitHubClient, GitHubClientError
from export.csv_tasks import export_task_rows
from export.markdown import (
    export_executive_output_markdown,
    export_executive_summary_markdown,
)
from export.workspace_graph import (
    artifact_export_data,
    artifact_to_news_article_markdown,
    build_workspace_graph,
    graph_to_completeness_markdown,
    graph_to_completeness_review,
    graph_to_executive_markdown,
    graph_to_markdown,
    graph_to_mermaid,
    graph_to_mmd_json,
    graph_to_news_article_markdown,
    graph_to_opml,
    graph_to_task_rows,
    graph_to_team_roadmap,
    graph_to_team_roadmap_markdown,
    select_latest_ai_draft_artifact,
    select_branch,
)
from integrations.miro.client import MiroClient
from integrations.miro.exporter import (
    export_branch_to_miro_payload,
    export_sme_review_board_payload,
)
from integrations.miro.native_mindmap import export_native_mindmap_payload
from integrations.miro.persistence import (
    apply_miro_external_refs_to_flow_json,
    miro_item_refs_from_result,
)
from integrations.monday.client import MondayClient
from integrations.monday.exporter import (
    export_tasks_to_monday_payload,
    select_monday_task_nodes,
)
from integrations.monday.persistence import (
    apply_monday_external_refs_to_flow_json,
    apply_monday_status_projection_to_flow_json,
    monday_item_refs_from_result,
    monday_refs_from_flow_json,
    monday_status_projections_from_result,
)
from documents.ingestion import (
    ALLOWED_DOCUMENT_EXTENSIONS,
    build_ai_intake_source_document,
    build_source_set_metadata,
    chunk_source_segments,
    DocumentIngestionError,
    file_sha256,
    ingest_supported_document,
    normalize_relative_source_path,
    sanitize_filename,
    source_document_with_source_set_metadata,
    source_document_from_upload,
    validate_ai_intake_bytes,
    validate_upload_bytes,
)
from documents.source_refs import attach_source_refs_to_mindmap
from structured_data import build_structured_data_artifacts
from ai.roles import (
    UnknownSourceIntakeRole,
    build_source_intake_instruction as build_role_source_intake_instruction,
    clean_source_intake_value as clean_role_source_intake_value,
    resolve_source_intake_role_label,
)
from ai_model_policy import normalize_model_name
from ai_helpers import (
    accept_ai_draft_revision,
    add_source_to_ai_draft_session,
    append_ai_draft_revision,
    build_ai_draft_revision,
    build_ai_draft_session,
    build_ai_action_run,
    discard_ai_draft_session,
    generate_ai_action_preview,
    generate_ai_draft_session_with_provider,
    generate_helper_preview,
    generate_node_info_message_with_provider,
    generate_source_reconciliation_preview,
    generate_source_librarian_preview,
    normalize_ai_draft_scope,
    revise_ai_draft_session_with_provider,
    validate_ai_draft_session,
    list_prompt_profiles,
)
from graph.ai_contract import (
    append_ai_graph_prompt_contract,
    parse_ai_mindmap_response,
    validate_ai_mindmap_contract,
)
from graph.schemas import GraphSchemaError, validate_workspace_brief
from openai_sources import (
    generate_component_answer,
    generate_component_follow_up_questions,
    generate_document_mindmap,
    generate_document_summary,
    generate_audio_mindmap,
    generate_image_mindmap,
    generate_video_mindmap,
    generate_web_mindmap,
    transcribe_audio,
)

load_dotenv()
os.environ.setdefault("ANONYMIZED_TELEMETRY", "False")

mongo_db_url = os.getenv("mongo_db_url")
openai_api_key_str = os.getenv("openai_api_key")
gemini_api_key_str = os.getenv("gemini_api_key")
gcp_project_id_str = os.getenv("gcp_project_id")
aws_access_key_id_str = os.getenv("aws_access_key_id")
aws_secret_access_key_str = os.getenv("aws_secret_access_key")
bucket_name = os.getenv("bucket_name")
gcp_service_account_file = os.getenv(
    "gcp_service_account_file", "./ai-interview-poc-2b5cf8540f16.json"
)
OPENAI_DEFAULT_MODEL = os.getenv("openai_default_model", "gpt-5.5")
OPENAI_REASONING_MODEL = os.getenv("openai_reasoning_model", "gpt-5.4")
OPENAI_EMBEDDING_MODEL = os.getenv(
    "openai_embedding_model", "text-embedding-3-large"
)
LEGACY_ASSISTANTS_FALLBACK_ENV = "DOCMAP_ALLOW_LEGACY_ASSISTANTS"


def clean_source_intake_value(value: str | None, max_length: int = 2000) -> str:
    return clean_role_source_intake_value(value, max_length)


def resolve_assistants_model(model_name: str | None = None) -> str:
    requested_model = clean_source_intake_value(model_name, 120)
    if not requested_model:
        return OPENAI_DEFAULT_MODEL
    try:
        return normalize_model_name(requested_model)
    except ValueError as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=str(exc),
        ) from exc


def legacy_assistants_fallback_enabled() -> bool:
    value = os.getenv(LEGACY_ASSISTANTS_FALLBACK_ENV, "").strip().lower()
    return value in {"1", "true", "yes", "on"}


def require_legacy_assistants_fallback(source_type: str, *, purpose: str) -> None:
    if legacy_assistants_fallback_enabled():
        return
    raise HTTPException(
        status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
        detail=(
            f"Responses-based {purpose} could not run for {source_type} because prepared document sections "
            f"were unavailable. Set {LEGACY_ASSISTANTS_FALLBACK_ENV}=true to allow the temporary "
            "Assistants file-search fallback."
        ),
    )


def resolve_source_intake_role(intake_role: str | None = None) -> str:
    try:
        return resolve_source_intake_role_label(intake_role)
    except UnknownSourceIntakeRole as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Unknown DOCX intake role.",
        ) from exc


def build_source_intake_instruction(
    intake_role: str | None = None,
    intake_prompt: str | None = None,
) -> str:
    try:
        return build_role_source_intake_instruction(intake_role, intake_prompt)
    except UnknownSourceIntakeRole as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Unknown DOCX intake role.",
        ) from exc


def repair_flow_snapshot_for_persistence(
    flow_json: str,
    *,
    flow_id: str = "",
    flow_name: str = "Untitled Workspace",
    flow_type: str = "",
    summary: str = "",
) -> str:
    try:
        snapshot = json.loads(flow_json or "{}")
    except json.JSONDecodeError:
        return flow_json
    if not isinstance(snapshot, dict) or not isinstance(snapshot.get("nodes"), list):
        return flow_json
    workspace_brief = snapshot.get("workspace_brief")
    if workspace_brief is not None:
        validate_workspace_brief(workspace_brief)

    graph = build_workspace_graph(
        {
            "_id": flow_id,
            "flow_name": flow_name,
            "summary": summary,
            "flow_type": flow_type,
            "flow_json": json.dumps(snapshot),
        }
    )
    graph_nodes_by_id = {node.get("id"): node for node in graph.get("nodes", [])}

    changed = False
    next_nodes = []
    for raw_node in snapshot.get("nodes", []):
        graph_node = graph_nodes_by_id.get(raw_node.get("id"))
        if not graph_node:
            next_nodes.append(raw_node)
            continue

        next_node = _apply_graph_review_state_to_react_node(raw_node, graph_node)
        changed = changed or next_node != raw_node
        next_nodes.append(next_node)

    if not changed:
        return flow_json

    snapshot["nodes"] = next_nodes
    return json.dumps(snapshot)


def _apply_graph_review_state_to_react_node(raw_node: dict, graph_node: dict) -> dict:
    if not isinstance(raw_node, dict):
        return raw_node

    status_value = graph_node.get("status")
    source_refs = graph_node.get("source_refs")
    if not status_value and not source_refs:
        return raw_node

    next_node = deepcopy(raw_node)
    data = next_node.setdefault("data", {})
    if not isinstance(data, dict):
        next_node["data"] = {}
        data = next_node["data"]

    nested_data = data.get("data")
    if not isinstance(nested_data, dict):
        nested_data = None

    if status_value:
        data["status"] = status_value
        if nested_data is not None:
            nested_data["status"] = status_value

    if isinstance(source_refs, list) and source_refs:
        data["source_refs"] = source_refs
        if nested_data is not None:
            nested_data["source_refs"] = source_refs

    return next_node

credentials = None
model_vertexai = None
if gcp_project_id_str and Path(gcp_service_account_file).exists():
    credentials = service_account.Credentials.from_service_account_file(
        gcp_service_account_file
    )
    vertexai.init(
        project=gcp_project_id_str, credentials=credentials, location="us-central1"
    )
    model_vertexai = GenerativeModel("gemini-2.0-flash")

genai.configure(api_key=gemini_api_key_str)

model = genai.GenerativeModel("gemini-2.0-flash")

connection = sqlite3.connect("sqlite_data.db")

GPT_4O_MAX_TOKENS = 128000

UPLOAD_DIR = "uploaded_pdfs"

openai.api_key = openai_api_key_str

operation_progress: dict[str, dict] = {}


def update_operation_progress(
    operation_id: str | None,
    *,
    phase: str,
    message: str,
    progress: int,
    detail: str = "",
    status_value: str = "running",
) -> None:
    if not operation_id:
        return

    operation_progress[operation_id] = {
        "operation_id": operation_id,
        "phase": phase,
        "message": message,
        "detail": detail,
        "progress": max(0, min(100, progress)),
        "status": status_value,
        "updated_at": utc_timestamp(),
    }

class SQLBot(ChromaDB_VectorStore, OpenAI_Chat):
    def __init__(self, config=None):
        ChromaDB_VectorStore.__init__(self, config=config)
        OpenAI_Chat.__init__(self, config=config)


class CSVBot(ChromaDB_VectorStore, OpenAI_Chat):
    def __init__(self, config=None):
        ChromaDB_VectorStore.__init__(self, config=config)
        OpenAI_Chat.__init__(self, config=config)


class LazyVannaBot:
    def __init__(self, bot_class, path: str, sqlite_path: str, train_sql_schema: bool = False):
        self.bot_class = bot_class
        self.path = path
        self.sqlite_path = sqlite_path
        self.train_sql_schema = train_sql_schema
        self._bot = None
        self._api_key = None

    def _get_bot(self):
        api_key = get_setting("openai_api_key")
        if not api_key:
            raise MissingConfigurationError(
                "Missing required environment variable(s): openai_api_key."
            )

        if self._bot is not None and self._api_key == api_key:
            return self._bot

        bot = self.bot_class(
            config={
                "api_key": api_key,
                "model": OPENAI_DEFAULT_MODEL,
                "temperature": 0,
                "path": self.path,
                "client": "persistent",
                "n_results": 1,
            }
        )
        bot.connect_to_sqlite(self.sqlite_path)

        if self.train_sql_schema:
            df_ddl = bot.run_sql("SELECT type, sql FROM sqlite_master WHERE sql IS NOT NULL")
            for ddl in df_ddl["sql"].to_list():
                bot.train(ddl=ddl)

        self._bot = bot
        self._api_key = api_key
        return bot

    def __getattr__(self, name):
        return getattr(self._get_bot(), name)


sqlBot = LazyVannaBot(SQLBot, "./SQLVectorStore", "sqlite_data.db", train_sql_schema=True)
csvBot = LazyVannaBot(CSVBot, "./CSVVectorStore", "csv_data.db")

app = FastAPI()


class LocalRepoScanRequest(BaseModel):
    root: str
    repo_label: str = ""
    changed_paths: list[str] = Field(default_factory=list)
    max_files: int = Field(default=500, ge=1, le=5000)
    max_file_bytes: int = Field(default=256_000, ge=1, le=2_000_000)
    large_file_line_threshold: int = Field(default=500, ge=50, le=10_000)


class GitHubRepoScanRequest(BaseModel):
    owner: str
    repo: str
    ref: str = "main"
    path: str = ""
    repo_label: str = ""
    changed_paths: list[str] = Field(default_factory=list)
    max_files: int = Field(default=200, ge=1, le=1000)
    max_file_bytes: int = Field(default=256_000, ge=1, le=2_000_000)
    large_file_line_threshold: int = Field(default=500, ge=50, le=10_000)

origins = [
    "http://127.0.0.1:5173",
    "http://localhost:5173",
    "http://127.0.0.1:4173",
    "http://localhost:4173",
]


def _cors_error_headers_for_origin(origin: str | None) -> dict[str, str]:
    if origin not in origins:
        return {}
    return {
        "Access-Control-Allow-Origin": origin,
        "Access-Control-Allow-Credentials": "true",
        "Vary": "Origin",
    }

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.get("/operations/{operation_id}")
def get_operation_progress(operation_id: str):
    progress = operation_progress.get(operation_id)
    if not progress:
        return {
            "operation_id": operation_id,
            "phase": "queued",
            "message": "Waiting for backend to begin.",
            "detail": "",
            "progress": 5,
            "status": "queued",
            "updated_at": utc_timestamp(),
        }
    return progress


@app.middleware("http")
async def apply_local_user_settings(request, call_next):
    request_settings = {
        "openai_api_key": request.headers.get("x-docmap-openai-api-key") or "",
        "miro_api_token": request.headers.get("x-docmap-miro-api-token") or "",
        "monday_api_token": request.headers.get("x-docmap-monday-api-token") or "",
    }
    token = set_request_settings(
        {key: value for key, value in request_settings.items() if value}
    )
    previous_openai_key = openai.api_key
    if request_settings["openai_api_key"]:
        openai.api_key = request_settings["openai_api_key"]
    try:
        return await call_next(request)
    except Exception as exc:
        traceback.print_exc()
        return JSONResponse(
            status_code=500,
            content={
                "detail": {
                    "message": "Backend request failed unexpectedly.",
                    "error_type": exc.__class__.__name__,
                    "error": str(exc),
                }
            },
            headers=_cors_error_headers_for_origin(request.headers.get("origin")),
        )
    finally:
        openai.api_key = previous_openai_key
        reset_request_settings(token)


client = MongoClient(mongo_db_url, serverSelectionTimeoutMS=2000)

try:
    client.admin.command("ping")
    print("Pinged your deployment. You successfully connected to MongoDB!")
except Exception as e:
    print(e)

db = client["MindMap"]
flow_collection = db["flows"]
component_collection = db["components"]
node_collection = db["nodes"]
ai_draft_session_collection = db["ai_draft_sessions"]
LOCAL_FLOW_STORE_PATH = Path(
    os.getenv("DOCMAP_LOCAL_FLOW_STORE", "docmap_flows.json")
)


@app.get("/api/capabilities")
def get_capabilities():
    return code_intelligence_capability_contract()


def _github_http_exception(exc: GitHubClientError) -> HTTPException:
    status_code = exc.status_code or status.HTTP_502_BAD_GATEWAY
    if status_code not in {
        status.HTTP_401_UNAUTHORIZED,
        status.HTTP_403_FORBIDDEN,
        status.HTTP_404_NOT_FOUND,
        status.HTTP_429_TOO_MANY_REQUESTS,
    }:
        status_code = status.HTTP_502_BAD_GATEWAY
    headers = {"Retry-After": exc.retry_after} if exc.retry_after and status_code == status.HTTP_429_TOO_MANY_REQUESTS else None
    return HTTPException(
        status_code=status_code,
        detail={
            "message": str(exc),
            "reason_code": exc.reason_code,
        },
        headers=headers,
    )


@app.post("/api/code-intelligence/local-repo/scan")
def scan_local_repo_endpoint(request: LocalRepoScanRequest):
    try:
        repo_root = resolve_allowed_local_repo_root(request.root)
    except CodeIntelligenceCapabilityError as exc:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail=str(exc)) from exc

    try:
        return scan_local_repo(
            repo_root,
            repo_label=request.repo_label,
            max_files=request.max_files,
            max_file_bytes=request.max_file_bytes,
            large_file_line_threshold=request.large_file_line_threshold,
        )
    except ValueError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc


@app.post("/api/code-intelligence/local-repo/report.md")
def scan_local_repo_report_endpoint(request: LocalRepoScanRequest):
    graph = scan_local_repo_endpoint(request)
    return Response(
        content=code_intelligence_to_markdown(graph),
        media_type="text/markdown",
    )


@app.post("/api/code-intelligence/local-repo/artifacts")
def scan_local_repo_artifacts_endpoint(request: LocalRepoScanRequest):
    graph = scan_local_repo_endpoint(request)
    return build_code_intelligence_artifacts(graph, changed_paths=request.changed_paths)


@app.post("/api/code-intelligence/github/scan")
def scan_github_repo_endpoint(
    request: GitHubRepoScanRequest,
    x_docmap_github_token: str = Header(default=""),
):
    try:
        require_code_intelligence_enabled()
    except CodeIntelligenceCapabilityError as exc:
        raise HTTPException(status_code=status.HTTP_403_FORBIDDEN, detail=str(exc)) from exc

    if not x_docmap_github_token.strip():
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail="GitHub token is required.")

    repo_name = f"{request.owner}/{request.repo}"
    try:
        return scan_github_repo(
            GitHubClient(x_docmap_github_token),
            repo=repo_name,
            ref=request.ref,
            path_prefix=request.path,
            max_files=request.max_files,
            max_file_bytes=request.max_file_bytes,
            large_file_line_threshold=request.large_file_line_threshold,
        )
    except GitHubRepoScanError as exc:
        raise HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(exc)) from exc
    except GitHubClientError as exc:
        raise _github_http_exception(exc) from exc


@app.post("/api/code-intelligence/github/report.md")
def scan_github_repo_report_endpoint(
    request: GitHubRepoScanRequest,
    x_docmap_github_token: str = Header(default=""),
):
    graph = scan_github_repo_endpoint(request, x_docmap_github_token)
    return Response(
        content=code_intelligence_to_markdown(graph),
        media_type="text/markdown",
    )


@app.post("/api/code-intelligence/github/artifacts")
def scan_github_repo_artifacts_endpoint(
    request: GitHubRepoScanRequest,
    x_docmap_github_token: str = Header(default=""),
):
    graph = scan_github_repo_endpoint(request, x_docmap_github_token)
    return build_code_intelligence_artifacts(graph, changed_paths=request.changed_paths)
LOCAL_AI_DRAFT_SESSION_STORE_PATH = Path(
    os.getenv("DOCMAP_LOCAL_AI_DRAFT_SESSION_STORE", "docmap_ai_draft_sessions.json")
)


def local_flow_store_path() -> Path:
    return Path(__file__).resolve().parent / LOCAL_FLOW_STORE_PATH


def local_ai_draft_session_store_path() -> Path:
    return Path(__file__).resolve().parent / LOCAL_AI_DRAFT_SESSION_STORE_PATH


def load_local_flows() -> list[dict]:
    path = local_flow_store_path()
    if not path.exists():
        return []
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return []


def load_local_ai_draft_sessions() -> list[dict]:
    path = local_ai_draft_session_store_path()
    if not path.exists():
        return []
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError):
        return []


def save_local_flows(flows: list[dict]) -> None:
    path = local_flow_store_path()
    path.write_text(json.dumps(flows, indent=2, default=str), encoding="utf-8")


def save_local_ai_draft_sessions(sessions: list[dict]) -> None:
    path = local_ai_draft_session_store_path()
    path.write_text(json.dumps(sessions, indent=2, default=str), encoding="utf-8")


def normalize_flow_record(flow: dict) -> dict:
    return {
        "flow_id": str(flow.get("_id") or flow.get("flow_id")),
        "flow_name": flow.get("flow_name") or "Untitled workspace",
        "flow_json": flow.get("flow_json") or "",
        "summary": flow.get("summary") or "",
        "flow_type": flow.get("flow_type") or "manual",
    }


def local_create_flow(flow_data: dict) -> dict:
    flows = load_local_flows()
    flow = {
        "_id": str(ObjectId()),
        **flow_data,
    }
    flows.append(flow)
    save_local_flows(flows)
    return flow


def local_find_flow(flow_id: str) -> dict | None:
    for flow in load_local_flows():
        if str(flow.get("_id") or flow.get("flow_id")) == flow_id:
            return flow
    return None


def local_update_flow(flow_id: str, updates: dict) -> bool:
    flows = load_local_flows()
    updated = False
    for flow in flows:
        if str(flow.get("_id") or flow.get("flow_id")) == flow_id:
            flow.update(updates)
            updated = True
            break
    if updated:
        save_local_flows(flows)
    return updated


def local_delete_flow(flow_id: str) -> bool:
    flows = load_local_flows()
    next_flows = [
        flow for flow in flows
        if str(flow.get("_id") or flow.get("flow_id")) != flow_id
    ]
    if len(next_flows) == len(flows):
        return False
    save_local_flows(next_flows)
    return True


def local_save_ai_draft_session(session: dict) -> dict:
    sessions = load_local_ai_draft_sessions()
    saved = False
    for index, existing in enumerate(sessions):
        if existing.get("session_id") == session.get("session_id"):
            sessions[index] = session
            saved = True
            break
    if not saved:
        sessions.append(session)
    save_local_ai_draft_sessions(sessions)
    return session


def local_find_ai_draft_session(flow_id: str, session_id: str) -> dict | None:
    for session in load_local_ai_draft_sessions():
        if (
            str(session.get("workspace_id")) == flow_id
            and session.get("session_id") == session_id
        ):
            return session
    return None


def promote_local_flow_to_mongo(flow_id: str) -> dict | None:
    if not ObjectId.is_valid(flow_id):
        return None

    local_flow = local_find_flow(flow_id)
    if not local_flow:
        return None

    flow_data = {
        "_id": ObjectId(flow_id),
        "flow_name": local_flow.get("flow_name") or "Untitled workspace",
        "flow_json": local_flow.get("flow_json") or "",
        "summary": local_flow.get("summary") or "",
        "flow_type": local_flow.get("flow_type") or "manual",
    }

    try:
        flow_collection.insert_one(flow_data)
    except PyMongoError:
        try:
            return flow_collection.find_one({"_id": ObjectId(flow_id)})
        except PyMongoError:
            return None

    return flow_data

# AWS S3 setup
s3_client = boto3.client(
    "s3",
    aws_access_key_id=aws_access_key_id_str,
    aws_secret_access_key=aws_secret_access_key_str,
    region_name="ap-south-1",
)


persistent_client = chromadb.PersistentClient()


class LazyOpenAIEmbeddings:
    def __init__(self):
        self._client = None
        self._api_key = None

    def _get_client(self):
        api_key = get_setting("openai_api_key")
        if not api_key:
            raise MissingConfigurationError(
                "Missing required environment variable(s): openai_api_key."
            )
        if self._client is None or self._api_key != api_key:
            self._client = OpenAIEmbeddings(
                model=OPENAI_EMBEDDING_MODEL, api_key=api_key
            )
            self._api_key = api_key
        return self._client

    def __getattr__(self, name):
        return getattr(self._get_client(), name)


class LazyChromaStore:
    def __init__(self):
        self._store = None
        self._api_key = None

    def _get_store(self):
        api_key = get_setting("openai_api_key")
        if not api_key:
            raise MissingConfigurationError(
                "Missing required environment variable(s): openai_api_key."
            )
        if self._store is None or self._api_key != api_key:
            self._store = Chroma(
                client=persistent_client,
                collection_name="pdfs",
                embedding_function=OpenAIEmbeddings(
                    model=OPENAI_EMBEDDING_MODEL, api_key=api_key
                ),
            )
            self._api_key = api_key
        return self._store

    def __getattr__(self, name):
        return getattr(self._get_store(), name)


class LazySemanticChunker:
    def __init__(self):
        self._chunker = None
        self._api_key = None

    def _get_chunker(self):
        api_key = get_setting("openai_api_key")
        if not api_key:
            raise MissingConfigurationError(
                "Missing required environment variable(s): openai_api_key."
            )
        if self._chunker is None or self._api_key != api_key:
            self._chunker = SemanticChunker(
                OpenAIEmbeddings(model=OPENAI_EMBEDDING_MODEL, api_key=api_key)
            )
            self._api_key = api_key
        return self._chunker

    def __getattr__(self, name):
        return getattr(self._get_chunker(), name)


class LazyChatOpenAI:
    def __init__(self):
        self._llm = None
        self._api_key = None

    def _get_llm(self):
        api_key = get_setting("openai_api_key")
        if not api_key:
            raise MissingConfigurationError(
                "Missing required environment variable(s): openai_api_key."
            )
        if self._llm is None or self._api_key != api_key:
            self._llm = ChatOpenAI(model=OPENAI_REASONING_MODEL, api_key=api_key)
            self._api_key = api_key
        return self._llm

    def __getattr__(self, name):
        return getattr(self._get_llm(), name)

    def __ror__(self, other):
        return other | self._get_llm()


embedding_function = LazyOpenAIEmbeddings()
vector_store_from_client = LazyChromaStore()
PDFCollection = persistent_client.get_or_create_collection("pdfs")
text_splitter = LazySemanticChunker()
llm = LazyChatOpenAI()


def read_upload_bytes(file: UploadFile) -> bytes:
    file.file.seek(0)
    file_bytes = file.file.read()
    file.file.seek(0)
    return file_bytes


def prepare_source_upload(file: UploadFile, flow_id: str, expected_extension: str | None = None) -> dict:
    file_bytes = read_upload_bytes(file)
    upload = validate_upload_bytes(file.filename, file_bytes)

    if expected_extension and upload["extension"] != expected_extension:
        raise DocumentIngestionError(f"Only {expected_extension.upper()} files are allowed.")

    try:
        existing_component = component_collection.find_one(
            {"file_hash": upload["file_hash"], "flow_id": ObjectId(flow_id)}
        )
        existing_versions = component_collection.count_documents(
            {
                "flow_id": ObjectId(flow_id),
                "source_document.filename": upload["filename"],
            }
        )
    except PyMongoError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=(
                "Source uploads need MongoDB so document metadata and source references can be saved. "
                "Start MongoDB, then reopen or create a workspace and try the DOCX upload again."
            ),
        ) from exc

    if existing_component:
        source_document = dict(existing_component.get("source_document") or {})
        if not source_document:
            source_document = source_document_from_upload(upload, version=existing_component.get("version", 1) or 1)
        return {
            "upload": upload,
            "file_bytes": file_bytes,
            "source_document": source_document,
            "source_segments": existing_component.get("source_segments", []),
            "document_chunks": existing_component.get("document_chunks", []),
            "existing_component": existing_component,
            "reused_existing_source": True,
        }

    return ingest_supported_document(
        upload["filename"],
        file_bytes,
        version=existing_versions + 1,
    )


def prepare_ai_intake_upload(file: UploadFile, flow_id: str) -> dict:
    file_bytes = read_upload_bytes(file)
    upload = validate_ai_intake_bytes(file.filename, file_bytes)

    if upload["extension"] in ALLOWED_DOCUMENT_EXTENSIONS:
        return prepare_source_upload(file, flow_id)

    try:
        existing_component = component_collection.find_one(
            {"file_hash": upload["file_hash"], "flow_id": ObjectId(flow_id)}
        )
        existing_versions = component_collection.count_documents(
            {
                "flow_id": ObjectId(flow_id),
                "source_document.filename": upload["filename"],
            }
        )
    except PyMongoError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=(
                "Source uploads need MongoDB so document metadata and source references can be saved. "
                "Start MongoDB, then reopen or create a workspace and try the upload again."
            ),
        ) from exc

    if existing_component:
        source_document = dict(existing_component.get("source_document") or {})
        if not source_document:
            source_document = build_ai_intake_source_document(
                upload["filename"],
                file_bytes,
                version=existing_component.get("version", 1) or 1,
            )
        return {
            "upload": upload,
            "file_bytes": file_bytes,
            "source_document": source_document,
            "source_segments": existing_component.get("source_segments", []),
            "document_chunks": existing_component.get("document_chunks", []),
            "existing_component": existing_component,
            "reused_existing_source": True,
        }

    source_document = build_ai_intake_source_document(
        upload["filename"],
        file_bytes,
        version=existing_versions + 1,
    )
    return {
        "upload": upload,
        "file_bytes": file_bytes,
        "source_document": source_document,
        "source_segments": [],
        "document_chunks": [],
    }


def binary_source_context(
    *,
    filename: str,
    file_bytes: bytes,
    source_type: str,
    flow_id: str,
) -> dict:
    sanitized_filename = sanitize_filename(filename)
    file_hash = file_sha256(file_bytes)
    existing_versions = component_collection.count_documents(
        {
            "flow_id": ObjectId(flow_id),
            "source_document.filename": sanitized_filename,
        }
    )
    upload = {
        "filename": sanitized_filename,
        "original_filename": filename or sanitized_filename,
        "extension": source_type,
        "size": len(file_bytes),
        "file_hash": file_hash,
    }
    source_document = source_document_from_upload(upload, version=existing_versions + 1)
    return {
        "upload": upload,
        "file_bytes": file_bytes,
        "source_document": source_document,
        "source_segments": [],
        "document_chunks": [],
    }


def virtual_source_context(*, label: str, source_type: str, flow_id: str) -> dict:
    file_bytes = label.encode("utf-8")
    return binary_source_context(
        filename=f"{source_type}-{file_sha256(file_bytes)[:12]}.{source_type}.txt",
        file_bytes=file_bytes,
        source_type=source_type,
        flow_id=flow_id,
    )


def source_metadata_fields(source_context: dict) -> dict:
    source_document = source_context["source_document"]
    return {
        "name": source_document["filename"],
        "original_name": source_document["original_filename"],
        "file_hash": source_document["file_hash"],
        "source_document_id": source_document["id"],
        "source_document": source_document,
        "source_segments": source_context.get("source_segments", []),
        "document_chunks": source_context["document_chunks"],
    }


def source_ref_from_chunk(source_document: dict, chunk: dict | None = None) -> dict:
    chunk = chunk if isinstance(chunk, dict) else {}
    text = str(chunk.get("text") or "")
    return {
        "document_id": source_document.get("id", ""),
        "chunk_id": chunk.get("id", ""),
        "page": chunk.get("page"),
        "section": chunk.get("heading") or "",
        "quote_snippet": text[:280],
        "source_type": source_document.get("type", ""),
    }


def source_library_record(source_context: dict, component_id: Any | None = None) -> dict:
    source_document = source_context["source_document"]
    return {
        "id": source_document.get("id", ""),
        "document_id": source_document.get("id", ""),
        "title": source_document.get("original_filename")
        or source_document.get("filename")
        or "Uploaded source",
        "filename": source_document.get("filename", ""),
        "original_filename": source_document.get("original_filename", ""),
        "type": source_document.get("type", "source"),
        "type_label": str(source_document.get("type", "source")).upper(),
        "status": "parsed",
        "component_id": str(component_id or ""),
        "flow_id": "",
        "file_hash": source_document.get("file_hash", ""),
        "size": source_document.get("size", 0),
        "version": source_document.get("version", 1),
        "chunks": source_context.get("document_chunks", []),
        "segments": source_context.get("source_segments", []),
        "normalized_document_id": source_document.get("id", ""),
        "metadata": source_document,
    }


def fallback_source_summary(source_context: dict) -> str:
    source_document = source_context["source_document"]
    chunks = source_context.get("document_chunks", [])
    chunk_count = len(chunks)
    page_values = sorted(
        {
            chunk.get("page")
            for chunk in chunks
            if isinstance(chunk, dict) and chunk.get("page") is not None
        }
    )
    page_text = (
        f" across {len(page_values)} page{'s' if len(page_values) != 1 else ''}"
        if page_values
        else ""
    )
    headings = [
        str(chunk.get("heading") or "").strip()
        for chunk in chunks
        if isinstance(chunk, dict) and str(chunk.get("heading") or "").strip()
    ][:4]
    heading_text = f" Key sections: {', '.join(headings)}." if headings else ""
    return (
        f"{source_document.get('original_filename') or source_document.get('filename') or 'The source'} "
        f"was parsed into {chunk_count} source chunk{'s' if chunk_count != 1 else ''}{page_text}. "
        "The OpenAI request timed out, so TraceSpace saved the parsed source context "
        "for follow-up questions and review instead of discarding the upload."
        f"{heading_text}"
    )


def fallback_source_mindmap(source_context: dict, flow_id: str, component_id: Any | None = None) -> dict:
    source_document = source_context["source_document"]
    chunks = [
        chunk
        for chunk in source_context.get("document_chunks", [])
        if isinstance(chunk, dict) and str(chunk.get("text") or "").strip()
    ]
    filename = source_document.get("original_filename") or source_document.get("filename") or "Uploaded source"
    id_prefix = re.sub(r"[^A-Za-z0-9_-]+", "-", source_document.get("id", "source")).strip("-") or "source"
    source_node_id = f"{id_prefix}-timeout-root"
    summary_node_id = f"{id_prefix}-timeout-summary"
    review_node_id = f"{id_prefix}-timeout-review"
    next_steps_node_id = f"{id_prefix}-timeout-next-steps"
    source_ref = source_ref_from_chunk(source_document, chunks[0] if chunks else None)
    source_refs = [source_ref] if source_ref.get("document_id") else []
    nodes = [
        {
            "id": source_node_id,
            "type": "dataSource",
            "position": {"x": 0, "y": 80},
            "data": {
                "name": source_document.get("type", "source"),
                "content": filename,
                "flow_id": flow_id,
                "component_id": str(component_id or ""),
                "source_document_id": source_document.get("id", ""),
                "source_document": source_document,
                "document_chunks": source_context.get("document_chunks", []),
                "source_segments": source_context.get("source_segments", []),
                "processing_type": "responses_timeout_fallback",
            },
            "deletable": False,
        },
        {
            "id": summary_node_id,
            "type": "response",
            "position": {"x": 420, "y": -20},
            "data": {
                "title": "Parsed Source Saved",
                "node_type": "source_summary",
                "status": "needs_review",
                "assumption": False,
                "source_refs": source_refs,
                "data": {
                    "summ": fallback_source_summary(source_context),
                    "query": "",
                    "df": [],
                    "graph": {},
                    "source_refs": source_refs,
                },
            },
            "deletable": True,
        },
        {
            "id": review_node_id,
            "type": "response",
            "position": {"x": 420, "y": 170},
            "data": {
                "title": "AI Derivation Timed Out",
                "node_type": "needs_review",
                "status": "needs_review",
                "priority": "high",
                "assumption": True,
                "source_refs": [],
                "data": {
                    "summ": (
                        "The full source-backed workspace draft did not finish within 120 seconds. "
                        "Use Ask AI on this saved source, retry derivation with a faster model, or derive a smaller branch."
                    ),
                    "query": "",
                    "df": [],
                    "graph": {},
                    "source_refs": [],
                },
            },
            "deletable": True,
        },
        {
            "id": next_steps_node_id,
            "type": "response",
            "position": {"x": 840, "y": 170},
            "data": {
                "title": "Next Step: Ask From Source Context",
                "node_type": "task",
                "status": "needs_review",
                "priority": "medium",
                "assumption": True,
                "source_refs": [],
                "data": {
                    "summ": (
                        "Start with a targeted question such as: summarize the workflow, list risks, "
                        "extract decisions, or create a concise launch-readiness map from this source."
                    ),
                    "query": "",
                    "df": [],
                    "graph": {},
                    "source_refs": [],
                },
            },
            "deletable": True,
        },
    ]

    for index, chunk in enumerate(chunks[:3], start=1):
        chunk_ref = source_ref_from_chunk(source_document, chunk)
        title = str(chunk.get("heading") or "").strip() or f"Source Chunk {index}"
        nodes.append(
            {
                "id": f"{id_prefix}-timeout-chunk-{index}",
                "type": "response",
                "position": {"x": 840, "y": -120 + index * 120},
                "data": {
                    "title": title[:90],
                    "node_type": "reference",
                    "status": "ai_generated",
                    "assumption": False,
                    "source_refs": [chunk_ref],
                    "data": {
                        "summ": str(chunk.get("text") or "").strip()[:700],
                        "query": "",
                        "df": [],
                        "graph": {},
                        "source_refs": [chunk_ref],
                    },
                },
                "deletable": True,
            }
        )

    edges = [
        {"id": f"{id_prefix}-timeout-edge-1", "source": source_node_id, "target": summary_node_id, "type": "smoothstep", "animated": True},
        {"id": f"{id_prefix}-timeout-edge-2", "source": source_node_id, "target": review_node_id, "type": "smoothstep", "animated": True},
        {"id": f"{id_prefix}-timeout-edge-3", "source": review_node_id, "target": next_steps_node_id, "type": "smoothstep", "animated": True},
        *[
            {
                "id": f"{id_prefix}-timeout-edge-chunk-{index}",
                "source": summary_node_id,
                "target": f"{id_prefix}-timeout-chunk-{index}",
                "type": "smoothstep",
                "animated": True,
            }
            for index in range(1, min(len(chunks), 3) + 1)
        ],
    ]

    return {
        "nodes": nodes,
        "edges": edges,
        "viewport": {"x": 80, "y": 120, "zoom": 0.7},
        "source_library": [source_library_record(source_context, component_id)],
        "metadata": {
            "ai_graph_contract_version": "1",
            "source_type": source_document.get("type", ""),
            "source_label": filename,
            "fallback_reason": "openai_timeout",
            "fallback": True,
        },
    }


def parse_source_set_relative_paths(raw_paths: list[str] | str | None, file_count: int) -> list[str]:
    if raw_paths is None:
        return []

    if isinstance(raw_paths, str):
        candidates = [raw_paths]
    else:
        candidates = [str(path) for path in raw_paths]

    if len(candidates) == 1:
        raw_value = candidates[0].strip()
        if raw_value.startswith("["):
            try:
                decoded = json.loads(raw_value)
            except json.JSONDecodeError as exc:
                raise DocumentIngestionError("relative_paths must be valid JSON when sent as an array string.") from exc
            if not isinstance(decoded, list):
                raise DocumentIngestionError("relative_paths JSON must be an array.")
            candidates = [str(path) for path in decoded]

    if candidates and len(candidates) != file_count:
        raise DocumentIngestionError("relative_paths must contain one path for each uploaded file.")

    return candidates


def prepare_source_set_uploads(
    files: list[UploadFile],
    *,
    flow_id: str,
    relative_paths: list[str] | str | None = None,
    source_set_id: str | None = None,
    source_set_label: str | None = None,
) -> dict:
    if not files:
        raise DocumentIngestionError("Upload at least one source document.")

    raw_paths = parse_source_set_relative_paths(relative_paths, len(files))
    prepared_uploads = []
    skipped_sources = []
    for index, file in enumerate(files):
        fallback_name = file.filename or f"source-{index + 1}"
        relative_path = normalize_relative_source_path(
            raw_paths[index] if raw_paths else "",
            fallback_filename=fallback_name,
        )
        sanitized_filename = sanitize_filename(fallback_name)
        extension = sanitized_filename.rsplit(".", 1)[-1].lower() if "." in sanitized_filename else ""
        if extension not in ALLOWED_DOCUMENT_EXTENSIONS:
            skipped_sources.append(
                source_set_skip_record(
                    filename=fallback_name,
                    relative_path=relative_path,
                    extension=extension,
                    reason_code="unsupported_extension",
                    message=(
                        "Unsupported file type for source-traceable folder review. "
                        "Supported types are PDF, DOCX, Markdown, and TXT."
                    ),
                )
            )
            continue

        file_bytes = read_upload_bytes(file)
        try:
            upload = validate_upload_bytes(file.filename, file_bytes)
        except DocumentIngestionError as exc:
            skipped_sources.append(
                source_set_skip_record(
                    filename=fallback_name,
                    relative_path=relative_path,
                    extension=extension,
                    reason_code=source_set_skip_reason(str(exc)),
                    message=str(exc),
                    size=len(file_bytes),
                )
            )
            continue
        relative_path = normalize_relative_source_path(
            raw_paths[index] if raw_paths else "",
            fallback_filename=upload["original_filename"] or upload["filename"],
        )
        prepared_uploads.append(
            {
                "file": file,
                "file_bytes": file_bytes,
                "upload": upload,
                "relative_path": relative_path,
            }
        )

    if not prepared_uploads:
        raise DocumentIngestionError(
            "No source-traceable documents were found in this folder. "
            "Upload PDF, DOCX, Markdown, or TXT files for source-set review."
        )

    existing_version_counts: dict[str, int] = {}
    try:
        for item in prepared_uploads:
            filename = item["upload"]["filename"]
            if filename not in existing_version_counts:
                existing_version_counts[filename] = component_collection.count_documents(
                    {
                        "flow_id": ObjectId(flow_id),
                        "source_document.filename": filename,
                    }
                )
    except PyMongoError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=(
                "Source-set uploads need MongoDB so document metadata and source references can be saved. "
                "Start MongoDB, then reopen or create a workspace and try again."
            ),
        ) from exc

    batch_version_offsets: dict[str, int] = {}
    ingested_contexts = []
    for item in prepared_uploads:
        filename = item["upload"]["filename"]
        batch_version_offsets[filename] = batch_version_offsets.get(filename, 0) + 1
        version = existing_version_counts.get(filename, 0) + batch_version_offsets[filename]
        try:
            context = ingest_supported_document(
                filename,
                item["file_bytes"],
                version=version,
            )
        except DocumentIngestionError as exc:
            skipped_sources.append(
                source_set_skip_record(
                    filename=item["upload"]["original_filename"],
                    relative_path=item["relative_path"],
                    extension=item["upload"]["extension"],
                    reason_code=source_set_skip_reason(str(exc)),
                    message=str(exc),
                    size=item["upload"].get("size", 0),
                )
            )
            continue
        context["relative_path"] = item["relative_path"]
        ingested_contexts.append(context)

    if not ingested_contexts:
        raise DocumentIngestionError(
            "No source-set files produced source-aware sections. "
            "Check whether the PDFs have extractable text or whether the DOCX/TXT/Markdown files are empty."
        )

    source_set = build_source_set_metadata(
        [item["relative_path"] for item in ingested_contexts],
        source_set_id=source_set_id,
        label=source_set_label,
    )
    source_set.update(
        {
            "selected_count": len(files),
            "skipped_count": len(skipped_sources),
            "skipped_sources": skipped_sources,
            "supported_extensions": sorted(ALLOWED_DOCUMENT_EXTENSIONS),
        }
    )

    prepared_contexts = []
    for context in ingested_contexts:
        source_document = source_document_with_source_set_metadata(
            context["source_document"],
            relative_path=context["relative_path"],
            source_set=source_set,
        )
        context.update(
            {
                "source_document": source_document,
                "relative_path": source_document["relative_path"],
                "folder": source_document.get("folder", ""),
                "source_set": source_set,
            }
        )
        prepared_contexts.append(context)

    return {
        "source_set": source_set,
        "sources": prepared_contexts,
        "skipped_sources": skipped_sources,
    }


def source_set_skip_reason(message: str) -> str:
    lowered = message.lower()
    if "empty" in lowered or "did not contain extractable text" in lowered:
        return "empty_or_no_extractable_text"
    if "exceeds" in lowered or "upload limit" in lowered:
        return "too_large"
    if "malformed" in lowered:
        return "malformed_document"
    if "unsupported file type" in lowered:
        return "unsupported_extension"
    return "ingestion_error"


def source_set_skip_record(
    *,
    filename: str,
    relative_path: str,
    extension: str,
    reason_code: str,
    message: str,
    size: int = 0,
) -> dict:
    return {
        "filename": sanitize_filename(filename),
        "original_filename": filename,
        "relative_path": relative_path or sanitize_filename(filename),
        "extension": extension,
        "size": size,
        "reason_code": reason_code,
        "message": message,
        "status": "skipped",
    }


def source_set_component_metadata(source_context: dict, flow_id: str) -> dict:
    source_document = source_context["source_document"]
    return {
        "flow_id": ObjectId(flow_id),
        "file_id": "",
        "assistant_id": "",
        "vector_store_id": "",
        "size": source_document["size"],
        "type": source_document["type"],
        "status": source_document.get("status", "uploaded"),
        "processing_type": "source_set_ingestion",
        "summary": "",
        "relative_path": source_document.get("relative_path", ""),
        "path": source_document.get("path", ""),
        "folder": source_document.get("folder", ""),
        "source_set_id": source_document.get("source_set_id", ""),
        "source_set": source_document.get("source_set", {}),
        **source_metadata_fields(source_context),
    }


def uploaded_source_payload(component_id: Any, source_context: dict) -> dict:
    source_document = source_context["source_document"]
    return {
        "component_id": str(component_id),
        "name": source_document.get("filename", ""),
        "original_name": source_document.get("original_filename", ""),
        "type": source_document.get("type", ""),
        "status": source_document.get("status", "uploaded"),
        "file_hash": source_document.get("file_hash", ""),
        "size": source_document.get("size", 0),
        "relative_path": source_document.get("relative_path", ""),
        "path": source_document.get("path", ""),
        "folder": source_document.get("folder", ""),
        "source_document_id": source_document.get("id", ""),
        "source_document": source_document,
        "source_set_id": source_document.get("source_set_id", ""),
        "source_set": source_document.get("source_set", {}),
        "chunk_count": len(source_context.get("document_chunks", [])),
        "segment_count": len(source_context.get("source_segments", [])),
    }


def ingestion_http_error(error: DocumentIngestionError) -> HTTPException:
    return HTTPException(status_code=status.HTTP_400_BAD_REQUEST, detail=str(error))


def source_segments_from_page_records(records: list[dict]) -> list[dict]:
    segments = []
    cursor_by_page: dict[int, int] = {}

    for record in records:
        text = record.get("content") or record.get("data") or ""
        if not text:
            continue

        page = record.get("page_number")
        cursor = cursor_by_page.get(page, 0)
        segments.append(
            {
                "text": text,
                "page": page,
                "heading": None,
                "start_char": cursor,
                "end_char": cursor + len(text),
            }
        )
        cursor_by_page[page] = cursor + len(text) + 2

    return segments


RESPONSES_COMPONENT_TYPES = {"pdf", "txt", "md", "html", "docx", "pptx", "web"}
COMPONENT_QA_RESPONSE_MODELS = {
    "pdf": PDFNodeQueryResponse,
    "txt": TXTNodeQueryResponse,
    "md": MDNodeQueryResponse,
    "html": HTMLNodeQueryResponse,
    "docx": DOCXNodeQueryResponse,
    "pptx": PPTXNodeQueryResponse,
    "web": WebNodeQueryResponse,
}
COMPONENT_QA_NODE_TYPES = {
    "pdf": "PDFNode",
    "txt": "TXTNode",
    "md": "MDNode",
    "html": "HTMLNode",
    "docx": "DOCXNode",
    "pptx": "PPTXNode",
    "web": "WebNode",
}


def component_context_text(record: dict | None, max_length: int = 48000) -> str:
    if not record:
        return ""

    chunks = record.get("document_chunks")
    if isinstance(chunks, list) and chunks:
        parts = []
        for chunk in chunks:
            if not isinstance(chunk, dict):
                continue
            text = str(chunk.get("text") or "").strip()
            if not text:
                continue
            parts.append(
                f"[chunk_id={chunk.get('id', '')}; page={chunk.get('page') or ''}; "
                f"section={chunk.get('heading') or chunk.get('section') or ''}]\n{text}"
            )
        context = "\n\n".join(parts)
        if context.strip():
            return context[:max_length]

    segments = record.get("source_segments")
    if isinstance(segments, list) and segments:
        context = "\n\n".join(
            str(segment.get("text") or "").strip()
            for segment in segments
            if isinstance(segment, dict) and str(segment.get("text") or "").strip()
        )
        if context.strip():
            return context[:max_length]

    summary = record.get("summary")
    if isinstance(summary, list):
        context = " ".join(str(item) for item in summary if item)
    elif summary is None:
        context = ""
    else:
        context = str(summary)

    if not context.strip():
        context = str(record.get("content") or record.get("name") or "")
    return context.strip()[:max_length]


def get_component_record_or_404(flow_id: str, component_id: str, component_type: str) -> dict:
    record = component_collection.find_one(
        {
            "flow_id": ObjectId(flow_id),
            "_id": ObjectId(component_id),
            "type": component_type,
        }
    )
    if not record:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail=f"{component_type.upper()} component not found.",
        )
    return record


def update_component_persona(record_id: str, *, instructions: str, persona_name: str) -> None:
    component_collection.update_one(
        {"_id": ObjectId(record_id)},
        {"$set": {"instructions": instructions, "persona_name": persona_name}},
    )


def answer_component_with_responses(
    request: Any,
    component_type: str,
    response_model: Any,
):
    record = get_component_record_or_404(
        request.flow_id,
        request.component_id,
        component_type,
    )
    context = component_context_text(record)
    if not context:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail=f"{component_type.upper()} component has no source context to answer from.",
        )
    answer = generate_component_answer(
        question=request.query,
        context=query_with_workspace_brief(context, request.workspace_brief),
        persona=record.get("persona_name") or "TraceSpace reviewer",
        instructions=record.get("instructions") or "",
        model=getattr(request, "model_name", None),
        workspace_brief=request.workspace_brief,
    )
    entries = [
        response_model(
            id=request.node_id,
            type=COMPONENT_QA_NODE_TYPES.get(component_type, "ResponseNode"),
            data={
                "question": request.query,
                "summ": answer["summ"],
                "df": validate_dataframe(answer["df"]),
                "graph": answer["graph"],
                "flow_id": request.flow_id,
                "component_id": request.component_id,
                "component_type": component_type,
            },
        )
    ]
    if request.request_type == "question":
        entries.append(
            response_model(
                id=str(ObjectId()),
                type="question",
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": component_type,
                },
            )
        )
    return entries


def follow_up_entries(
    *,
    flow_id: str,
    component_id: str,
    component_type: str,
    questions: list[str],
) -> list[ComponentFollowUpQueryResponse]:
    entries = [
        ComponentFollowUpQueryResponse(
            id=str(ObjectId()),
            flow_id=flow_id,
            data={
                "question": question,
                "component_id": component_id,
                "component_type": component_type,
            },
            type="followUp",
            position={"x": 0, "y": 0},
        )
        for question in questions
    ]
    entries.append(
        ComponentFollowUpQueryResponse(
            id=str(ObjectId()),
            flow_id=flow_id,
            position={"x": 0, "y": 0},
            data={
                "question": "",
                "component_id": component_id,
                "component_type": component_type,
            },
            type="question",
        )
    )
    return entries


def follow_up_questions_with_responses(
    request: ComponentFollowUpQueryRequest,
) -> list[ComponentFollowUpQueryResponse]:
    record = get_component_record_or_404(
        request.flow_id,
        request.component_id,
        request.component_type,
    )
    update_component_persona(
        request.component_id,
        instructions=request.instructions,
        persona_name=request.persona_name,
    )
    context = component_context_text(record)
    if not context:
        return follow_up_entries(
            flow_id=request.flow_id,
            component_id=request.component_id,
            component_type=request.component_type,
            questions=[],
        )
    questions = generate_component_follow_up_questions(
        context=context,
        persona=request.persona_name,
        instructions=request.instructions,
        model=request.model_name,
    )
    return follow_up_entries(
        flow_id=request.flow_id,
        component_id=request.component_id,
        component_type=request.component_type,
        questions=questions,
    )


def ground_mindmap_with_source_refs(response_json: dict, source_context: dict) -> dict:
    response_json = validate_ai_mindmap_contract(response_json)
    return attach_source_refs_to_mindmap(
        response_json,
        source_context["source_document"],
        source_context["document_chunks"],
    )


def parse_ai_mindmap_or_422(raw_response: str | dict) -> dict:
    try:
        return parse_ai_mindmap_response(raw_response)
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={"message": "AI graph output failed schema validation.", "errors": exc.errors},
        ) from exc


def workspace_brief_has_context(workspace_brief: dict | None) -> bool:
    if not workspace_brief:
        return False

    if workspace_brief.get("configured") is True:
        return True

    desired_outputs = workspace_brief.get("desired_outputs") or []
    non_default_outputs = [output for output in desired_outputs if output != "mind_map"]

    return any(
        [
            str(workspace_brief.get("goal") or "").strip(),
            str(workspace_brief.get("audience") or "").strip(),
            str(workspace_brief.get("domain_context") or "").strip(),
            str(workspace_brief.get("review_rules") or "").strip(),
            non_default_outputs,
        ]
    )


def query_with_workspace_brief(query: str, workspace_brief: dict | None) -> str:
    if not workspace_brief_has_context(workspace_brief):
        return query

    brief_lines = [
        "Use this structured workspace brief while answering.",
        f"Goal: {workspace_brief.get('goal') or 'Not specified'}",
        f"Audience: {workspace_brief.get('audience') or 'Not specified'}",
        f"Domain context: {workspace_brief.get('domain_context') or 'Not specified'}",
        f"Desired outputs: {', '.join(workspace_brief.get('desired_outputs') or []) or 'Not specified'}",
        f"Source mode: {workspace_brief.get('source_mode') or 'source_plus_context'}",
        f"Assumptions allowed: {bool(workspace_brief.get('assumptions_allowed'))}",
        f"Preset: {workspace_brief.get('preset') or 'custom'}",
        f"Output style: {workspace_brief.get('output_style') or 'technical_reference_map'}",
        f"Preferred node types: {', '.join(workspace_brief.get('node_types') or []) or 'Not specified'}",
        f"Review policy: {', '.join(workspace_brief.get('review_policy') or []) or 'Not specified'}",
        f"Review rules: {workspace_brief.get('review_rules') or 'Not specified'}",
        "",
        "When the brief adds context beyond the source, keep the answer explicit about what is source-backed versus assumption-based.",
        "Respect review policy by marking uncited, low-confidence, or assumption-based nodes as reviewable when requested.",
        "",
        f"User question: {query}",
    ]
    return "\n".join(brief_lines)


def _bounded_json_for_prompt(value: Any, limit: int = 10000) -> str:
    try:
        text = json.dumps(value, indent=2, default=str)
    except TypeError:
        text = str(value)
    if len(text) <= limit:
        return text
    return f"{text[: max(0, limit - 32)]}\n... truncated for prompt budget"


def query_with_follow_up_memory(query: str, request: dict[str, Any] | None) -> str:
    if not isinstance(request, dict):
        return query

    metadata = request.get("metadata") if isinstance(request.get("metadata"), dict) else {}
    memory_context = request.get("memory_context")
    if not isinstance(memory_context, dict):
        memory_context = metadata.get("follow_up_memory")
    if not isinstance(memory_context, dict):
        memory_context = {}

    change_intent = (
        request.get("change_intent")
        or metadata.get("change_intent")
        or memory_context.get("change_intent")
        or ""
    )
    if not memory_context and not change_intent:
        return query

    memory_lines = [
        "Use this follow-up AI memory while answering.",
        "The memory describes what the user had selected, the current graph context, source refs, prior draft/session state, and whether to update, supplement, or compare.",
        f"Change intent: {change_intent or 'supplement'}",
    ]
    if memory_context:
        memory_lines.extend(
            [
                "Follow-up memory context JSON:",
                _bounded_json_for_prompt(memory_context),
            ]
        )
    memory_lines.extend(["", f"User question: {query}"])
    return "\n".join(memory_lines)


def calculate_file_hash(file):
    hasher = sha256()

    if isinstance(file, bytes):
        file = BytesIO(file)

    while chunk := file.read(8192):
        hasher.update(chunk)

    file.seek(0)
    return hasher.hexdigest()


def validate_dataframe(df):
    try:
        if not isinstance(df, list) and not all(isinstance(item, dict) for item in df):
            return []
        else:
            return df

    except ValueError:
        return []


def upload_to_s3(file_bytes, bucket, key):
    try:
        require_settings("aws_access_key_id", "aws_secret_access_key", "bucket_name")
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc

    try:
        s3_client.put_object(Bucket=bucket, Key=key, Body=file_bytes)
    except ClientError as e:
        raise HTTPException(status_code=500, detail=f"S3 Upload Error: {str(e)}")


def extract_text_and_tables(key):
    try:
        require_settings("aws_access_key_id", "aws_secret_access_key", "bucket_name")
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc

    client = boto3.client(
        "textract",
        aws_access_key_id=aws_access_key_id_str,
        aws_secret_access_key=aws_secret_access_key_str,
        region_name="ap-south-1",
    )

    response = client.start_document_analysis(
        DocumentLocation={"S3Object": {"Bucket": bucket_name, "Name": key}},
        FeatureTypes=["TABLES", "FORMS"],
    )

    print(response)

    job_id = response["JobId"]
    print("job_id")
    while True:
        response = client.get_document_analysis(JobId=job_id)
        status = response["JobStatus"]
        if status in ["SUCCEEDED", "FAILED"]:
            break
        time.sleep(5)

    if status == "FAILED":
        raise Exception("Document analysis failed")

    all_blocks = response["Blocks"]
    next_token = response.get("NextToken", None)
    print(all_blocks)

    while next_token:
        response = client.get_document_analysis(JobId=job_id, NextToken=next_token)
        all_blocks.extend(response["Blocks"])
        next_token = response.get("NextToken", None)

    response["Blocks"] = all_blocks

    doc = Document(response)
    lines = [line.text for page in doc.pages for line in page.lines if line.text]
    tables = []
    key_values = {}

    for page in doc.pages:
        for table in page.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text if cell.text else "" for cell in row.cells]
                table_data.append(row_data)
            df = pd.DataFrame(table_data)
            tables.append(df)

        for field in page.form.fields:
            key = field.key.text if field.key and field.key.text else ""
            value = field.value.text if field.value and field.value.text else ""
            key_values[key] = value

    date = next((line for line in lines if "Date" in line), None)
    return lines, tables, key_values, date


def sanitize_path(path):
    """Sanitize the path to remove any invalid characters."""
    return re.sub(r"[^\w\s-]", "_", path).strip()


def camelot_pdf_processing(flow_id, file, flow_type):
    try:
        source_context = prepare_source_upload(file, flow_id, expected_extension="pdf")
        source_document = source_context["source_document"]
        file_bytes = source_context["file_bytes"]
        sanitized_flow_id = sanitize_path(flow_id)
        sanitized_filename = source_document["filename"]

        flow_dir = os.path.join(UPLOAD_DIR, sanitized_flow_id)

        # Ensure the directory exists
        os.makedirs(flow_dir, exist_ok=True)

        # Create the full file path
        file_path = os.path.join(flow_dir, sanitized_filename)

        print(file_path)

        # Save the uploaded file
        with open(file_path, "wb") as f:
            f.write(file_bytes)

        print(f"File saved at: {file_path}")

        print(file_bytes)

        file_hash = source_document["file_hash"]
        print(file_hash)

        # Now, upload the file to S3
        folder = f"uploads/{flow_id}/"
        s3_key = folder + sanitized_filename

        # Pass the local file path to upload_to_s3
        upload_to_s3(file_bytes, bucket_name, s3_key)
        print("File uploaded to S3")

        ocr_config = {
            "languages": ["eng"],
            "strategy": "fast",
        }
        elements = partition_pdf(filename=file_path, **ocr_config)

        content_data = []

        # Extract content (text, title, etc.) from the PDF
        for element in elements:
            if isinstance(element, (Text, Title, NarrativeText, ListItem, UnstructuredHeader)):
                content_data.append(
                    {"data": element.text, "page_number": element.metadata.page_number}
                )

        tables = camelot.read_pdf(
            file_path,
            pages="all",
            flavor="stream",
            edge_tol=900,
            row_tol=6,
            strip_text="\n",
        )

        print(f"Total tables found: {len(tables)}")

        result_tbl_list = []

        # Extract tables from the PDF
        for i, table in enumerate(tables):
            page_number = (
                table.page
            )  # Get the page number from which the table was extracted
            print(f"Table {i + 1} extracted from page: {page_number}\n")
            print(table.df)
            result_tbl_list.append(
                {"data": table.df.to_string(index=False), "page_number": table.page}
            )

        # Combine content_data and result_tbl_list based on the page number
        combined_data = []

        for content_item in content_data:
            page_number = content_item["page_number"]
            # Check if a matching table exists for the current page number
            matching_table = next(
                (
                    table
                    for table in result_tbl_list
                    if table["page_number"] == page_number
                ),
                None,
            )

            # Store content along with the page number and table (if exists)
            if matching_table:
                combined_content = (
                    f"{content_item['data']}\n\nTable:\n{matching_table['data']}"
                )
                combined_data.append(
                    {"page_number": page_number, "content": combined_content}
                )
            else:
                combined_data.append(
                    {"page_number": page_number, "content": content_item["data"]}
                )

        print("-==============-")

        # Print the final combined data
        print(combined_data)

        combined_data_str = json.dumps(combined_data, indent=4)
        source_context["source_segments"] = source_segments_from_page_records(combined_data)
        source_context["document_chunks"] = chunk_source_segments(
            source_context["source_segments"], source_document["id"]
        )

        chunks = do_semantic_chunking(combined_data_str)
        print(chunks)

        summary = process_pdf_summary(chunks)

        if flow_type == "manual":

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "size": len(file_bytes),
                "s3_path": s3_key,
                "type": "pdf",
                "processing_type": "custom",
                "summary": summary,
                **source_metadata_fields(source_context),
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id
            EmbeddingsDocuments = []
            for i in range(len(chunks)):
                metadata = {
                    "component_id": str(component_id),
                    "file_name": source_document["filename"],
                    "flow_id": flow_id,
                    "source_document_id": source_document["id"],
                }
                EmbeddingsDocuments.append(
                    LangDocument(page_content=chunks[i].page_content, metadata=metadata)
                )
            uuids = [str(uuid4()) for _ in range(len(EmbeddingsDocuments))]

            vector_store_from_client.add_documents(documents=EmbeddingsDocuments, ids=uuids)
            return {"component_id": str(component_id), "type": "pdf"}

        else:

            template = """You are tasked with generating a JSON mind map for given summary of the pdf document and that should be compatible with React Flow for rendering a flow diagram. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be a maximum of 5 `response` nodes.

                2. **Node Relationships:**
                - The `dataSource` node should be connected to all `response` nodes.
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "pdf", !!!DOESN"T CHANGES
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{filename}"  // Empty object or file metadata
                        }}

                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "MDNode | WebNode | MultipleQA | other",
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<summary or answer>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "pdf"
                            }}
                        }}

                5. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                6. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                Here is the PDF summary for which you need to generate the mind map:
                {summary_pdf}

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - If fewer than 5 `response` nodes are required, adjust accordingly.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

            template = append_ai_graph_prompt_contract(template)
            prompt = PromptTemplate.from_template(template)

            lm_chain = prompt | llm

            answer = lm_chain.invoke(
                    {"summary_pdf": summary, "flow_id": flow_id, "filename": source_document["filename"]}
            )

            responseList = answer.content

            print(responseList)

            response_json = parse_ai_mindmap_or_422(responseList)

            print(response_json)

            response_json = ground_mindmap_with_source_refs(response_json, source_context)

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "type": "pdf",
                "processing_type": "gpt",
                "mindmap_json": response_json,
                **source_metadata_fields(source_context),
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id

            return {
                "component_id": str(component_id),
                "type": "pdf",
                "mindmap_json": response_json,
                "flow_type": "automatic"
            }

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=e)


def extract_text_from_upload(file: UploadFile) -> str:
    file.file.seek(0)
    content = file.file.read()
    extension = file.filename.split(".")[-1].lower()

    if extension == "txt":
        return content.decode("utf-8", errors="ignore")

    elif extension == "md":
        html = markdown.markdown(content.decode("utf-8", errors="ignore"))
        return BeautifulSoup(html, "html.parser").get_text()

    elif extension == "html":
        return BeautifulSoup(content, "html.parser").get_text()

    elif extension == "docx":
        doc = Document(io.BytesIO(content))
        return "\n".join([p.text for p in doc.paragraphs])

    elif extension == "pptx":
        prs = Presentation(io.BytesIO(content))
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)

    else:
        raise ValueError(f"Unsupported file type: {extension}")


def is_within_gpt4o_token_limit(file: UploadFile) -> bool:
    try:
        text = extract_text_from_upload(file)
        encoding = tiktoken.get_encoding("cl100k_base")
        token_count = len(encoding.encode(text))
        print("token_count : ", token_count)
        return token_count <= GPT_4O_MAX_TOKENS
    except Exception as e:
        print(f"Error checking token count: {e}")
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Malformed or unreadable document: {str(e)}",
        ) from e


def process_pdf_summary(chunks):
    map_template = """Write a concise summary of the following content:
                    {content}
                    Summary:
                    """
    map_prompt = PromptTemplate.from_template(map_template)
    map_chain = LLMChain(prompt=map_prompt, llm=llm)

    reduce_template = """The following is a set of summaries:
                        {doc_summaries}
                        Summarize the above summaries with all the key details.
                        Summary:"""
    reduce_prompt = PromptTemplate.from_template(reduce_template)
    reduce_chain = LLMChain(prompt=reduce_prompt, llm=llm)

    stuff_chain = StuffDocumentsChain(
        llm_chain=reduce_chain, document_variable_name="doc_summaries"
    )

    reduce_documents_chain = ReduceDocumentsChain(
        combine_documents_chain=stuff_chain,
    )

    map_reduce_chain = MapReduceDocumentsChain(
        llm_chain=map_chain,
        document_variable_name="content",
        reduce_documents_chain=reduce_documents_chain,
    )

    small_chunks = chunks[:24]

    summary = map_reduce_chain.run(small_chunks)
    print(summary)
    return summary


def use_aws_textract(file, flow_id, flow_type):
    source_context = prepare_source_upload(file, flow_id, expected_extension="pdf")
    source_document = source_context["source_document"]
    file_bytes = source_context["file_bytes"]
    file_hash = source_document["file_hash"]
    print(file_hash)
    file_name = source_document["filename"]
    folder = f"uploads/{flow_id}/"
    s3_key = folder + file_name
    upload_to_s3(file_bytes, bucket_name, s3_key)
    print("uploaded")
    data_for_chunking = "Data : "
    try:
        lines, tables, key_values, date = extract_text_and_tables(s3_key)
        for i, table in enumerate(tables):
            print(f"Table {i+1}:\n")
            print(table.to_string() + "\n\n")
            data_for_chunking = data_for_chunking + f"Table {i+1}:\n"
            data_for_chunking = data_for_chunking + table.to_string() + "\n\n"
        print("Extracted Key-Value Pairs:\n")
        for key, value in key_values.items():
            print(f"{key}: {value}\n")
            data_for_chunking = data_for_chunking + f"{key}: {value}\n"
        if date:
            print(f"Date: {date}\n")
            data_for_chunking = data_for_chunking + f"Date: {date}\n"
    except Exception as e:
        print(e)
        raise HTTPException(status_code=500, detail=f"Text extraction error: {str(e)}")
    source_context["source_segments"] = [
        {
            "text": data_for_chunking,
            "page": None,
            "heading": "Textract output",
            "start_char": 0,
            "end_char": len(data_for_chunking),
        }
    ]
    source_context["document_chunks"] = chunk_source_segments(
        source_context["source_segments"], source_document["id"]
    )
    chunks = do_semantic_chunking(data_for_chunking)
    print(chunks)

    summary = process_pdf_summary(chunks)

    if flow_type == "manual":

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "size": len(file_bytes),
            "s3_path": s3_key,
            "type": "pdf",
            "processing_type": "aws",
            "summary": summary,
            **source_metadata_fields(source_context),
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id
        EmbeddingsDocuments = []
        for i in range(len(chunks)):
            metadata = {
                "component_id": str(component_id),
                "file_name": source_document["filename"],
                "flow_id": flow_id,
                "source_document_id": source_document["id"],
            }
            EmbeddingsDocuments.append(
                LangDocument(page_content=chunks[i].page_content, metadata=metadata)
            )
        uuids = [str(uuid4()) for _ in range(len(EmbeddingsDocuments))]

        vector_store_from_client.add_documents(documents=EmbeddingsDocuments, ids=uuids)
        return {"component_id": str(component_id), "type": "pdf"}

    else:

        template = """You are tasked with generating a JSON mind map for given summary of the pdf document and that should be compatible with React Flow for rendering a flow diagram. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be a maximum of 5 `response` nodes.

                2. **Node Relationships:**
                - The `dataSource` node should be connected to all `response` nodes.
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "pdf", !!!DOESN"T CHANGES
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{filename}"  // Empty object or file metadata
                        }}

                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "MDNode | WebNode | MultipleQA | other",
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<summary or answer>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "pdf"
                            }}
                        }}

                5. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                6. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                Here is the PDF summary for which you need to generate the mind map:
                {summary_pdf}

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - If fewer than 5 `response` nodes are required, adjust accordingly.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

        template = append_ai_graph_prompt_contract(template)
        prompt = PromptTemplate.from_template(template)

        lm_chain = prompt | llm

        answer = lm_chain.invoke(
            {"summary_pdf": summary, "flow_id": flow_id, "filename": source_document["filename"]}
        )

        responseList = answer.content

        print(responseList)

        response_json = parse_ai_mindmap_or_422(responseList)

        print(response_json)

        response_json = ground_mindmap_with_source_refs(response_json, source_context)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "type": "pdf",
            "processing_type": "gpt",
            "mindmap_json": response_json,
            **source_metadata_fields(source_context),
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id

        return {
            "component_id": str(component_id),
            "type": "pdf",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }


def do_semantic_chunking(docs):
    documents = text_splitter.create_documents([docs])
    print("Number of chunks created: ", len(documents))
    for i in range(len(documents)):
        print()
        print(f"CHUNK : {i+1}")
        print(documents[i].page_content)
    return documents


def get_relevant_passage(query: str, flow_id: str, component_id: str, n_results: int):
    results = vector_store_from_client.similarity_search(
        query=query,
        k=2,
        filter={"$and": [{"component_id": component_id}, {"flow_id": flow_id}]},
    )
    return [doc.page_content for doc in results]


def fetch_question_answer_from_node_collection(parent_id: str, flow_id: str):
    try:
        record = node_collection.find_one(
            {
                "_id": ObjectId(parent_id),
                "flow_id": ObjectId(flow_id),
                "is_delete": "false",
            }
        )

        print("Fetched record:", record)

        if not record:
            return None, None

        print("Question value:", record.get("question"))

        question = record.get("question", None)
        answer = None

        if record["type"] == "pdf":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "txt":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "md":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "html":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "docx":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "pptx":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "image":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])

            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "audio":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "youtube":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "video":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "sql":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "csv":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "web":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        elif record["type"] == "MultipleQA":
            answer = "Answer: " + str(record.get("summ", "Answer not found."))
            df_list = record.get("df", [])
            df_string = " | ".join([str(item) for item in df_list])

            answer += " DataFrame: " + df_string

        print("Answer:", answer)

        return question, answer
    except Exception as e:
        traceback.print_exc()
        print(f"Error: {e}")
        return None, None


@app.post("/create-flow/")
def create_flow(flow: dict):
    try:
        flow_type = flow.get("flow_type") or "manual"
        flow_data = {
            "flow_name": flow.get("flow_name") or "New Flow",
            "flow_json": flow.get("flow_json") or "",
            "summary": flow.get("summary") or "",
            "flow_type": flow_type,
        }
        try:
            flow_id = flow_collection.insert_one(flow_data).inserted_id
        except PyMongoError:
            flow_id = local_create_flow(flow_data)["_id"]
        return {
            "flow_id": str(flow_id),
            "flow_name": flow_data["flow_name"],
            "flow_type": flow_type,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error creating flow: {str(e)}")


@app.delete("/delete-flow/{flow_id}")
def delete_flow(flow_id: str):
    try:
        if not ObjectId.is_valid(flow_id):
            if local_delete_flow(flow_id):
                return {"status": "success"}
            raise HTTPException(status_code=404, detail="Flow not found")

        flow_object_id = ObjectId(flow_id)

        try:
            components = component_collection.find({"flow_id": flow_object_id})
            for component in components:
                component_id = component["_id"]
                node_collection.delete_many({"component_id": component_id})
            component_collection.delete_many({"flow_id": flow_object_id})

            result = flow_collection.delete_one({"_id": flow_object_id})
            if result.deleted_count == 0:
                local_delete_flow(flow_id)
        except PyMongoError:
            local_delete_flow(flow_id)

        return {"status": "success"}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error deleting flow: {str(e)}")


@app.get("/flows/", response_model=List[Flow])
def list_flows():
    try:
        flows = flow_collection.find()
        mongo_flows = [normalize_flow_record(flow) for flow in flows]
        local_flows = [normalize_flow_record(flow) for flow in load_local_flows()]
        mongo_ids = {flow["flow_id"] for flow in mongo_flows}
        return mongo_flows + [
            flow for flow in local_flows if flow["flow_id"] not in mongo_ids
        ]
    except PyMongoError:
        return [normalize_flow_record(flow) for flow in load_local_flows()]


@app.get("/flows/{flow_id}", response_model=Flow)
def get_flow(flow_id: str):
    flow = get_workspace_flow_or_404(flow_id)
    return normalize_flow_record(flow)


def _snapshot_debug_counts(flow_json: str) -> dict[str, Any]:
    summary: dict[str, Any] = {
        "bytes": len((flow_json or "").encode("utf-8")),
    }
    try:
        snapshot = json.loads(flow_json or "{}")
    except (TypeError, json.JSONDecodeError) as exc:
        summary["parse_error"] = exc.__class__.__name__
        return summary
    if not isinstance(snapshot, dict):
        summary["parse_error"] = "not_object"
        return summary

    for key in ("nodes", "edges", "activity_events", "ai_action_runs", "automations"):
        value = snapshot.get(key)
        summary[key] = len(value) if isinstance(value, list) else 0

    source_library = snapshot.get("source_library")
    if isinstance(source_library, list):
        summary["source_library_items"] = len(source_library)
    elif isinstance(source_library, dict):
        documents = source_library.get("documents")
        summary["source_library_items"] = len(documents) if isinstance(documents, list) else 0
    else:
        summary["source_library_items"] = 0
    return summary


def _print_flow_update_debug(
    *,
    flow_id: str,
    storage: str,
    started_at: float,
    input_summary: dict[str, Any],
    persisted_summary: dict[str, Any],
    result: Any = None,
) -> None:
    payload = {
        "flow_id": flow_id,
        "storage": storage,
        "elapsed_ms": round((time.perf_counter() - started_at) * 1000, 1),
        "input": input_summary,
        "persisted": persisted_summary,
    }
    if result is not None:
        payload["mongo"] = {
            "matched": getattr(result, "matched_count", None),
            "modified": getattr(result, "modified_count", None),
            "upserted_id": str(getattr(result, "upserted_id", "") or ""),
        }
    print(f"[flow-update] {json.dumps(payload, separators=(',', ':'))}")


@app.put("/flow-update/")
def update_flow(update_data: Flow):
    started_at = time.perf_counter()
    input_summary = _snapshot_debug_counts(update_data.flow_json)
    try:
        repaired_flow_json = repair_flow_snapshot_for_persistence(
            update_data.flow_json,
            flow_id=update_data.flow_id,
            flow_name=update_data.flow_name,
            flow_type=update_data.flow_type,
            summary=update_data.summary,
        )
        updates = {
            "flow_name": update_data.flow_name,
            "flow_json": repaired_flow_json,
            "flow_type": update_data.flow_type,
            "summary": update_data.summary,
        }
        result = None
        if ObjectId.is_valid(update_data.flow_id):
            try:
                result = flow_collection.update_one(
                    {"_id": ObjectId(update_data.flow_id)},
                    {"$set": updates},
                )
            except PyMongoError:
                result = None
        persisted_summary = _snapshot_debug_counts(repaired_flow_json)

        if result is None or result.matched_count == 0:
            promoted_flow = promote_local_flow_to_mongo(update_data.flow_id)
            if promoted_flow:
                result = flow_collection.update_one(
                    {"_id": ObjectId(update_data.flow_id)},
                    {"$set": updates},
                )
                if result.matched_count:
                    local_update_flow(update_data.flow_id, updates)
                    _print_flow_update_debug(
                        flow_id=update_data.flow_id,
                        storage="mongo_promoted",
                        started_at=started_at,
                        input_summary=input_summary,
                        persisted_summary=persisted_summary,
                        result=result,
                    )
                    return {
                        "flow_id": str(update_data.flow_id),
                        "message": "Flow updated successfully",
                    }
            if local_update_flow(update_data.flow_id, updates):
                _print_flow_update_debug(
                    flow_id=update_data.flow_id,
                    storage="local",
                    started_at=started_at,
                    input_summary=input_summary,
                    persisted_summary=persisted_summary,
                    result=result,
                )
                return {
                    "flow_id": str(update_data.flow_id),
                    "message": "Flow updated successfully",
                }
            raise HTTPException(
                status_code=status.HTTP_404_NOT_FOUND, detail="Flow not found"
            )

        _print_flow_update_debug(
            flow_id=update_data.flow_id,
            storage="mongo",
            started_at=started_at,
            input_summary=input_summary,
            persisted_summary=persisted_summary,
            result=result,
        )
        return {
            "flow_id": str(update_data.flow_id),
            "message": "Flow updated successfully",
        }

    except GraphSchemaError as e:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail={"message": "Flow snapshot schema validation failed.", "errors": e.errors},
        )
    except Exception as e:
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"An error occurred: {str(e)}",
        )


def get_workspace_graph_or_404(flow_id: str) -> dict:
    flow = get_workspace_flow_or_404(flow_id)
    return build_workspace_graph(flow, source_components=get_source_components(flow_id))


def get_source_components(flow_id: str) -> list[dict]:
    if not ObjectId.is_valid(flow_id):
        return []

    try:
        return list(
            component_collection.find(
                {
                    "flow_id": ObjectId(flow_id),
                    "$or": [
                        {"source_document": {"$exists": True}},
                        {"source_document_id": {"$exists": True}},
                        {"document_chunks": {"$exists": True}},
                    ],
                }
            )
        )
    except PyMongoError:
        return []


def get_workspace_flow_or_404(flow_id: str) -> dict:
    flow = None
    if ObjectId.is_valid(flow_id):
        try:
            flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        except PyMongoError:
            flow = None

    if not flow:
        flow = promote_local_flow_to_mongo(flow_id) or local_find_flow(flow_id)

    if not flow:
        raise HTTPException(status_code=404, detail="Workspace not found.")

    return flow


def get_upload_flow_or_400(flow_id: str) -> dict:
    if not flow_id or flow_id == "undefined" or not ObjectId.is_valid(flow_id):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Create or open a workspace before uploading a source document.",
        )

    try:
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    except PyMongoError as exc:
        if local_find_flow(flow_id):
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail=(
                    "Source uploads need MongoDB so document metadata and source references can be saved. "
                    "Start MongoDB, then reopen or create a workspace and try the DOCX upload again."
                ),
            ) from exc
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail="MongoDB is unavailable. Start MongoDB, then try the source upload again.",
        ) from exc

    if not flow and local_find_flow(flow_id):
        flow = promote_local_flow_to_mongo(flow_id)
        if not flow:
            raise HTTPException(
                status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
                detail=(
                    "Source uploads need MongoDB so document metadata and source references can be saved. "
                    "MongoDB is running, but this local workspace could not be promoted for upload. "
                    "Reopen or create the workspace and try again."
                ),
            )

    if not flow:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Workspace not found.")

    return flow


def get_workspace_branch_or_404(flow_id: str, node_id: str) -> dict:
    graph = get_workspace_graph_or_404(flow_id)

    if not any(node["id"] == node_id for node in graph["nodes"]):
        raise HTTPException(status_code=404, detail="Branch root node not found.")

    return select_branch(graph, node_id)


def monday_task_nodes_from_graph(graph: dict) -> list[dict]:
    return select_monday_task_nodes(graph)


def utc_timestamp() -> str:
    return datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z"


@app.get("/api/workspaces/{flow_id}/exports/json")
def export_workspace_json(flow_id: str):
    return get_workspace_graph_or_404(flow_id)


@app.get("/api/workspaces/{flow_id}/sources")
def get_workspace_sources(flow_id: str):
    return get_workspace_graph_or_404(flow_id)["source_library"]


@app.post("/api/workspaces/{flow_id}/sources/source-set")
def upload_workspace_source_set(
    flow_id: str,
    files: list[UploadFile] = File(...),
    relative_paths: list[str] | None = Form(None),
    source_set_id: str | None = Form(None),
    source_set_label: str | None = Form(None),
):
    flow = get_upload_flow_or_400(flow_id)
    try:
        prepared = prepare_source_set_uploads(
            files,
            flow_id=flow_id,
            relative_paths=relative_paths,
            source_set_id=source_set_id,
            source_set_label=source_set_label,
        )
    except DocumentIngestionError as exc:
        raise ingestion_http_error(exc) from exc

    uploaded_sources = []
    source_components = []
    try:
        for source_context in prepared["sources"]:
            component_metadata = source_set_component_metadata(source_context, flow_id)
            component_id = component_collection.insert_one(component_metadata).inserted_id
            saved_component = {**component_metadata, "_id": component_id}
            source_components.append(saved_component)
            uploaded_sources.append(uploaded_source_payload(component_id, source_context))
    except PyMongoError as exc:
        raise HTTPException(
            status_code=status.HTTP_503_SERVICE_UNAVAILABLE,
            detail=(
                "Source-set upload could not save document metadata. "
                "Check MongoDB and try the folder upload again."
            ),
        ) from exc

    existing_components = get_source_components(flow_id)
    existing_ids = {str(component.get("_id")) for component in existing_components}
    all_source_components = [
        *existing_components,
        *[
            component
            for component in source_components
            if str(component.get("_id")) not in existing_ids
        ],
    ]
    source_library = build_workspace_graph(
        flow,
        source_components=all_source_components,
    )["source_library"]

    return {
        "uploaded_sources": uploaded_sources,
        "skipped_sources": prepared.get("skipped_sources", []),
        "source_set": {
            **prepared["source_set"],
            "source_count": len(uploaded_sources),
            "skipped_count": len(prepared.get("skipped_sources", [])),
        },
        "source_library": source_library,
    }


@app.get("/api/workspaces/{flow_id}/completeness-review")
def get_workspace_completeness_review(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return graph.get("views", {}).get("completeness_review") or graph_to_completeness_review(graph)


@app.get("/api/workspaces/{flow_id}/team-roadmap")
def get_workspace_team_roadmap(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return graph_to_team_roadmap(graph)


@app.post("/api/workspaces/{flow_id}/sources/{source_id}/reconcile/preview")
def preview_source_reconciliation(flow_id: str, source_id: str, request: dict[str, Any] | None = None):
    graph = get_workspace_graph_or_404(flow_id)
    scope = request.get("scope") if isinstance(request, dict) else None
    try:
        return generate_source_reconciliation_preview(
            graph,
            source_id=source_id,
            scope=scope if isinstance(scope, dict) else {"type": "source", "source_id": source_id},
        )
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "Source reconciliation preview failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc


@app.get("/api/workspaces/{flow_id}/exports/markdown")
def export_workspace_markdown(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return Response(
        content=graph_to_markdown(graph),
        media_type="text/markdown",
    )


@app.get("/api/workspaces/{flow_id}/exports/completeness-review.md")
def export_workspace_completeness_review_markdown(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return Response(
        content=graph_to_completeness_markdown(graph),
        media_type="text/markdown",
    )


@app.get("/api/workspaces/{flow_id}/exports/executive.md")
@app.get("/api/workspaces/{flow_id}/exports/executive-output.md")
def export_workspace_executive_markdown(flow_id: str):
    artifact = _latest_ai_draft_artifact(flow_id, {"executive_summary", "executive_output"})
    if artifact and artifact.get("artifact_type") == "executive_summary":
        content = export_executive_summary_markdown(artifact_export_data(artifact))
    elif artifact and artifact.get("artifact_type") == "executive_output":
        content = export_executive_output_markdown(artifact_export_data(artifact))
    else:
        graph = get_workspace_graph_or_404(flow_id)
        content = graph_to_executive_markdown(graph)
    return Response(content=content, media_type="text/markdown")


@app.get("/api/workspaces/{flow_id}/exports/article.md")
@app.get("/api/workspaces/{flow_id}/exports/news-article.md")
def export_workspace_news_article_markdown(flow_id: str):
    artifact = _latest_ai_draft_artifact(flow_id, {"news_article"})
    if artifact:
        content = artifact_to_news_article_markdown(artifact)
    else:
        graph = get_workspace_graph_or_404(flow_id)
        content = graph_to_news_article_markdown(graph)
    return Response(
        content=content,
        media_type="text/markdown",
    )


def _latest_ai_draft_artifact(flow_id: str, artifact_types: set[str]) -> dict[str, Any] | None:
    return select_latest_ai_draft_artifact(
        list_ai_draft_sessions_for_workspace(flow_id),
        artifact_types,
    )


@app.get("/api/workspaces/{flow_id}/exports/team-roadmap.md")
def export_workspace_team_roadmap_markdown(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return Response(
        content=graph_to_team_roadmap_markdown(graph),
        media_type="text/markdown",
    )


@app.get("/api/workspaces/{flow_id}/exports/csv")
def export_workspace_csv(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    rows = graph_to_task_rows(graph)
    return Response(
        content=export_task_rows(rows),
        media_type="text/csv",
        headers={
            "Content-Disposition": f'attachment; filename="{flow_id}-tasks.csv"'
        },
    )


@app.get("/api/workspaces/{flow_id}/exports/opml")
def export_workspace_opml(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return Response(
        content=graph_to_opml(graph),
        media_type="application/xml",
    )


@app.get("/api/workspaces/{flow_id}/exports/mmd-json")
def export_workspace_mmd_json(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return graph_to_mmd_json(graph)


@app.get("/api/workspaces/{flow_id}/exports/mermaid")
def export_workspace_mermaid(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return Response(
        content=graph_to_mermaid(graph),
        media_type="text/plain",
    )


@app.get("/api/workspaces/{flow_id}/branches/{node_id}/exports/json")
def export_branch_json(flow_id: str, node_id: str):
    return get_workspace_branch_or_404(flow_id, node_id)


@app.get("/api/ai/prompt-profiles")
def get_ai_prompt_profiles():
    return {"profiles": list_prompt_profiles()}


def save_ai_draft_session(session: dict) -> dict:
    normalized = validate_ai_draft_session(session)
    try:
        ai_draft_session_collection.update_one(
            {
                "workspace_id": normalized["workspace_id"],
                "session_id": normalized["session_id"],
            },
            {"$set": normalized},
            upsert=True,
        )
    except PyMongoError:
        local_save_ai_draft_session(normalized)
    return normalized


def list_ai_draft_sessions_for_workspace(flow_id: str) -> list[dict]:
    sessions: list[dict] = []
    try:
        sessions = list(
            ai_draft_session_collection.find(
                {"workspace_id": flow_id},
                {"_id": 0},
            )
        )
    except PyMongoError:
        sessions = []
    if not sessions:
        sessions = [
            session
            for session in load_local_ai_draft_sessions()
            if session.get("workspace_id") == flow_id
        ]
    return [validate_ai_draft_session(session) for session in sessions]


def _ai_usage_from_metadata(metadata: dict | None) -> dict:
    if not isinstance(metadata, dict):
        return {}
    usage = metadata.get("usage") if isinstance(metadata.get("usage"), dict) else metadata
    return {
        "input_tokens": _int_ai_usage(usage.get("input_tokens")),
        "output_tokens": _int_ai_usage(usage.get("output_tokens")),
        "total_tokens": _int_ai_usage(usage.get("total_tokens") or usage.get("estimated_tokens")),
        "estimated_cost_usd": usage.get("estimated_cost_usd"),
        "cost_source": usage.get("cost_source") or metadata.get("usage_cost_source"),
    }


def _int_ai_usage(value) -> int:
    try:
        return max(0, int(float(value or 0)))
    except (TypeError, ValueError):
        return 0


def _add_ai_usage(left: dict, right: dict) -> dict:
    result = {
        "input_tokens": int(left.get("input_tokens") or 0) + int(right.get("input_tokens") or 0),
        "output_tokens": int(left.get("output_tokens") or 0) + int(right.get("output_tokens") or 0),
        "total_tokens": int(left.get("total_tokens") or 0) + int(right.get("total_tokens") or 0),
    }
    costs = []
    for value in (left.get("estimated_cost_usd"), right.get("estimated_cost_usd")):
        if isinstance(value, str) and value.startswith("$"):
            try:
                costs.append(float(value[1:]))
            except ValueError:
                pass
    if costs:
        result["estimated_cost_usd"] = f"${sum(costs):.4f}"
    return result


def summarize_ai_usage_for_workspace(flow_id: str) -> dict:
    sessions = list_ai_draft_sessions_for_workspace(flow_id)
    totals = {"input_tokens": 0, "output_tokens": 0, "total_tokens": 0}
    by_session = []
    for session in sessions:
        session_total = {"input_tokens": 0, "output_tokens": 0, "total_tokens": 0}
        revision_usages = []
        for revision in session.get("revisions", []):
            usage = _ai_usage_from_metadata(revision.get("metadata"))
            if not usage.get("total_tokens"):
                continue
            revision_usages.append(
                {
                    "revision_id": revision.get("revision_id", ""),
                    "created_at": revision.get("created_at", ""),
                    "model": revision.get("model") or revision.get("metadata", {}).get("model", ""),
                    **usage,
                }
            )
            session_total = _add_ai_usage(session_total, usage)
        if not revision_usages:
            session_total = _ai_usage_from_metadata(session.get("metadata"))
        totals = _add_ai_usage(totals, session_total)
        by_session.append(
            {
                "session_id": session.get("session_id", ""),
                "status": session.get("status", ""),
                "selected_model": session.get("selected_model", ""),
                "created_at": session.get("created_at", ""),
                **session_total,
                "revisions": revision_usages,
            }
        )
    return {
        "workspace_id": flow_id,
        **totals,
        "session_count": len(sessions),
        "sessions": by_session,
    }


def get_ai_draft_session_or_404(flow_id: str, session_id: str) -> dict:
    session = None
    try:
        session = ai_draft_session_collection.find_one(
            {"workspace_id": flow_id, "session_id": session_id},
            {"_id": 0},
        )
    except PyMongoError:
        session = None
    if not session:
        session = local_find_ai_draft_session(flow_id, session_id)
    if not session:
        raise HTTPException(status_code=404, detail="AI draft session not found.")
    try:
        return validate_ai_draft_session(session)
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI draft session failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc


def _flow_snapshot(flow: dict) -> dict:
    try:
        snapshot = json.loads(flow.get("flow_json") or "{}")
    except json.JSONDecodeError:
        snapshot = {}
    if not isinstance(snapshot, dict):
        snapshot = {}
    snapshot.setdefault("nodes", [])
    snapshot.setdefault("edges", [])
    snapshot.setdefault("viewport", {})
    return snapshot


def _persist_flow_snapshot(flow_id: str, snapshot: dict) -> None:
    updates = {"flow_json": json.dumps(snapshot)}
    updated = False
    if ObjectId.is_valid(flow_id):
        try:
            result = flow_collection.update_one(
                {"_id": ObjectId(flow_id)},
                {"$set": updates},
            )
            updated = result.matched_count > 0
        except PyMongoError:
            updated = False
    if not updated:
        updated = local_update_flow(flow_id, updates)
    if not updated:
        raise HTTPException(status_code=404, detail="Workspace not found.")


def _react_node_from_graph_node(node: dict, index: int) -> dict:
    node_type = node.get("node_type") or "concept"
    title = str(node.get("title") or "")
    body = str(node.get("summary") or node.get("body") or "")
    source_refs = node.get("source_refs", [])
    if not isinstance(source_refs, list):
        source_refs = []
    external_refs = node.get("external_refs", {})
    if not isinstance(external_refs, (dict, list)):
        external_refs = {}
    metadata = node.get("metadata", {})
    if not isinstance(metadata, dict):
        metadata = {}
    return {
        "id": node.get("id", ""),
        "type": "response",
        "position": metadata.get("position") or {"x": 120 + (index % 4) * 260, "y": 160 + index * 120},
        "data": {
            "title": title,
            "body": body,
            "summary": body,
            "summ": body,
            "node_type": node_type,
            "status": node.get("status", ""),
            "priority": node.get("priority", ""),
            "owner_id": node.get("owner_id", ""),
            "due_date": node.get("due_date", ""),
            "confidence": node.get("confidence"),
            "source_refs": source_refs,
            "external_refs": external_refs,
            "metadata": metadata,
            "manual": True,
            "display": {
                "collapsed": False,
                "layoutMode": metadata.get("layout_mode") or "vertical-children",
            },
            "data": {
                "title": title,
                "body": body,
                "summary": body,
                "summ": body or title,
                "query": "",
                "df": [],
                "graph": {},
                "source_refs": source_refs,
                "status": node.get("status", ""),
            },
        },
        "deletable": True,
        "targetPosition": "left",
        "sourcePosition": "right",
    }


def _react_edge_from_graph_edge(edge: dict) -> dict:
    metadata = edge.get("metadata", {}) if isinstance(edge.get("metadata"), dict) else {}
    relationship_type = edge.get("relationship_type", "contains")
    return {
        "id": edge.get("id", ""),
        "source": edge.get("source_node_id", ""),
        "target": edge.get("target_node_id", ""),
        "type": metadata.get("react_flow_type", ""),
        "animated": metadata.get("animated", False),
        "relationship_type": relationship_type,
        "confidence": metadata.get("confidence", edge.get("confidence", "")),
        "review_state": metadata.get("review_state", edge.get("review_state", "")),
        "source_refs": edge.get("source_refs", []) if isinstance(edge.get("source_refs"), list) else [],
        "data": {
            "relationship_type": relationship_type,
            "confidence": metadata.get("confidence", edge.get("confidence", "")),
            "review_state": metadata.get("review_state", edge.get("review_state", "")),
            "source_signal": metadata.get("source_signal", ""),
            "rationale": metadata.get("rationale", ""),
            "assumptions": metadata.get("assumptions", []),
            "artifact_id": metadata.get("artifact_id", ""),
        },
    }


def _accepted_graph_focus_viewport(snapshot: dict, accepted_node_ids: list[str]) -> dict:
    accepted_ids = {str(node_id) for node_id in accepted_node_ids if str(node_id)}
    if not accepted_ids:
        return snapshot.get("viewport", {}) if isinstance(snapshot.get("viewport"), dict) else {}

    positions = []
    for node in snapshot.get("nodes", []):
        if not isinstance(node, dict) or str(node.get("id", "")) not in accepted_ids:
            continue
        position = node.get("position") if isinstance(node.get("position"), dict) else {}
        try:
            positions.append((float(position.get("x", 0)), float(position.get("y", 0))))
        except (TypeError, ValueError):
            continue

    if not positions:
        return snapshot.get("viewport", {}) if isinstance(snapshot.get("viewport"), dict) else {}

    min_x = min(position[0] for position in positions)
    max_x = max(position[0] for position in positions)
    min_y = min(position[1] for position in positions)
    max_y = max(position[1] for position in positions)
    width = max(max_x - min_x + 320, 320)
    height = max(max_y - min_y + 220, 220)
    zoom = min(1, max(0.65, min(1080 / width, 620 / height)))
    center_x = min_x + (max_x - min_x) / 2
    center_y = min_y + (max_y - min_y) / 2
    return {
        "x": round(640 - center_x * zoom, 2),
        "y": round(320 - center_y * zoom, 2),
        "zoom": round(zoom, 3),
    }


def _append_accepted_graph_to_flow_snapshot(
    flow: dict,
    accept_result: dict,
    accepted_graph: dict,
) -> dict:
    snapshot = _flow_snapshot(flow)
    node_lookup = {
        node.get("id"): node
        for node in accepted_graph.get("nodes", [])
        if isinstance(node, dict)
    }
    edge_lookup = {
        edge.get("id"): edge
        for edge in accepted_graph.get("edges", [])
        if isinstance(edge, dict)
    }
    existing_node_ids = {
        node.get("id")
        for node in snapshot.get("nodes", [])
        if isinstance(node, dict)
    }
    existing_edge_ids = {
        edge.get("id")
        for edge in snapshot.get("edges", [])
        if isinstance(edge, dict)
    }

    for operation in accept_result.get("patch_operations", []):
        if not isinstance(operation, dict):
            continue
        op = operation.get("op")
        if op == "remove_node":
            node_id = operation.get("node_id")
            snapshot["nodes"] = [
                node for node in snapshot.get("nodes", [])
                if not isinstance(node, dict) or node.get("id") != node_id
            ]
            existing_node_ids.discard(node_id)
            continue
        if op == "remove_edge":
            edge_id = operation.get("edge_id")
            snapshot["edges"] = [
                edge for edge in snapshot.get("edges", [])
                if not isinstance(edge, dict) or edge.get("id") != edge_id
            ]
            existing_edge_ids.discard(edge_id)
            continue
        if op == "update_node":
            node_id = operation.get("node_id")
            graph_node = node_lookup.get(node_id)
            if not graph_node:
                continue
            updated_node = _react_node_from_graph_node(graph_node, len(snapshot.get("nodes", [])) + 1)
            if node_id not in existing_node_ids:
                snapshot["nodes"].append(updated_node)
                existing_node_ids.add(node_id)
                continue
            snapshot["nodes"] = [
                (
                    {
                        **node,
                        **updated_node,
                        "position": node.get("position") or updated_node.get("position"),
                    }
                    if isinstance(node, dict) and node.get("id") == node_id
                    else node
                )
                for node in snapshot.get("nodes", [])
            ]
            continue

    for index, node_id in enumerate(accept_result.get("accepted_node_ids", []), start=1):
        if node_id in existing_node_ids:
            continue
        node = node_lookup.get(node_id)
        if node:
            snapshot["nodes"].append(_react_node_from_graph_node(node, index))
            existing_node_ids.add(node_id)
    for edge_id in accept_result.get("accepted_edge_ids", []):
        if edge_id in existing_edge_ids:
            continue
        edge = edge_lookup.get(edge_id)
        if edge:
            snapshot["edges"].append(_react_edge_from_graph_edge(edge))
            existing_edge_ids.add(edge_id)
    snapshot["viewport"] = _accepted_graph_focus_viewport(
        snapshot,
        accept_result.get("accepted_node_ids", []),
    )
    return snapshot


def _draft_revision_from_request(
    session: dict,
    graph: dict,
    request: dict[str, Any],
) -> dict:
    if (
        isinstance(request.get("draft_nodes"), list)
        or isinstance(request.get("draft_items"), list)
        or isinstance(request.get("generated_artifacts"), list)
    ):
        return build_ai_draft_revision(
            session=session,
            prompt=request.get("prompt") or request.get("custom_prompt") or "",
            draft_nodes=request.get("draft_nodes") or [],
            draft_edges=request.get("draft_edges") or [],
            draft_annotations=request.get("draft_annotations") or [],
            draft_items=request.get("draft_items"),
            generated_artifacts=request.get("generated_artifacts") or [],
            model=request.get("model") or session.get("selected_model", ""),
            metadata=request.get("metadata") if isinstance(request.get("metadata"), dict) else {},
        )

    preview = generate_ai_action_preview(
        graph,
        workspace_id=session["workspace_id"],
        role=request.get("role") or session.get("role") or "Custom",
        action=request.get("action") or "custom_prompt",
        scope=session.get("scope") or {"type": "workspace"},
        custom_prompt=_requested_prompt(request),
        created_by=request.get("created_by") or "user",
        model=request.get("model"),
    )
    return build_ai_draft_revision(
        session=session,
        prompt=request.get("prompt") or request.get("custom_prompt") or "",
        draft_nodes=preview.get("draft_nodes", []),
        draft_edges=preview.get("draft_edges", []),
        draft_annotations=preview.get("draft_annotations", []),
        generated_artifacts=preview.get("generated_artifacts", []),
        model=preview.get("metadata", {}).get("model", ""),
        validation_report=preview.get("validation_report"),
        metadata={
            "ai_action_id": preview.get("ai_action_id", ""),
            "model_reason": preview.get("metadata", {}).get("model_reason", ""),
            "preview_mode": preview.get("metadata", {}).get("preview_mode", ""),
            "output_shape": (request.get("metadata") or {}).get("output_shape", "")
            if isinstance(request.get("metadata"), dict)
            else "",
            "requested_visual": (request.get("metadata") or {}).get("requested_visual", "")
            if isinstance(request.get("metadata"), dict)
            else "",
        },
    )


def _has_client_supplied_draft(request: dict[str, Any]) -> bool:
    return (
        isinstance(request.get("draft_nodes"), list)
        or isinstance(request.get("draft_items"), list)
        or isinstance(request.get("generated_artifacts"), list)
    )


def _requested_desired_outputs(request: dict[str, Any]) -> list[str] | None:
    desired_outputs = request.get("desired_outputs")
    if isinstance(desired_outputs, list):
        return [str(output) for output in desired_outputs if str(output).strip()]
    return None


def _requested_source_chunks(request: dict[str, Any]) -> list[dict[str, Any]]:
    source_chunks = request.get("source_chunks")
    if not isinstance(source_chunks, list):
        return []
    return [chunk for chunk in source_chunks if isinstance(chunk, dict)]


def _requested_prompt(request: dict[str, Any]) -> str:
    prompt_with_brief = query_with_workspace_brief(
        request.get("prompt") or request.get("custom_prompt") or "",
        request.get("workspace_brief")
        if isinstance(request.get("workspace_brief"), dict)
        else (request.get("metadata") or {}).get("workspace_brief")
        if isinstance(request.get("metadata"), dict)
        else None,
    )
    return query_with_follow_up_memory(prompt_with_brief, request)


def _display_prompt(request: dict[str, Any]) -> str:
    return str(request.get("prompt") or request.get("custom_prompt") or "")


def _requested_model_policy(request: dict[str, Any]) -> str | None:
    policy = request.get("model_policy")
    if isinstance(policy, str):
        return policy
    if isinstance(policy, dict):
        value = policy.get("policy")
        return str(value) if value else None
    return None


def _requested_model(request: dict[str, Any]) -> str | None:
    model = request.get("model")
    if not isinstance(model, str):
        return None
    normalized = model.strip()
    if not normalized or normalized.lower() == "auto":
        return None
    return normalized


def _ai_draft_request_debug_summary(request: dict[str, Any], scope: dict[str, Any] | None = None) -> dict[str, Any]:
    prompt = _requested_prompt(request)
    desired_outputs = _requested_desired_outputs(request)
    source_chunks = _requested_source_chunks(request)
    source_refs = request.get("source_refs") if isinstance(request.get("source_refs"), list) else []
    normalized_scope = scope or (
        request.get("scope") if isinstance(request.get("scope"), dict) else {"type": "workspace"}
    )
    return {
        "role": request.get("role") or "Ask AI",
        "action": request.get("action") or request.get("intent") or "custom_prompt",
        "scope_type": normalized_scope.get("type", "workspace"),
        "desired_outputs": desired_outputs,
        "model": request.get("model") or "auto",
        "model_policy": _requested_model_policy(request) or "balanced",
        "prompt_chars": len(prompt or ""),
        "source_chunks": len(source_chunks),
        "source_refs": len(source_refs),
        "client_supplied_draft": _has_client_supplied_draft(request),
    }


def _ai_draft_session_debug_summary(session: dict[str, Any]) -> dict[str, Any]:
    revisions = session.get("revisions") if isinstance(session.get("revisions"), list) else []
    latest_revision = revisions[-1] if revisions and isinstance(revisions[-1], dict) else {}
    metadata = latest_revision.get("metadata") if isinstance(latest_revision.get("metadata"), dict) else {}
    return {
        "session_id": session.get("session_id", ""),
        "revision_id": latest_revision.get("revision_id", ""),
        "model": latest_revision.get("model") or session.get("selected_model") or metadata.get("actual_model") or "",
        "draft_nodes": len(latest_revision.get("draft_nodes", [])) if isinstance(latest_revision.get("draft_nodes"), list) else 0,
        "draft_edges": len(latest_revision.get("draft_edges", [])) if isinstance(latest_revision.get("draft_edges"), list) else 0,
        "draft_items": len(latest_revision.get("draft_items", [])) if isinstance(latest_revision.get("draft_items"), list) else 0,
        "generated_artifacts": len(latest_revision.get("generated_artifacts", [])) if isinstance(latest_revision.get("generated_artifacts"), list) else 0,
        "source_refs": len(session.get("source_refs", [])) if isinstance(session.get("source_refs"), list) else 0,
        "source_context_mode": metadata.get("source_context_mode", ""),
        "source_chunks_included": metadata.get("source_chunks_included", 0),
        "source_context_truncated": bool(metadata.get("source_context_truncated")),
        "input_tokens": metadata.get("input_tokens", 0),
        "output_tokens": metadata.get("output_tokens", 0),
        "total_tokens": metadata.get("total_tokens", 0),
    }


def _print_ai_draft_debug(
    *,
    flow_id: str,
    status_label: str,
    started_at: float,
    graph_elapsed_ms: float,
    generation_elapsed_ms: float,
    save_elapsed_ms: float,
    request_summary: dict[str, Any],
    session: dict[str, Any] | None = None,
) -> None:
    payload = {
        "flow_id": flow_id,
        "status": status_label,
        "elapsed_ms": round((time.perf_counter() - started_at) * 1000, 1),
        "graph_ms": round(graph_elapsed_ms, 1),
        "generation_ms": round(generation_elapsed_ms, 1),
        "save_ms": round(save_elapsed_ms, 1),
        "request": request_summary,
    }
    if session is not None:
        payload["session"] = _ai_draft_session_debug_summary(session)
    print(f"[ai-draft-session] {json.dumps(payload, separators=(',', ':'))}")


@app.post("/api/workspaces/{flow_id}/ai/draft-sessions")
def create_ai_draft_session(
    flow_id: str,
    request: dict[str, Any] | None = None,
):
    started_at = time.perf_counter()
    request = request or {}
    graph_started_at = time.perf_counter()
    graph = get_workspace_graph_or_404(flow_id)
    graph_elapsed_ms = (time.perf_counter() - graph_started_at) * 1000
    scope = normalize_ai_draft_scope(request.get("scope") if isinstance(request.get("scope"), dict) else {"type": "workspace"})
    request_summary = _ai_draft_request_debug_summary(request, scope)
    if not _has_client_supplied_draft(request):
        try:
            generation_started_at = time.perf_counter()
            generated_session = generate_ai_draft_session_with_provider(
                graph,
                workspace_id=flow_id,
                prompt=_requested_prompt(request),
                display_prompt=_display_prompt(request),
                scope=scope,
                role=request.get("role") or "Ask AI",
                model_policy=_requested_model_policy(request),
                model=_requested_model(request),
                desired_outputs=_requested_desired_outputs(request),
                source_chunks=_requested_source_chunks(request),
                metadata=request.get("metadata") if isinstance(request.get("metadata"), dict) else {},
            )
            generated_session["ai_action_run"] = build_ai_action_run(
                workspace_id=flow_id,
                scope=scope if scope.get("type") in {"workspace", "branch", "node"} else {"type": "workspace"},
                role=request.get("role") or "Ask AI",
                action=request.get("action") or "custom_prompt",
                custom_prompt=request.get("prompt") or request.get("custom_prompt"),
                input_source_refs=generated_session.get("source_refs") or graph.get("source_refs") or [],
                created_by=request.get("created_by") or "user",
                generated_node_ids=[
                    node.get("id", "")
                    for node in generated_session.get("revisions", [{}])[-1].get("draft_nodes", [])
                    if isinstance(node, dict)
                ],
            )
            metadata = request.get("metadata") if isinstance(request.get("metadata"), dict) else {}
            if metadata:
                generated_session.setdefault("metadata", {}).update(metadata)
            generation_elapsed_ms = (time.perf_counter() - generation_started_at) * 1000
            save_started_at = time.perf_counter()
            saved_session = save_ai_draft_session(validate_ai_draft_session(generated_session))
            save_elapsed_ms = (time.perf_counter() - save_started_at) * 1000
            _print_ai_draft_debug(
                flow_id=flow_id,
                status_label="generated",
                started_at=started_at,
                graph_elapsed_ms=graph_elapsed_ms,
                generation_elapsed_ms=generation_elapsed_ms,
                save_elapsed_ms=save_elapsed_ms,
                request_summary=request_summary,
                session=saved_session,
            )
            return saved_session
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc) from exc
        except GraphSchemaError as exc:
            print(
                "AI draft generation schema validation failed:",
                json.dumps(exc.errors, indent=2),
            )
            raise HTTPException(
                status_code=422,
                detail={
                    "message": "AI draft generation failed schema validation.",
                    "errors": exc.errors,
                },
            ) from exc
        except Exception as exc:
            traceback.print_exc()
            raise HTTPException(
                status_code=502,
                detail={
                    "message": "AI draft generation failed while calling the model provider.",
                    "error_type": exc.__class__.__name__,
                    "error": str(exc),
                },
            ) from exc

    action_run = build_ai_action_run(
        workspace_id=flow_id,
        scope=scope if scope.get("type") in {"workspace", "branch", "node"} else {"type": "workspace"},
        role=request.get("role") or "Custom",
        action=request.get("action") or "custom_prompt",
        custom_prompt=request.get("prompt") or request.get("custom_prompt"),
        input_source_refs=graph.get("source_refs") or [],
        created_by=request.get("created_by") or "user",
    )
    session = build_ai_draft_session(
        workspace_id=flow_id,
        prompt=request.get("prompt") or request.get("custom_prompt") or "",
        scope=scope,
        role=request.get("role") or "Custom",
        intent=request.get("intent") or request.get("action") or "custom_prompt",
        model_policy=request.get("model_policy") if isinstance(request.get("model_policy"), dict) else {},
        selected_model=request.get("model") or "",
        model_reason=request.get("model_reason") or "",
        source_refs=graph.get("source_refs") or [],
        ai_action_run=action_run,
        created_by=request.get("created_by") or "user",
        metadata=request.get("metadata") if isinstance(request.get("metadata"), dict) else {},
    )
    revision = _draft_revision_from_request(session, graph, request)
    session = append_ai_draft_revision(
        session,
        revision,
        prompt=request.get("prompt") or request.get("custom_prompt") or "",
        created_by=request.get("created_by") or "user",
    )
    if revision.get("model"):
        session["selected_model"] = revision["model"]
    if revision.get("metadata", {}).get("model_reason"):
        session["model_reason"] = revision["metadata"]["model_reason"]
    save_started_at = time.perf_counter()
    saved_session = save_ai_draft_session(session)
    save_elapsed_ms = (time.perf_counter() - save_started_at) * 1000
    _print_ai_draft_debug(
        flow_id=flow_id,
        status_label="client_supplied",
        started_at=started_at,
        graph_elapsed_ms=graph_elapsed_ms,
        generation_elapsed_ms=0,
        save_elapsed_ms=save_elapsed_ms,
        request_summary=request_summary,
        session=saved_session,
    )
    return saved_session


@app.post("/api/workspaces/{flow_id}/ai/node-message")
def create_node_info_message(
    flow_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    graph = get_workspace_graph_or_404(flow_id)
    scope = normalize_ai_draft_scope(request.get("scope") if isinstance(request.get("scope"), dict) else {"type": "workspace"})
    try:
        return generate_node_info_message_with_provider(
            graph,
            prompt=_requested_prompt(request),
            scope=scope,
            role=request.get("role") or "Ask AI",
            model_policy=_requested_model_policy(request),
            model=_requested_model(request),
            source_chunks=_requested_source_chunks(request),
            message_history=request.get("message_history") if isinstance(request.get("message_history"), list) else [],
            metadata=request.get("metadata") if isinstance(request.get("metadata"), dict) else {},
        )
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc
    except Exception as exc:
        traceback.print_exc()
        raise HTTPException(
            status_code=502,
            detail={
                "message": "AI node message failed while calling the model provider.",
                "error_type": exc.__class__.__name__,
                "error": str(exc),
            },
        ) from exc


@app.get("/api/workspaces/{flow_id}/ai/draft-sessions/{session_id}")
def get_ai_draft_session(flow_id: str, session_id: str):
    return get_ai_draft_session_or_404(flow_id, session_id)


@app.get("/api/workspaces/{flow_id}/ai/usage")
def get_workspace_ai_usage(flow_id: str):
    get_workspace_graph_or_404(flow_id)
    return summarize_ai_usage_for_workspace(flow_id)


@app.post("/api/workspaces/{flow_id}/ai/draft-sessions/{session_id}/revisions")
def create_ai_draft_revision(
    flow_id: str,
    session_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    session = get_ai_draft_session_or_404(flow_id, session_id)
    if session.get("status") != "drafting":
        raise HTTPException(status_code=409, detail="Only active draft sessions can be revised.")
    graph = get_workspace_graph_or_404(flow_id)
    if not _has_client_supplied_draft(request):
        try:
            session = revise_ai_draft_session_with_provider(
                session,
                graph,
                prompt=_requested_prompt(request),
                display_prompt=_display_prompt(request),
                model_policy=_requested_model_policy(request),
                model=_requested_model(request),
                desired_outputs=_requested_desired_outputs(request),
                source_chunks=_requested_source_chunks(request),
            )
            metadata = request.get("metadata") if isinstance(request.get("metadata"), dict) else {}
            if metadata:
                session.setdefault("metadata", {}).update(metadata)
            return save_ai_draft_session(session)
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc) from exc
        except GraphSchemaError as exc:
            raise HTTPException(
                status_code=422,
                detail={
                    "message": "AI draft revision failed schema validation.",
                    "errors": exc.errors,
                },
            ) from exc
        except Exception as exc:
            traceback.print_exc()
            raise HTTPException(
                status_code=502,
                detail={
                    "message": "AI draft revision failed while calling the model provider.",
                    "error_type": exc.__class__.__name__,
                    "error": str(exc),
                },
            ) from exc

    revision = _draft_revision_from_request(session, graph, request)
    session = append_ai_draft_revision(
        session,
        revision,
        prompt=request.get("prompt") or request.get("custom_prompt") or "",
        created_by=request.get("created_by") or "user",
    )
    if revision.get("model"):
        session["selected_model"] = revision["model"]
    if revision.get("metadata", {}).get("model_reason"):
        session["model_reason"] = revision["metadata"]["model_reason"]
    return save_ai_draft_session(session)


@app.post("/api/workspaces/{flow_id}/ai/draft-sessions/{session_id}/sources")
def add_ai_draft_session_source(
    flow_id: str,
    session_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    session = get_ai_draft_session_or_404(flow_id, session_id)
    if session.get("status") != "drafting":
        raise HTTPException(status_code=409, detail="Only active draft sessions can add sources.")
    source_chunks = request.get("source_chunks")
    if not isinstance(source_chunks, list):
        source_chunks = []
    graph = get_workspace_graph_or_404(flow_id)
    if request.get("source_id") and not source_chunks:
        source_id = str(request.get("source_id"))
        source_library = graph.get("source_library", {}) if isinstance(graph.get("source_library"), dict) else {}
        for document in source_library.get("documents", []) if isinstance(source_library.get("documents"), list) else []:
            if not isinstance(document, dict):
                continue
            if str(document.get("id") or document.get("document_id") or "") != source_id:
                continue
            source_chunks = [
                {
                    **chunk,
                    "document_id": source_id,
                    "source_ref": {
                        "document_id": source_id,
                        "chunk_id": chunk.get("id", ""),
                        "page": chunk.get("page"),
                        "section": chunk.get("heading", ""),
                        "quote_snippet": chunk.get("snippet", ""),
                        "confidence": "medium",
                    },
                }
                for chunk in document.get("chunks", [])
                if isinstance(chunk, dict)
            ]
            break
    try:
        session = add_source_to_ai_draft_session(
            session,
            graph,
            source_chunks=source_chunks,
            prompt=request.get("prompt"),
            model_policy=request.get("model_policy"),
            model=request.get("model"),
        )
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI draft source reconciliation failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc
    return save_ai_draft_session(session)


@app.post("/api/workspaces/{flow_id}/ai/draft-sessions/{session_id}/discard")
def discard_ai_draft_session_endpoint(
    flow_id: str,
    session_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    session = get_ai_draft_session_or_404(flow_id, session_id)
    session = discard_ai_draft_session(
        session,
        discarded_by=request.get("discarded_by") or request.get("created_by") or "user",
    )
    return save_ai_draft_session(session)


@app.post("/api/workspaces/{flow_id}/ai/draft-sessions/{session_id}/accept")
def accept_ai_draft_session_endpoint(
    flow_id: str,
    session_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    session = get_ai_draft_session_or_404(flow_id, session_id)
    if session.get("status") != "drafting":
        raise HTTPException(status_code=409, detail="Only active draft sessions can be accepted.")
    flow = get_workspace_flow_or_404(flow_id)
    graph = build_workspace_graph(flow, source_components=get_source_components(flow_id))
    previous_flow_json = flow.get("flow_json", "")
    try:
        accepted_graph, session, accept_result = accept_ai_draft_revision(
            graph,
            session,
            revision_id=request.get("revision_id"),
            accept_mode=request.get("mode") or request.get("accept_mode") or "append",
            selected_item_ids=request.get("selected_item_ids") if isinstance(request.get("selected_item_ids"), list) else [],
            accepted_by=request.get("accepted_by") or request.get("created_by") or "user",
        )
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI draft accept failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc
    accept_result.setdefault("metadata", {})["undo_snapshot"] = previous_flow_json
    if session.get("accept_history"):
        session["accept_history"][-1].setdefault("metadata", {})["undo_snapshot"] = previous_flow_json
    snapshot = _append_accepted_graph_to_flow_snapshot(flow, accept_result, accepted_graph)
    _persist_flow_snapshot(flow_id, snapshot)
    save_ai_draft_session(session)
    response = {
        **accept_result,
        "graph": snapshot,
        "session": session,
        "accept_result": accept_result,
    }
    return response


@app.post("/api/workspaces/{flow_id}/ai/actions/preview")
def preview_ai_action(
    flow_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    scope = request.get("scope") if isinstance(request.get("scope"), dict) else {}
    if request.get("node_id") and not scope:
        scope = {"type": "node", "node_id": request["node_id"]}
    if request.get("branch_node_id") and not scope:
        scope = {"type": "branch", "node_id": request["branch_node_id"]}
    if not scope:
        scope = {"type": "workspace"}

    graph = get_workspace_graph_or_404(flow_id)
    try:
        return generate_ai_action_preview(
            graph,
            workspace_id=flow_id,
            role=request.get("role") or request.get("role_id") or "",
            action=request.get("action") or "",
            scope=scope,
            custom_prompt=request.get("custom_prompt"),
            created_by=request.get("created_by") or "user",
            model=request.get("model"),
        )
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI action preview failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc


@app.post("/api/workspaces/{flow_id}/ai/actions/node/{node_id}/preview")
def preview_node_ai_action(
    flow_id: str,
    node_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    request["scope"] = {"type": "node", "node_id": node_id}
    return preview_ai_action(flow_id, request)


@app.post("/api/workspaces/{flow_id}/ai/actions/branch/{node_id}/preview")
def preview_branch_ai_action(
    flow_id: str,
    node_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    request["scope"] = {"type": "branch", "node_id": node_id}
    return preview_ai_action(flow_id, request)


@app.post("/api/workspaces/{flow_id}/ai/actions/workspace/preview")
def preview_workspace_ai_action(
    flow_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    request["scope"] = {"type": "workspace"}
    return preview_ai_action(flow_id, request)


@app.post("/api/workspaces/{flow_id}/ai/helpers/{helper_id}/preview")
def preview_ai_helper(
    flow_id: str,
    helper_id: str,
    request: dict[str, Any] | None = None,
):
    request = request or {}
    helper_id = helper_id.replace("-", "_")
    action = request.get("action") or ""
    scope = request.get("scope") if isinstance(request.get("scope"), dict) else {}
    if request.get("branch_node_id") and not scope:
        scope = {"type": "branch", "node_id": request["branch_node_id"]}
    if not scope:
        scope = {"type": "workspace"}

    if scope.get("type") == "branch" and scope.get("node_id"):
        graph = get_workspace_branch_or_404(flow_id, scope["node_id"])
    else:
        graph = get_workspace_graph_or_404(flow_id)

    try:
        return generate_helper_preview(
            helper_id,
            action,
            graph,
            scope=scope,
            use_ai=bool(request.get("use_ai", True)),
            allow_deterministic_fallback=bool(
                request.get("allow_deterministic_fallback", True)
            ),
            model=request.get("model"),
        )
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI helper preview failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc


@app.post("/api/workspaces/{flow_id}/ai/source-librarian/preview")
def preview_source_librarian(flow_id: str, request: dict[str, Any] | None = None):
    request = request or {}
    action = request.get("action") or "source_repair"
    scope = request.get("scope") if isinstance(request.get("scope"), dict) else {}
    if request.get("branch_node_id") and not scope:
        scope = {"type": "branch", "node_id": request["branch_node_id"]}
    if not scope:
        scope = {"type": "workspace"}

    if scope.get("type") == "branch" and scope.get("node_id"):
        graph = get_workspace_branch_or_404(flow_id, scope["node_id"])
    else:
        graph = get_workspace_graph_or_404(flow_id)

    try:
        return generate_source_librarian_preview(
            graph,
            action=action,
            scope=scope,
            use_ai=bool(request.get("use_ai", True)),
            allow_deterministic_fallback=bool(
                request.get("allow_deterministic_fallback", True)
            ),
            model=request.get("model"),
        )
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI helper preview failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc


@app.post("/api/workspaces/{flow_id}/ai/reviewer/preview")
def preview_reviewer(flow_id: str, request: dict[str, Any] | None = None):
    request = request or {}
    action = request.get("action") or "missing_information"
    scope = request.get("scope") if isinstance(request.get("scope"), dict) else {}
    if request.get("branch_node_id") and not scope:
        scope = {"type": "branch", "node_id": request["branch_node_id"]}
    if not scope:
        scope = {"type": "workspace"}

    if scope.get("type") == "branch" and scope.get("node_id"):
        graph = get_workspace_branch_or_404(flow_id, scope["node_id"])
    else:
        graph = get_workspace_graph_or_404(flow_id)

    try:
        return generate_helper_preview(
            "reviewer",
            action,
            graph,
            scope=scope,
            use_ai=bool(request.get("use_ai", True)),
            allow_deterministic_fallback=bool(
                request.get("allow_deterministic_fallback", True)
            ),
            model=request.get("model"),
        )
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI helper preview failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc


@app.post("/api/workspaces/{flow_id}/ai/project-planner/preview")
def preview_project_planner(flow_id: str, request: dict[str, Any] | None = None):
    request = request or {}
    action = request.get("action") or "task_projection"
    scope = request.get("scope") if isinstance(request.get("scope"), dict) else {}
    if request.get("branch_node_id") and not scope:
        scope = {"type": "branch", "node_id": request["branch_node_id"]}
    if not scope:
        scope = {"type": "workspace"}

    if scope.get("type") == "branch" and scope.get("node_id"):
        graph = get_workspace_branch_or_404(flow_id, scope["node_id"])
    else:
        graph = get_workspace_graph_or_404(flow_id)

    try:
        return generate_helper_preview(
            "project_planner",
            action,
            graph,
            scope=scope,
            use_ai=bool(request.get("use_ai", True)),
            allow_deterministic_fallback=bool(
                request.get("allow_deterministic_fallback", True)
            ),
            model=request.get("model"),
        )
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI helper preview failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc


@app.post("/api/workspaces/{flow_id}/ai/integration-operator/preview")
def preview_integration_operator(flow_id: str, request: dict[str, Any] | None = None):
    request = request or {}
    action = request.get("action") or "handoff_readiness"
    scope = request.get("scope") if isinstance(request.get("scope"), dict) else {}
    if request.get("branch_node_id") and not scope:
        scope = {"type": "branch", "node_id": request["branch_node_id"]}
    if not scope:
        scope = {"type": "workspace"}

    if scope.get("type") == "branch" and scope.get("node_id"):
        graph = get_workspace_branch_or_404(flow_id, scope["node_id"])
    else:
        graph = get_workspace_graph_or_404(flow_id)

    try:
        return generate_helper_preview(
            "integration_operator",
            action,
            graph,
            scope=scope,
            use_ai=bool(request.get("use_ai", True)),
            allow_deterministic_fallback=bool(
                request.get("allow_deterministic_fallback", True)
            ),
            model=request.get("model"),
        )
    except MissingConfigurationError as exc:
        raise configuration_http_error(exc) from exc
    except GraphSchemaError as exc:
        raise HTTPException(
            status_code=422,
            detail={
                "message": "AI helper preview failed schema validation.",
                "errors": exc.errors,
            },
        ) from exc
    except RuntimeError as exc:
        raise HTTPException(status_code=502, detail=str(exc)) from exc


@app.post("/api/workspaces/{flow_id}/export/miro")
def preview_workspace_miro_export(flow_id: str):
    graph = get_workspace_graph_or_404(flow_id)
    return export_branch_to_miro_payload(
        graph["nodes"],
        graph["edges"],
        graph["workspace"],
        target="workspace_board_preview",
    )


@app.post("/api/workspaces/{flow_id}/export/miro/board")
def export_workspace_to_miro_board(
    flow_id: str,
    board_id: str = Query(...),
    dry_run: bool = Query(True),
):
    if not dry_run:
        try:
            require_settings("miro_api_token")
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc)

    flow = get_workspace_flow_or_404(flow_id)
    graph = build_workspace_graph(flow)
    payload = export_branch_to_miro_payload(
        graph["nodes"],
        graph["edges"],
        graph["workspace"],
        target="workspace_board",
    )
    client = MiroClient(
        token=get_setting("miro_api_token") or "",
        base_url=get_setting("miro_api_base_url") or "https://api.miro.com/v2",
    )
    result = client.export_frame_payload(board_id, payload, dry_run=dry_run)

    if dry_run:
        return result

    pushed_at = datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z"
    refs_by_node_id = miro_item_refs_from_result(board_id, result, pushed_at)
    updated_flow_json = apply_miro_external_refs_to_flow_json(
        flow.get("flow_json", ""),
        refs_by_node_id,
    )
    flow_collection.update_one(
        {"_id": ObjectId(flow_id)},
        {"$set": {"flow_json": updated_flow_json}},
    )

    return {
        **result,
        "external_refs": refs_by_node_id,
        "persisted": bool(refs_by_node_id),
    }


@app.post("/api/workspaces/{flow_id}/export/miro/sme-review")
def export_workspace_to_miro_sme_review_board(
    flow_id: str,
    board_id: str = Query(...),
    dry_run: bool = Query(True),
):
    if not dry_run:
        try:
            require_settings("miro_api_token")
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc)

    flow = get_workspace_flow_or_404(flow_id)
    graph = build_workspace_graph(flow)
    payload = export_sme_review_board_payload(graph)
    client = MiroClient(
        token=get_setting("miro_api_token") or "",
        base_url=get_setting("miro_api_base_url") or "https://api.miro.com/v2",
    )
    result = client.export_frame_payload(board_id, payload, dry_run=dry_run)

    if dry_run:
        return result

    pushed_at = utc_timestamp()
    refs_by_node_id = miro_item_refs_from_result(board_id, result, pushed_at)
    updated_flow_json = apply_miro_external_refs_to_flow_json(
        flow.get("flow_json", ""),
        refs_by_node_id,
    )
    flow_collection.update_one(
        {"_id": ObjectId(flow_id)},
        {"$set": {"flow_json": updated_flow_json}},
    )

    return {
        **result,
        "external_refs": refs_by_node_id,
        "persisted": bool(refs_by_node_id),
    }


@app.post("/api/workspaces/{flow_id}/export/miro/native-mindmap")
def preview_workspace_miro_native_mindmap_export(
    flow_id: str,
    board_id: str = Query(...),
):
    graph = get_workspace_graph_or_404(flow_id)
    payload = export_native_mindmap_payload(graph)
    client = MiroClient(
        token=os.getenv("miro_api_token", ""),
        base_url=os.getenv("miro_api_base_url", "https://api.miro.com/v2"),
    )
    return client.export_native_mindmap_payload(board_id, payload, dry_run=True)


@app.post("/api/workspaces/{flow_id}/branches/{node_id}/export/miro")
def preview_branch_miro_export(flow_id: str, node_id: str):
    branch = get_workspace_branch_or_404(flow_id, node_id)
    return export_branch_to_miro_payload(
        branch["nodes"],
        branch["edges"],
        branch["workspace"],
        target="selected_branch_frame",
    )


@app.post("/api/workspaces/{flow_id}/branches/{node_id}/export/miro/frame")
def export_branch_to_miro_frame(
    flow_id: str,
    node_id: str,
    board_id: str = Query(...),
    dry_run: bool = Query(True),
):
    if not dry_run:
        try:
            require_settings("miro_api_token")
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc)

    flow = get_workspace_flow_or_404(flow_id)
    graph = build_workspace_graph(flow)
    if not any(node["id"] == node_id for node in graph["nodes"]):
        raise HTTPException(status_code=404, detail="Branch root node not found.")

    branch = select_branch(graph, node_id)
    payload = export_branch_to_miro_payload(
        branch["nodes"],
        branch["edges"],
        branch["workspace"],
        target="selected_branch_frame",
    )
    client = MiroClient(
        token=get_setting("miro_api_token") or "",
        base_url=get_setting("miro_api_base_url") or "https://api.miro.com/v2",
    )
    result = client.export_frame_payload(board_id, payload, dry_run=dry_run)

    if dry_run:
        return result

    pushed_at = datetime.datetime.utcnow().replace(microsecond=0).isoformat() + "Z"
    refs_by_node_id = miro_item_refs_from_result(board_id, result, pushed_at)
    updated_flow_json = apply_miro_external_refs_to_flow_json(
        flow.get("flow_json", ""),
        refs_by_node_id,
    )
    flow_collection.update_one(
        {"_id": ObjectId(flow_id)},
        {"$set": {"flow_json": updated_flow_json}},
    )

    return {
        **result,
        "external_refs": refs_by_node_id,
        "persisted": bool(refs_by_node_id),
    }


@app.post("/api/workspaces/{flow_id}/export/monday")
def preview_workspace_monday_export(
    flow_id: str,
    confirmed: bool = Query(False),
    board_id: str = Query(""),
    group_id: str = Query(""),
    template_id: str = Query(""),
):
    graph = get_workspace_graph_or_404(flow_id)
    return export_tasks_to_monday_payload(
        monday_task_nodes_from_graph(graph),
        graph["workspace"],
        confirmed=confirmed,
        board_id=board_id,
        group_id=group_id,
        scope="workspace",
        created_at=utc_timestamp(),
        template_id=template_id,
    )


@app.post("/api/integrations/monday/preflight/existing-group")
def preflight_monday_existing_group(
    board_id: str = Query(...),
    group_id: str = Query(...),
    dry_run: bool = Query(True),
    template_id: str = Query(""),
):
    if not dry_run:
        try:
            require_settings("monday_api_token")
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc)

    client = MondayClient(
        token=get_setting("monday_api_token") or "",
        base_url=get_setting("monday_api_base_url") or "https://api.monday.com/v2",
    )
    return client.preflight_existing_group(
        board_id,
        group_id,
        template_id=template_id,
        dry_run=dry_run,
    )


@app.post("/api/workspaces/{flow_id}/export/monday/existing-group")
def export_workspace_to_monday_existing_group(
    flow_id: str,
    board_id: str = Query(...),
    group_id: str = Query(...),
    dry_run: bool = Query(True),
    confirmed: bool = Query(False),
    template_id: str = Query(""),
):
    if not dry_run and not confirmed:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Confirm the monday export before creating items.",
        )
    if not dry_run:
        try:
            require_settings("monday_api_token")
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc)

    flow = get_workspace_flow_or_404(flow_id)
    graph = build_workspace_graph(flow)
    payload = export_tasks_to_monday_payload(
        monday_task_nodes_from_graph(graph),
        graph["workspace"],
        confirmed=confirmed,
        board_id=board_id,
        group_id=group_id,
        scope="workspace",
        created_at=utc_timestamp(),
        template_id=template_id,
    )
    client = MondayClient(
        token=get_setting("monday_api_token") or "",
        base_url=get_setting("monday_api_base_url") or "https://api.monday.com/v2",
    )
    preflight_result = None
    if not dry_run:
        preflight_result = client.preflight_existing_group(
            board_id,
            group_id,
            template_id=template_id,
            dry_run=False,
        )
        if not preflight_result.get("preflight", {}).get("ok"):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail={
                    "message": "monday target preflight failed.",
                    "preflight": preflight_result.get("preflight", {}),
                },
            )
    result = client.export_existing_group_items(payload, dry_run=dry_run)

    if dry_run:
        return result

    pushed_at = utc_timestamp()
    refs_by_node_id = monday_item_refs_from_result(
        board_id,
        group_id,
        result,
        pushed_at,
    )
    updated_flow_json = apply_monday_external_refs_to_flow_json(
        flow.get("flow_json", ""),
        refs_by_node_id,
    )
    flow_collection.update_one(
        {"_id": ObjectId(flow_id)},
        {"$set": {"flow_json": updated_flow_json}},
    )

    return {
        **result,
        "preflight": preflight_result.get("preflight", {}) if preflight_result else {},
        "external_refs": refs_by_node_id,
        "persisted": bool(refs_by_node_id),
    }


@app.post("/api/workspaces/{flow_id}/branches/{node_id}/export/monday")
def preview_branch_monday_export(
    flow_id: str,
    node_id: str,
    confirmed: bool = Query(False),
    board_id: str = Query(""),
    group_id: str = Query(""),
    template_id: str = Query(""),
):
    branch = get_workspace_branch_or_404(flow_id, node_id)
    return export_tasks_to_monday_payload(
        monday_task_nodes_from_graph(branch),
        branch["workspace"],
        confirmed=confirmed,
        board_id=board_id,
        group_id=group_id,
        scope="branch",
        root_node_id=node_id,
        created_at=utc_timestamp(),
        template_id=template_id,
    )


@app.post("/api/workspaces/{flow_id}/branches/{node_id}/export/monday/existing-group")
def export_branch_to_monday_existing_group(
    flow_id: str,
    node_id: str,
    board_id: str = Query(...),
    group_id: str = Query(...),
    dry_run: bool = Query(True),
    confirmed: bool = Query(False),
    template_id: str = Query(""),
):
    if not dry_run and not confirmed:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Confirm the monday export before creating items.",
        )
    if not dry_run:
        try:
            require_settings("monday_api_token")
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc)

    flow = get_workspace_flow_or_404(flow_id)
    graph = build_workspace_graph(flow)
    if not any(node["id"] == node_id for node in graph["nodes"]):
        raise HTTPException(status_code=404, detail="Branch root node not found.")

    branch = select_branch(graph, node_id)
    payload = export_tasks_to_monday_payload(
        monday_task_nodes_from_graph(branch),
        branch["workspace"],
        confirmed=confirmed,
        board_id=board_id,
        group_id=group_id,
        scope="branch",
        root_node_id=node_id,
        created_at=utc_timestamp(),
        template_id=template_id,
    )
    client = MondayClient(
        token=get_setting("monday_api_token") or "",
        base_url=get_setting("monday_api_base_url") or "https://api.monday.com/v2",
    )
    preflight_result = None
    if not dry_run:
        preflight_result = client.preflight_existing_group(
            board_id,
            group_id,
            template_id=template_id,
            dry_run=False,
        )
        if not preflight_result.get("preflight", {}).get("ok"):
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail={
                    "message": "monday target preflight failed.",
                    "preflight": preflight_result.get("preflight", {}),
                },
            )
    result = client.export_existing_group_items(payload, dry_run=dry_run)

    if dry_run:
        return result

    pushed_at = utc_timestamp()
    refs_by_node_id = monday_item_refs_from_result(
        board_id,
        group_id,
        result,
        pushed_at,
    )
    updated_flow_json = apply_monday_external_refs_to_flow_json(
        flow.get("flow_json", ""),
        refs_by_node_id,
    )
    flow_collection.update_one(
        {"_id": ObjectId(flow_id)},
        {"$set": {"flow_json": updated_flow_json}},
    )

    return {
        **result,
        "preflight": preflight_result.get("preflight", {}) if preflight_result else {},
        "external_refs": refs_by_node_id,
        "persisted": bool(refs_by_node_id),
    }


@app.post("/api/workspaces/{flow_id}/sync/monday/status")
def pull_monday_status_to_workspace(
    flow_id: str,
    dry_run: bool = Query(True),
    apply: bool = Query(False),
):
    if not dry_run:
        try:
            require_settings("monday_api_token")
        except MissingConfigurationError as exc:
            raise configuration_http_error(exc)

    flow = get_workspace_flow_or_404(flow_id)
    refs_by_node_id = monday_refs_from_flow_json(flow.get("flow_json", ""))
    client = MondayClient(
        token=get_setting("monday_api_token") or "",
        base_url=get_setting("monday_api_base_url") or "https://api.monday.com/v2",
    )
    result = client.pull_item_statuses(refs_by_node_id, dry_run=dry_run)

    if dry_run:
        return {
            **result,
            "tracked_node_count": len(refs_by_node_id),
        }

    pulled_at = utc_timestamp()
    status_projections = monday_status_projections_from_result(
        result,
        refs_by_node_id,
        pulled_at,
    )
    if apply:
        updated_flow_json = apply_monday_status_projection_to_flow_json(
            flow.get("flow_json", ""),
            status_projections,
        )
        flow_collection.update_one(
            {"_id": ObjectId(flow_id)},
            {"$set": {"flow_json": updated_flow_json}},
        )

    return {
        **result,
        "status_projections": status_projections,
        "status_updates": status_projections,
        "applied": apply and bool(status_projections),
        "applied_as": "external_status_projection",
    }


def get_summary_from_openai(
    file: UploadFile,
    flow_id: str,
    flow_type: str,
    operation_id: str | None = None,
    intake_role: str | None = None,
    intake_prompt: str | None = None,
    intake_model: str | None = None,
):
    try:
        update_operation_progress(
            operation_id,
            phase="checking_settings",
            message="Checking AI settings",
            detail="Verifying OpenAI configuration before upload.",
            progress=18,
        )
        require_settings("openai_api_key")
    except MissingConfigurationError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Missing AI settings",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise configuration_http_error(exc) from exc

    update_operation_progress(
        operation_id,
        phase="extracting",
        message="Reading source document",
        detail="Validating the source and preparing document metadata.",
        progress=28,
    )
    source_context = prepare_ai_intake_upload(file, flow_id)
    source_document = source_context["source_document"]
    file_bytes = source_context["file_bytes"]
    file_extension = source_document["type"]
    existing_component = source_context.get("existing_component")

    if len(file_bytes) == 0:
        raise ValueError("The uploaded file is actually empty!")

    if source_context.get("document_chunks"):
        intake_instruction = build_source_intake_instruction(intake_role, intake_prompt)
        update_operation_progress(
            operation_id,
            phase="ai_reading",
            message="AI is reading the source document",
            detail="The model is summarizing the document text prepared from your upload.",
            progress=62,
        )
        try:
            summary_text, ai_metadata = generate_document_summary(
                file_name=source_document["filename"],
                source_type=file_extension,
                chunks=source_context["document_chunks"],
                role_instruction=intake_instruction,
                model=intake_model,
            )
            processing_type = "responses"
        except HTTPException as exc:
            if exc.status_code != status.HTTP_504_GATEWAY_TIMEOUT:
                raise
            summary_text = fallback_source_summary(source_context)
            ai_metadata = {
                "provider": "deterministic_timeout_fallback",
                "model": intake_model or "",
                "reason": "OpenAI source summary timed out; saved parsed source context instead.",
            }
            processing_type = "responses_timeout_fallback"
        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "file_id": "",
            "assistant_id": "",
            "vector_store_id": "",
            "size": len(file_bytes),
            "type": file_extension,
            "processing_type": processing_type,
            "summary": summary_text,
            "ai_provider": ai_metadata,
            **source_metadata_fields(source_context),
        }
        if existing_component:
            component_id = existing_component["_id"]
            component_collection.update_one(
                {"_id": component_id},
                {
                    "$set": {
                        "summary": summary_text,
                        "processing_type": processing_type,
                        "ai_provider": ai_metadata,
                    }
                },
            )
        else:
            component_id = component_collection.insert_one(component_metadata).inserted_id
        update_operation_progress(
            operation_id,
            phase="complete",
            message="Source context is ready",
            detail=(
                "The parsed source was saved with a timeout fallback."
                if processing_type == "responses_timeout_fallback"
                else "The source component was saved to the workspace."
            ),
            progress=100,
            status_value="completed",
        )
        return {
            "component_id": str(component_id),
            "type": file_extension,
            flow_type: flow_type,
            "processing_type": processing_type,
            "timeout_fallback": processing_type == "responses_timeout_fallback",
            **source_metadata_fields(source_context),
        }

    require_legacy_assistants_fallback(file_extension, purpose="source summary")
    assistant_model = resolve_assistants_model(intake_model)
    intake_instruction = build_source_intake_instruction(intake_role, intake_prompt)

    update_operation_progress(
        operation_id,
        phase="uploading_to_ai",
        message="Uploading source to AI workspace",
        detail="Creating an AI file-search workspace for the source.",
        progress=42,
    )
    assistant = openai.beta.assistants.create(
        name="Summarize agent",
        instructions="Your task is to only summarize the document",
        model=assistant_model,
        tools=[{"type": "file_search"}],
    )
    vector_store = openai.beta.vector_stores.create(name=f"{file_extension}_{flow_id}")

    assistant = openai.beta.assistants.update(
        assistant_id=assistant.id,
        tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
    )

    if file_extension == "pdf":
        mime_type = "application/pdf"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "csv":
        mime_type = "application/json"
        df = pd.read_csv(BytesIO(file_bytes))
        json_data = df.to_dict(orient="records")
        json_str = json.dumps(json_data)
        messages_file = openai.files.create(
            file=("data.json", json_str, mime_type), purpose="assistants"
        )
    elif file_extension == "txt":
        mime_type = "text/plain"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "docx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "html":
        mime_type = "text/html"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "md":
        mime_type = "text/markdown"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "pptx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    else:
        raise ValueError("Unsupported file type")

    update_operation_progress(
        operation_id,
        phase="ai_reading",
        message="AI is reading the source",
        detail="The model is summarizing the uploaded document.",
        progress=62,
    )
    # Create thread for summarization task
    thread = openai.beta.threads.create(
        messages=[
            {
                "role": "user",
                "content": f"Generate a concise summary of the following document{intake_instruction}",
                "attachments": [
                    {"file_id": messages_file.id, "tools": [{"type": "file_search"}]}
                ],
            }
        ]
    )

    # Run the assistant and fetch the results
    run = openai.beta.threads.runs.create_and_poll(
        thread_id=thread.id, assistant_id=assistant.id
    )

    update_operation_progress(
        operation_id,
        phase="saving",
        message="Saving source metadata",
        detail="Persisting the source component and summary.",
        progress=86,
    )
    print(thread)
    print(run)

    messages = list(
        openai.beta.threads.messages.list(thread_id=thread.id, run_id=run.id)
    )
    print(messages)
    message_content = messages[0].content[0].text
    annotations = message_content.annotations

    # Annotate the summary content if necessary
    for index, annotation in enumerate(annotations):
        message_content.value = message_content.value.replace(
            annotation.text, f"[{index}]"
        )

    # Insert metadata into the database
    component_metadata = {
        "flow_id": ObjectId(flow_id),
        "file_id": messages_file.id,
        "assistant_id": assistant.id,
        "vector_store_id": vector_store.id,
        "size": len(file_bytes),
        "type": file_extension,
        "processing_type": "legacy_assistants",
        "ai_provider": {
            "provider": "assistants_legacy_fallback",
            "model": assistant_model,
            "reason": "Prepared document sections were unavailable, so source intake used the temporary Assistants file-search fallback.",
        },
        "summary": message_content.value,
        **source_metadata_fields(source_context),
    }

    # Store the document in MongoDB
    component_id = component_collection.insert_one(component_metadata).inserted_id

    update_operation_progress(
        operation_id,
        phase="complete",
        message="Source summary is ready",
        detail="The source component was saved to the workspace.",
        progress=100,
        status_value="completed",
    )
    # Return the component ID with the type
    return {"component_id": str(component_id), "type": file_extension, flow_type: flow_type}

def openai_mindmap_generator(
    file: UploadFile,
    flow_id: str,
    flow_type: str,
    operation_id: str | None = None,
    intake_role: str | None = None,
    intake_prompt: str | None = None,
    intake_model: str | None = None,
):
    try:
        update_operation_progress(
            operation_id,
            phase="checking_settings",
            message="Checking AI settings",
            detail="Verifying OpenAI configuration before deriving the workspace.",
            progress=18,
        )
        require_settings("openai_api_key")
    except MissingConfigurationError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Missing AI settings",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise configuration_http_error(exc) from exc

    update_operation_progress(
        operation_id,
        phase="extracting",
        message="Reading source document",
        detail="Validating the source and preparing document metadata.",
        progress=28,
    )
    source_context = prepare_ai_intake_upload(file, flow_id)
    source_document = source_context["source_document"]
    file_bytes = source_context["file_bytes"]
    file_extension = source_document["type"]
    existing_component = source_context.get("existing_component")

    if len(file_bytes) == 0:
        raise ValueError("The uploaded file is actually empty!")

    if source_context.get("document_chunks"):
        intake_role_label = resolve_source_intake_role(intake_role)
        intake_prompt_text = clean_source_intake_value(intake_prompt)
        intake_instruction = build_source_intake_instruction(
            intake_role_label,
            intake_prompt_text,
        )
        update_operation_progress(
            operation_id,
            phase="ai_deriving",
            message="AI is deriving workspace nodes",
            detail="The model is turning the source document into a reviewable TraceSpace draft.",
            progress=64,
        )
        try:
            response_json, ai_metadata = generate_document_mindmap(
                file_name=source_document["filename"],
                source_type=file_extension,
                flow_id=flow_id,
                chunks=source_context["document_chunks"],
                role_instruction=intake_instruction,
                model=intake_model,
            )
            response_json = ground_mindmap_with_source_refs(response_json, source_context)
            processing_type = "responses"
        except HTTPException as exc:
            if exc.status_code != status.HTTP_504_GATEWAY_TIMEOUT:
                raise
            response_json = fallback_source_mindmap(source_context, flow_id)
            ai_metadata = {
                "provider": "deterministic_timeout_fallback",
                "model": intake_model or "",
                "reason": "OpenAI graph derivation timed out; saved parsed source context with a reviewable starter graph.",
            }
            processing_type = "responses_timeout_fallback"
        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "file_id": "",
            "assistant_id": "",
            "vector_store_id": "",
            "size": len(file_bytes),
            "type": file_extension,
            "processing_type": processing_type,
            "mindmap_json": response_json,
            "ai_provider": ai_metadata,
            **source_metadata_fields(source_context),
        }
        if existing_component:
            component_id = existing_component["_id"]
            component_collection.update_one(
                {"_id": component_id},
                {
                    "$set": {
                        "mindmap_json": response_json,
                        "processing_type": processing_type,
                        "ai_provider": ai_metadata,
                    }
                },
            )
        else:
            component_id = component_collection.insert_one(component_metadata).inserted_id
        if processing_type == "responses_timeout_fallback":
            response_json = fallback_source_mindmap(source_context, flow_id, component_id)
            component_collection.update_one(
                {"_id": component_id},
                {"$set": {"mindmap_json": response_json}},
            )
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        update_operation_progress(
            operation_id,
            phase="complete",
            message="Source context is ready",
            detail=(
                "AI derivation timed out, so TraceSpace saved the parsed source with a reviewable starter graph."
                if processing_type == "responses_timeout_fallback"
                else "The generated mind map was saved to the workspace."
            ),
            progress=100,
            status_value="completed",
        )
        return {
            "flow_id": flow_id,
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": file_extension,
            "mindmap_json": response_json,
            "flow_type": flow_type,
            "processing_type": processing_type,
            "timeout_fallback": processing_type == "responses_timeout_fallback",
            **source_metadata_fields(source_context),
        }

    require_legacy_assistants_fallback(file_extension, purpose="graph generation")
    assistant_model = resolve_assistants_model(intake_model)
    intake_role_label = resolve_source_intake_role(intake_role)
    intake_prompt_text = clean_source_intake_value(intake_prompt)
    intake_instruction = build_source_intake_instruction(
        intake_role_label,
        intake_prompt_text,
    )

    update_operation_progress(
        operation_id,
        phase="uploading_to_ai",
        message="Uploading source to AI workspace",
        detail="Creating an AI file-search workspace for mind map derivation.",
        progress=42,
    )
    assistant = openai.beta.assistants.create(
        name="MindMap Builder",
        instructions="Your task is to create the mindmap of the document",
        model=assistant_model,
        tools=[{"type": "file_search"}],
    )
    vector_store = openai.beta.vector_stores.create(name=f"{file_extension}_mindmap_{flow_id}")

    assistant = openai.beta.assistants.update(
        assistant_id=assistant.id,
        tool_resources={"file_search": {"vector_store_ids": [vector_store.id]}},
    )

    if file_extension == "pdf":
        mime_type = "application/pdf"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "csv":
        mime_type = "application/json"
        df = pd.read_csv(BytesIO(file_bytes))
        json_data = df.to_dict(orient="records")
        json_str = json.dumps(json_data)
        messages_file = openai.files.create(
            file=("data.json", json_str, mime_type), purpose="assistants"
        )
    elif file_extension == "txt":
        mime_type = "text/plain"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "docx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "html":
        mime_type = "text/html"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "md":
        mime_type = "text/markdown"
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    elif file_extension == "pptx":
        mime_type = (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        messages_file = openai.files.create(
            file=(source_document["filename"], file_bytes, mime_type), purpose="assistants"
        )
    else:
        raise ValueError("Unsupported file type")
    print("======================================", file_extension)
    update_operation_progress(
        operation_id,
        phase="ai_deriving",
        message="AI is deriving workspace nodes",
        detail="The model is generating a React Flow compatible mind map.",
        progress=64,
    )
    thread = openai.beta.threads.create(
        messages=[
            {
                "role": "user",
                "content" : f"""
                You are tasked with generating a JSON mind map that is compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question

                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": {json.dumps(intake_role_label)},
                            "model_name": {json.dumps(assistant_model if intake_role_label or intake_prompt_text else "")},
                            "intake_model": {json.dumps(assistant_model)},
                            "intake_prompt": {json.dumps(intake_prompt_text)},
                            "name": "{file_extension}", !!!DOESN"T CHANGES
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{source_document["filename"]}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "{file_extension}",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "{file_extension}"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.
                {intake_instruction}

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                {append_ai_graph_prompt_contract("")}
                """,

                "attachments": [
                    {"file_id": messages_file.id, "tools": [{"type": "file_search"}]}
                ],
            }
        ]
    )

    # Run the assistant and fetch the results
    run = openai.beta.threads.runs.create_and_poll(
        thread_id=thread.id, assistant_id=assistant.id
    )

    messages = list(
        openai.beta.threads.messages.list(thread_id=thread.id, run_id=run.id)
    )
    message_content = messages[0].content[0].text
    annotations = message_content.annotations

    # Annotate the summary content if necessary
    for index, annotation in enumerate(annotations):
        message_content.value = message_content.value.replace(
            annotation.text, f"[{index}]"
        )

    print(message_content.value)

    update_operation_progress(
        operation_id,
        phase="grounding",
        message="Grounding generated nodes",
        detail="Checking graph shape and attaching source references where possible.",
        progress=82,
    )
    response_json = parse_ai_mindmap_or_422(message_content.value)
    response_json = ground_mindmap_with_source_refs(response_json, source_context)

    component_metadata = {
        "flow_id": ObjectId(flow_id),
        "file_id": messages_file.id,
        "assistant_id": assistant.id,
        "vector_store_id": vector_store.id,
        "size": len(file_bytes),
        "type": file_extension,
        "processing_type": "legacy_assistants",
        "ai_provider": {
            "provider": "assistants_legacy_fallback",
            "model": assistant_model,
            "reason": "Prepared document sections were unavailable, so graph generation used the temporary Assistants file-search fallback.",
        },
        "mindmap_json": response_json,
        **source_metadata_fields(source_context),
    }

    component_id = component_collection.insert_one(component_metadata).inserted_id
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    update_operation_progress(
        operation_id,
        phase="complete",
        message="Workspace structure is ready",
        detail="The generated mind map was saved to the workspace.",
        progress=100,
        status_value="completed",
    )
    return {
        "flow_id": flow_id,
        "flow_name": flow["flow_name"],
        "component_id": str(component_id),
        "type": file_extension,
        "mindmap_json": response_json,
        "flow_type": flow_type
    }


def get_page_len(file: UploadFile):
    try:
        f_bytes = read_upload_bytes(file)
        reader = pdfium.PdfDocument(f_bytes)
        return len(reader) > 100
    except Exception as exc:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Malformed PDF file: {str(exc)}",
        ) from exc


@app.post("/component-create-pdf")
def create_pdf_component(
    file: UploadFile,
    flow_id: str = Form(...),
    processing_type: str = Form(...),
    source_intent: str | None = Form(None),
    operation_id: str | None = Form(None),
    intake_role: str | None = Form(None),
    intake_prompt: str | None = Form(None),
    intake_model: str | None = Form(None),
):
    update_operation_progress(
        operation_id,
        phase="validating",
        message="Validating PDF upload",
        detail=file.filename or "",
        progress=12,
    )
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    try:
        upload = validate_upload_bytes(file.filename, read_upload_bytes(file))
    except DocumentIngestionError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="PDF validation failed",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise ingestion_http_error(exc) from exc

    processing_flow_type = (
        "automatic"
        if source_intent == "mindmap"
        else "manual"
        if source_intent == "context"
        else flow["flow_type"]
    )

    if upload["extension"] == "pdf":
        print(get_page_len(file))
        check_page_length = get_page_len(file)
        if processing_type == "gpt" and not check_page_length and processing_flow_type == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type='manual', operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        elif processing_type == "aws" and processing_flow_type == 'manual':
            update_operation_progress(
                operation_id,
                phase="extracting",
                message="Extracting PDF with AWS Textract",
                detail="Large PDF extraction can take a few minutes.",
                progress=35,
            )
            result = use_aws_textract(file, flow_id=flow_id, flow_type='manual')
            update_operation_progress(
                operation_id,
                phase="complete",
                message="PDF source is ready",
                detail="The source component was saved to the workspace.",
                progress=100,
                status_value="completed",
            )
            return result
        elif processing_type == "custom" and processing_flow_type == 'manual':
            update_operation_progress(
                operation_id,
                phase="extracting",
                message="Extracting PDF tables and text",
                detail="TraceSpace is parsing pages and tables before saving the source.",
                progress=35,
            )
            result = camelot_pdf_processing(flow_id, file, 'manual')
            update_operation_progress(
                operation_id,
                phase="complete",
                message="PDF source is ready",
                detail="The source component was saved to the workspace.",
                progress=100,
                status_value="completed",
            )
            return result
        elif processing_type == "gpt" and not check_page_length and processing_flow_type == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type='automatic', operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        elif processing_type == "aws" and processing_flow_type == 'automatic':
            update_operation_progress(
                operation_id,
                phase="extracting",
                message="Extracting PDF with AWS Textract",
                detail="Large PDF extraction can take a few minutes.",
                progress=35,
            )
            result = use_aws_textract(file, flow_id=flow_id, flow_type='automatic')
            update_operation_progress(
                operation_id,
                phase="complete",
                message="Generated PDF workspace is ready",
                detail="The derived mind map was saved to the workspace.",
                progress=100,
                status_value="completed",
            )
            return result
        elif processing_type == "custom" and processing_flow_type == 'automatic':
            update_operation_progress(
                operation_id,
                phase="extracting",
                message="Extracting PDF tables and text",
                detail="TraceSpace is parsing pages and tables before deriving nodes.",
                progress=35,
            )
            result = camelot_pdf_processing(flow_id, file, 'automatic')
            update_operation_progress(
                operation_id,
                phase="complete",
                message="Generated PDF workspace is ready",
                detail="The derived mind map was saved to the workspace.",
                progress=100,
                status_value="completed",
            )
            return result
        else:
            traceback.print_exc()
            update_operation_progress(
                operation_id,
                phase="failed",
                message="PDF exceeds selected processing limits",
                detail="Try AWS or custom PDF processing for larger files.",
                progress=100,
                status_value="failed",
            )
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        update_operation_progress(
            operation_id,
            phase="failed",
            message="PDF validation failed",
            detail="Only PDF files are allowed.",
            progress=100,
            status_value="failed",
        )
        raise HTTPException(status_code=400, detail="Only PDF files are allowed.")

@app.post("/component-create-img")
async def create_img_component(
    flow_id: str = Form(...),
    file: UploadFile = File(...),
    operation_id: str | None = Form(None),
):
    try:
        update_operation_progress(
            operation_id,
            phase="validating",
            message="Validating image upload",
            detail=file.filename or "",
            progress=12,
        )
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        MAX_IMAGE_SIZE_MB = 16
        ALLOWED_MIME_TYPES = {
            "image/jpeg",
            "image/png",
            "image/webp",
            "image/heic",
            "image/heif",
        }

        if file.content_type not in ALLOWED_MIME_TYPES:
            raise HTTPException(
                status_code=400, detail=f"Unsupported file type: {file.content_type}"
            )

        contents = await file.read()

        size_in_mb = len(contents) / (1024 * 1024)
        if size_in_mb > MAX_IMAGE_SIZE_MB:
            raise HTTPException(status_code=400, detail="Image exceeds 16MB size limit")

        image_base64 = base64.b64encode(contents).decode("utf-8")
        source_context = binary_source_context(
            filename=file.filename,
            file_bytes=contents,
            source_type="image",
            flow_id=flow_id,
        )

        if flow["flow_type"] == 'manual':

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "mime_type": file.content_type,
                "type": "image",
                "base64_image": image_base64,
                "processing_type": "openai",
                "instructions": "",
                "persona_name": "TraceSpace reviewer",
                **source_metadata_fields(source_context),
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id
            update_operation_progress(
                operation_id,
                phase="complete",
                message="Image source is ready",
                detail="The image component was saved to the workspace.",
                progress=100,
                status_value="completed",
            )

            return {
                "message": "Image component created successfully",
                "component_id": str(component_id),
                "type": "image",
            }

        update_operation_progress(
            operation_id,
            phase="ai_deriving",
            message="AI is reading the image",
            detail="Deriving reviewable workspace nodes from the image.",
            progress=62,
        )
        response_json = generate_image_mindmap(
            file_name=file.filename,
            mime_type=file.content_type,
            contents=contents,
            flow_id=flow_id,
            model=None,
        )

        response_json = ground_mindmap_with_source_refs(response_json, source_context)
        print(response_json)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "type": "image",
            "processing_type": "openai",
            "instructions": "",
            "persona_name": "TraceSpace reviewer",
            "mindmap_json": response_json,
            **source_metadata_fields(source_context),
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id
        update_operation_progress(
            operation_id,
            phase="complete",
            message="Generated image workspace is ready",
            detail="The derived mind map was saved to the workspace.",
            progress=100,
            status_value="completed",
        )

        return {
            "flow_id" : ObjectId(flow_id),
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "image",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }

    except HTTPException as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Image upload failed",
            detail=str(exc.detail),
            progress=100,
            status_value="failed",
        )
        raise
    except Exception as e:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Image upload failed",
            detail=str(e),
            progress=100,
            status_value="failed",
        )
        print(f"Error in /component-create-img endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-audio")
async def create_audio_component(
    flow_id: str = Form(...),
    file: UploadFile = File(...),
    operation_id: str | None = Form(None),
):
    try:
        update_operation_progress(
            operation_id,
            phase="validating",
            message="Validating audio upload",
            detail=file.filename or "",
            progress=12,
        )
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        MAX_AUDIO_SIZE_MB = 16
        ALLOWED_MIME_TYPES = {
            "audio/wav",
            "audio/mp3",
            "audio/aiff",
            "audio/aac",
            "audio/ogg",
            "audio/flac",
            "audio/mpeg",
        }

        if file.content_type not in ALLOWED_MIME_TYPES:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported audio file type: {file.content_type}",
            )

        contents = await file.read()
        size_in_mb = len(contents) / (1024 * 1024)

        if size_in_mb > MAX_AUDIO_SIZE_MB:
            raise HTTPException(status_code=400, detail="Audio exceeds 16MB size limit")

        audio_base64 = base64.b64encode(contents).decode("utf-8")
        source_context = binary_source_context(
            filename=file.filename,
            file_bytes=contents,
            source_type="audio",
            flow_id=flow_id,
        )

        if flow["flow_type"] == 'manual':

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "mime_type": file.content_type,
                "type": "audio",
                "base64_audio": audio_base64,
                "processing_type": "openai",
                "instructions": "",
                "persona_name": "TraceSpace reviewer",
                **source_metadata_fields(source_context),
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id
            update_operation_progress(
                operation_id,
                phase="complete",
                message="Audio source is ready",
                detail="The audio component was saved to the workspace.",
                progress=100,
                status_value="completed",
            )

            return {
                "message": "Audio component created successfully",
                "component_id": str(component_id),
                "type": "audio",
            }

        update_operation_progress(
            operation_id,
            phase="transcribing",
            message="Transcribing audio",
            detail="AI is converting the audio to text.",
            progress=52,
        )
        transcript = transcribe_audio(
            file_name=file.filename,
            mime_type=file.content_type,
            contents=contents,
        )
        update_operation_progress(
            operation_id,
            phase="ai_deriving",
            message="AI is deriving workspace nodes",
            detail="Building a mind map from the transcript.",
            progress=72,
        )
        response_json = generate_audio_mindmap(
            file_name=file.filename,
            transcript=transcript,
            flow_id=flow_id,
            model=None,
        )
        response_json = ground_mindmap_with_source_refs(response_json, source_context)
        print(response_json)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "type": "audio",
            "processing_type": "openai",
            "instructions": "",
            "persona_name": "TraceSpace reviewer",
            "transcript": transcript,
            "mindmap_json": response_json,
            **source_metadata_fields(source_context),
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id
        update_operation_progress(
            operation_id,
            phase="complete",
            message="Generated audio workspace is ready",
            detail="The derived mind map was saved to the workspace.",
            progress=100,
            status_value="completed",
        )

        return {
            "flow_id" : flow_id,
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "audio",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }

    except HTTPException as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Audio upload failed",
            detail=str(exc.detail),
            progress=100,
            status_value="failed",
        )
        raise
    except Exception as e:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Audio upload failed",
            detail=str(e),
            progress=100,
            status_value="failed",
        )
        print(f"Error in /component-create-audio endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-youtube")
def create_youtube_component(
    flow_id: str = Form(...),
    youtube_url: str = Form(...),
    operation_id: str | None = Form(None),
):
    try:
        update_operation_progress(
            operation_id,
            phase="validating",
            message="Preparing YouTube source",
            detail=youtube_url,
            progress=12,
        )
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})

        print(youtube_url)
        source_context = virtual_source_context(
            label=youtube_url,
            source_type="youtube",
            flow_id=flow_id,
        )

        if flow["flow_type"] == 'manual':
            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "youtube_url": youtube_url,
                "type": "youtube",
                "processing_type": "gemini",
                "instructions": "",
                "persona_name": "TraceSpace reviewer",
                **source_metadata_fields(source_context),
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id
            update_operation_progress(
                operation_id,
                phase="complete",
                message="YouTube source is ready",
                detail="The YouTube source was saved to the workspace.",
                progress=100,
                status_value="completed",
            )

            return {
                "message": "Youtube component created successfully",
                "component_id": str(component_id),
                "type": "youtube",
            }

        else:
            mime_type = "video/*"

            template = f"""
                You are tasked with generating a JSON mind map for give youtube URL and should be compatible with React Flow for rendering a flow diagram which should cover all the details and important aspects of the component for which multiple nodes can be required. The mind map should adhere to the following rules:

                1. **Node Types:**
                - There will always be one `dataSource` node, which serves as the root of the flow.
                - There will be `question` node which will be connected to the subsequent `response` node.
                - The `question` node can be connected to data sources or other `response` nodes.
                - There will be `response` for the above question

                2. **Node Relationships:**
                - `response` nodes may also connect to each other if it improves the logical flow or visualization.
                - `question` node will always have a `response` node
                - `dataSource` node will always be connected to a question node

                3. **Node Properties:**
                - Each node should have:
                    - `id` (unique identifier of 12 or 24 digit unique uuid or nanoid)
                    - `type` (`dataSource` or `response`)
                    - `position` (coordinates in the form {{ "x": <number>, "y": <number> }} for layout)
                    - `measured` (an object defining width and height):
                        {{
                            "width": <number>,
                            "height": <number>
                        }}
                    - `targetPosition` (position of the target connection, default to `"left"`)
                    - `sourcePosition` (position of the source connection, default to `"right"`)
                    - `selected` (boolean, default to `false`)
                    - `deletable` (boolean, default to `true` for `response` and `false` for `dataSource`)

                4. **Node Data Format:**
                - `dataSource` Node:
                    - `data` contains the following properties:
                        {{
                            "prompt": "<data source description>",
                            "name": "youtube", !!!DOESN"T CHANGES
                            "content": "<file name or content>",
                            "flow_id": "{flow_id}",
                            "file": "{youtube_url}"  // Empty object or file metadata
                        }}
                5. **Question Data Format:**
                - `question` Node:
                    - `data` contains the following properties:
                        {{
                            "question": "<the question asked for the response>",
                            "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "component_type" : "youtube",
                        }}
                6. **RESPONSE NODE FORMAT**
                - `response` Node:
                    - `data` contains nested properties:
                        {{
                            "id": "<unique identifier of 12 or 24 digit unique uuid or nanoid>",
                            "type": "response" !!DOESN'T CHANGE,
                            "data": {{
                                "question": "<question text, if applicable>",
                                "summ": "<!!give me a detailed answer for the above question>",
                                "df": [],
                                "graph": "",
                                "flow_id": "{flow_id}",
                                "component_id": "<component reference ID - unique identifier of 12 or 24 digit unique uuid or nanoid>",
                                "component_type": "youtube"
                            }}
                        }}

                7. **Connections:**
                - Connections between nodes should be represented by edges, with the following format:
                    - `id` (unique identifier for the edge)
                    - `source` (ID of the source node)
                    - `target` (ID of the target node)
                    - `type` (optional, defaults to `default`)
                    - 'animated' !!WILL ALWAYS BE TRUE

                8. **Viewport Configuration:**
                - Include a `viewport` object that specifies:
                    - `x` (horizontal position of the viewport)
                    - `y` (vertical position of the viewport)
                    - `zoom` (zoom level for initial rendering)

                ### Additional Considerations:
                - Ensure that the node positions are distributed properly to avoid overlap.
                - Prioritize connecting `response` nodes where it adds logical structure to the flow.

                ### IMPORTANT:
                - **RETURN ONLY THE VALID JSON OBJECT AND NO ADDITIONAL COMMENTS**.
                - Do **not** include any explanations, text, or additional information.
                - Maintain the format with double curly braces `{{` and `}}` as shown in the format.
                """

        update_operation_progress(
            operation_id,
            phase="ai_deriving",
            message="AI is reading the YouTube source",
            detail="Deriving reviewable workspace nodes from the video URL.",
            progress=64,
        )
        response = model_vertexai.generate_content(
            contents=[template, Part.from_uri(youtube_url, mime_type)]
        )

        response_json = parse_ai_mindmap_or_422(response.text)
        response_json = ground_mindmap_with_source_refs(response_json, source_context)

        print(response_json)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "youtube_url": youtube_url,
            "type": "youtube",
            "processing_type": "gemini",
            "instructions": "",
            "persona_name": "TraceSpace reviewer",
            "mindmap_json": response_json,
            **source_metadata_fields(source_context),
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id
        update_operation_progress(
            operation_id,
            phase="complete",
            message="Generated YouTube workspace is ready",
            detail="The derived mind map was saved to the workspace.",
            progress=100,
            status_value="completed",
        )

        return {
            "flow_id" : flow_id,
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "youtube",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }


    except HTTPException as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="YouTube source failed",
            detail=str(exc.detail),
            progress=100,
            status_value="failed",
        )
        raise
    except Exception as e:
        traceback.print_exc()
        update_operation_progress(
            operation_id,
            phase="failed",
            message="YouTube source failed",
            detail=str(e),
            progress=100,
            status_value="failed",
        )
        print(f"Error in /component-create-youtube endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-video")
async def create_video_component(
    flow_id: str = Form(...),
    file: UploadFile = File(...),
    operation_id: str | None = Form(None),
):
    try:
        update_operation_progress(
            operation_id,
            phase="validating",
            message="Validating video upload",
            detail=file.filename or "",
            progress=12,
        )
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        MAX_VIDEO_SIZE_MB = 16
        ALLOWED_MIME_TYPES = {
            "video/x-flv",
            "video/quicktime",
            "video/mpeg",
            "video/mpegs",
            "video/mpgs",
            "video/mpg",
            "video/mp4",
            "video/webm",
            "video/wmv",
            "video/3gpp",
        }

        if file.content_type not in ALLOWED_MIME_TYPES:
            raise HTTPException(
                status_code=400,
                detail=f"Unsupported video file type: {file.content_type}",
            )

        contents = await file.read()
        size_in_mb = len(contents) / (1024 * 1024)

        if size_in_mb > MAX_VIDEO_SIZE_MB:
            raise HTTPException(status_code=400, detail="Video exceeds 16MB size limit")

        video_base64 = base64.b64encode(contents).decode("utf-8")
        source_context = binary_source_context(
            filename=file.filename,
            file_bytes=contents,
            source_type="video",
            flow_id=flow_id,
        )

        if flow["flow_type"] == 'manual':

            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "name": file.filename,
                "mime_type": file.content_type,
                "type": "video",
                "base64_video": video_base64,
                "processing_type": "openai_local_frames",
                "instructions": "",
                "persona_name": "TraceSpace reviewer",
                **source_metadata_fields(source_context),
            }

            component_id = component_collection.insert_one(component_metadata).inserted_id
            update_operation_progress(
                operation_id,
                phase="complete",
                message="Video source is ready",
                detail="The video component was saved to the workspace.",
                progress=100,
                status_value="completed",
            )

            return {
                "message": "Video component created successfully",
                "component_id": str(component_id),
                "type": "video",
            }

        update_operation_progress(
            operation_id,
            phase="ai_deriving",
            message="AI is sampling and reading video",
            detail="Deriving reviewable workspace nodes from video frames and audio.",
            progress=64,
        )
        response_json = generate_video_mindmap(
            file_name=file.filename,
            mime_type=file.content_type,
            contents=contents,
            flow_id=flow_id,
            model=None,
        )

        response_json = ground_mindmap_with_source_refs(response_json, source_context)
        print(response_json)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "type": "video",
            "processing_type": "openai_local_video",
            "instructions": "",
            "persona_name": "TraceSpace reviewer",
            "audio_status": response_json.get("metadata", {})
            .get("video_audio", {})
            .get("status", ""),
            "transcript": response_json.get("metadata", {})
            .get("video_audio", {})
            .get("transcript", ""),
            "mindmap_json": response_json,
            **source_metadata_fields(source_context),
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id
        update_operation_progress(
            operation_id,
            phase="complete",
            message="Generated video workspace is ready",
            detail="The derived mind map was saved to the workspace.",
            progress=100,
            status_value="completed",
        )

        return {
            "flow_id" : flow_id,
            "flow_name": flow["flow_name"],
            "component_id": str(component_id),
            "type": "video",
            "mindmap_json": response_json,
            "flow_type": "automatic"
        }

    except HTTPException as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Video upload failed",
            detail=str(exc.detail),
            progress=100,
            status_value="failed",
        )
        raise
    except Exception as e:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Video upload failed",
            detail=str(e),
            progress=100,
            status_value="failed",
        )
        print(f"Error in /component-create-video endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-txt")
def create_txt_component(
    file: UploadFile,
    flow_id: str = Form(...),
    operation_id: str | None = Form(None),
    intake_role: str | None = Form(None),
    intake_prompt: str | None = Form(None),
    intake_model: str | None = Form(None),
):
    update_operation_progress(
        operation_id,
        phase="validating",
        message="Validating text upload",
        detail=file.filename or "",
        progress=12,
    )
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    try:
        upload = validate_upload_bytes(file.filename, read_upload_bytes(file))
    except DocumentIngestionError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Text validation failed",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise ingestion_http_error(exc) from exc

    if upload["extension"] == "txt":
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        else:
            traceback.print_exc()
            update_operation_progress(
                operation_id,
                phase="failed",
                message="Text source exceeds AI token limit",
                detail="Split the file into smaller sources and try again.",
                progress=100,
                status_value="failed",
            )
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Text validation failed",
            detail="Only TXT files are allowed.",
            progress=100,
            status_value="failed",
        )
        raise HTTPException(status_code=400, detail="Only TXT files are allowed.")


@app.post("/component-create-md")
def create_md_component(
    file: UploadFile,
    flow_id: str = Form(...),
    operation_id: str | None = Form(None),
    intake_role: str | None = Form(None),
    intake_prompt: str | None = Form(None),
    intake_model: str | None = Form(None),
):
    update_operation_progress(
        operation_id,
        phase="validating",
        message="Validating Markdown upload",
        detail=file.filename or "",
        progress=12,
    )
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    try:
        upload = validate_upload_bytes(file.filename, read_upload_bytes(file))
    except DocumentIngestionError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Markdown validation failed",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise ingestion_http_error(exc) from exc

    if upload["extension"] == "md":
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        else:
            traceback.print_exc()
            update_operation_progress(
                operation_id,
                phase="failed",
                message="Markdown source exceeds AI token limit",
                detail="Split the file into smaller sources and try again.",
                progress=100,
                status_value="failed",
            )
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Markdown validation failed",
            detail="Only Markdown files are allowed.",
            progress=100,
            status_value="failed",
        )
        raise HTTPException(status_code=400, detail="Only MarkDown files are allowed.")


@app.post("/component-create-pptx")
def create_pptx_component(
    file: UploadFile,
    flow_id: str = Form(...),
    operation_id: str | None = Form(None),
    intake_role: str | None = Form(None),
    intake_prompt: str | None = Form(None),
    intake_model: str | None = Form(None),
):
    update_operation_progress(
        operation_id,
        phase="validating",
        message="Validating PPTX upload",
        detail=file.filename or "",
        progress=12,
    )
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    try:
        upload = validate_ai_intake_bytes(file.filename, read_upload_bytes(file))
    except DocumentIngestionError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="PPTX validation failed",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise ingestion_http_error(exc) from exc

    if upload["extension"] == "pptx":
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        else:
            traceback.print_exc()
            update_operation_progress(
                operation_id,
                phase="failed",
                message="PPTX source exceeds AI token limit",
                detail="Split the deck into smaller sources and try again.",
                progress=100,
                status_value="failed",
            )
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        update_operation_progress(
            operation_id,
            phase="failed",
            message="PPTX validation failed",
            detail="Only PPTX files are allowed.",
            progress=100,
            status_value="failed",
        )
        raise HTTPException(status_code=400, detail="Only PPTX files are allowed.")


@app.post("/component-create-html")
def create_html_component(
    file: UploadFile,
    flow_id: str = Form(...),
    operation_id: str | None = Form(None),
    intake_role: str | None = Form(None),
    intake_prompt: str | None = Form(None),
    intake_model: str | None = Form(None),
):
    update_operation_progress(
        operation_id,
        phase="validating",
        message="Validating HTML upload",
        detail=file.filename or "",
        progress=12,
    )
    flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
    try:
        upload = validate_ai_intake_bytes(file.filename, read_upload_bytes(file))
    except DocumentIngestionError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="HTML validation failed",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise ingestion_http_error(exc) from exc

    if upload["extension"] == "html":
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and flow["flow_type"] == 'manual':
            return get_summary_from_openai(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        elif check_page_length and flow["flow_type"] == 'automatic':
            return openai_mindmap_generator(file, flow_id=flow_id, flow_type=flow["flow_type"], operation_id=operation_id, intake_role=intake_role, intake_prompt=intake_prompt, intake_model=intake_model)
        else:
            traceback.print_exc()
            update_operation_progress(
                operation_id,
                phase="failed",
                message="HTML source exceeds AI token limit",
                detail="Split the file into smaller sources and try again.",
                progress=100,
                status_value="failed",
            )
            return HTTPException(status_code=404, detail="Exceeded Page limit for GPT.")
    else:
        traceback.print_exc()
        update_operation_progress(
            operation_id,
            phase="failed",
            message="HTML validation failed",
            detail="Only HTML files are allowed.",
            progress=100,
            status_value="failed",
        )
        raise HTTPException(status_code=400, detail="Only HTML files are allowed.")


@app.post("/component-create-docx")
def create_docx_component(
    file: UploadFile,
    flow_id: str = Form(...),
    operation_id: str | None = Form(None),
    source_intent: str | None = Form(None),
    intake_role: str | None = Form(None),
    intake_model: str | None = Form(None),
    intake_prompt: str | None = Form(None),
):
    update_operation_progress(
        operation_id,
        phase="validating",
        message="Validating DOCX upload",
        detail=file.filename or "",
        progress=12,
    )
    try:
        flow = get_upload_flow_or_400(flow_id)
        upload = validate_upload_bytes(file.filename, read_upload_bytes(file))
        if upload["extension"] != "docx":
            raise DocumentIngestionError("Only DOCX files are allowed.")
    except DocumentIngestionError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="DOCX validation failed",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise ingestion_http_error(exc) from exc
    except HTTPException as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="DOCX upload could not start",
            detail=str(exc.detail),
            progress=100,
            status_value="failed",
        )
        raise

    try:
        processing_flow_type = (
            "automatic"
            if source_intent == "mindmap"
            else "manual"
            if source_intent == "context"
            else flow["flow_type"]
        )
        check_page_length = is_within_gpt4o_token_limit(file)
        if check_page_length and processing_flow_type == 'manual':
            return get_summary_from_openai(
                file,
                flow_id=flow_id,
                flow_type='manual',
                operation_id=operation_id,
                intake_role=intake_role,
                intake_prompt=intake_prompt,
                intake_model=intake_model,
            )
        elif check_page_length and processing_flow_type == 'automatic':
            return openai_mindmap_generator(
                file,
                flow_id=flow_id,
                flow_type='automatic',
                operation_id=operation_id,
                intake_role=intake_role,
                intake_prompt=intake_prompt,
                intake_model=intake_model,
            )
        else:
            traceback.print_exc()
            update_operation_progress(
                operation_id,
                phase="failed",
                message="DOCX source exceeds AI token limit",
                detail="Split the file into smaller sources and try again.",
                progress=100,
                status_value="failed",
            )
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail="DOCX is too large for OpenAI processing. Split it into smaller source files and try again.",
            )
    except HTTPException as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="DOCX processing failed",
            detail=str(exc.detail),
            progress=100,
            status_value="failed",
        )
        raise
    except MissingConfigurationError as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Missing AI settings",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise configuration_http_error(exc) from exc
    except Exception as exc:
        traceback.print_exc()
        update_operation_progress(
            operation_id,
            phase="failed",
            message="DOCX processing failed",
            detail=str(exc),
            progress=100,
            status_value="failed",
        )
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"DOCX processing failed: {str(exc)}",
        ) from exc


@app.post("/component-create-csv")
def create_csv_component(
    file: UploadFile = File(...),
    flow_id: str = Form(...),
    header_row: int = Form(...),
    operation_id: str | None = Form(None),
):
    try:
        update_operation_progress(
            operation_id,
            phase="validating",
            message="Preparing CSV source",
            detail=file.filename,
            progress=10,
        )
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        if flow["flow_type"] != 'manual':
            raise HTTPException(status_code=400, detail="Only Manual Mindmap is supported for CSV.")
        if not (file.filename or "").lower().endswith(".csv"):
            raise HTTPException(status_code=400, detail="Only CSV files are allowed.")

        file_bytes = file.file.read()
        update_operation_progress(
            operation_id,
            phase="extracting",
            message="Reading CSV table",
            detail="Checking duplicates and source metadata.",
            progress=25,
        )
        file_hash = calculate_file_hash(file_bytes)
        source_context = binary_source_context(
            filename=file.filename,
            file_bytes=file_bytes,
            source_type="csv",
            flow_id=flow_id,
        )

        existing_component = component_collection.find_one(
            {"file_hash": file_hash, "flow_id": ObjectId(flow_id)}
        )
        if existing_component:
            raise HTTPException(
                status_code=400, detail="File already exists in the system."
            )

        unique_table_name = f"tbl_{uuid4().hex[:8]}"

        file_name = file.filename
        folder = f"uploads/{flow_id}/"
        s3_key = folder + file_name
        upload_to_s3(file_bytes, bucket_name, s3_key)
        print("uploaded")
        update_operation_progress(
            operation_id,
            phase="importing",
            message="Importing CSV into the query engine",
            detail="Creating a temporary table and reading headers.",
            progress=48,
        )

        sql_con = sqlite3.connect("csv_data.db")
        buffer = BytesIO(file_bytes)
        df = pd.read_csv(
            buffer, skiprows=header_row, encoding="utf-8", encoding_errors="ignore"
        )
        print(df)
        buffer.close()
        file.file.close()
        print("CSV into SQLite")

        df.to_sql(name=unique_table_name, con=sql_con, if_exists="replace", index=False)
        sql_con.close()
        csvBot.connect_to_sqlite("csv_data.db")
        update_operation_progress(
            operation_id,
            phase="ai_reading",
            message="AI is reading the CSV schema",
            detail=f"{len(df.columns)} columns were detected.",
            progress=68,
        )

        df_ddl = csvBot.run_sql(
            f"SELECT sql FROM sqlite_master WHERE type = 'table' AND name = '{unique_table_name}'"
        )

        for ddl in df_ddl["sql"].to_list():
            csvBot.train(ddl=ddl)
        update_operation_progress(
            operation_id,
            phase="ai_deriving",
            message="AI is preparing CSV query context",
            detail="Training data source metadata for follow-up questions.",
            progress=84,
        )

        training_data = csvBot.get_training_data()
        print(training_data)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "name": file.filename,
            "table_name": unique_table_name,
            "file_hash": file_hash,
            "size": len(file_bytes),
            "type": "csv",
            "s3_path": s3_key,
            "created_at": datetime.datetime.utcnow(),
            **source_metadata_fields(source_context),
        }
        component_id = component_collection.insert_one(component_metadata).inserted_id
        update_operation_progress(
            operation_id,
            phase="complete",
            message="CSV source ready",
            detail=file.filename,
            progress=100,
            status_value="completed",
        )
        return {
            "component_id": str(component_id),
            "type": "csv",
            "message": "Component created successfully",
        }

    except HTTPException as e:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="CSV source could not be added",
            detail=str(e.detail),
            progress=100,
            status_value="failed",
        )
        raise e
    except Exception as e:
        print(f"Error in /component-create-csv endpoint: {e}")
        update_operation_progress(
            operation_id,
            phase="failed",
            message="CSV source could not be added",
            detail=str(e),
            progress=100,
            status_value="failed",
        )
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/component-create-crawl")
async def create_web_crawler(
    flow_id: str = Form(...),
    web_url: str = Form(...),
    operation_id: str | None = Form(None),
):
    try:
        update_operation_progress(
            operation_id,
            phase="validating",
            message="Preparing web source",
            detail=web_url,
            progress=12,
        )
        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        flow_type = flow["flow_type"]
        source_context = virtual_source_context(
            label=web_url,
            source_type="web",
            flow_id=flow_id,
        )
        if flow_type == "manual":
            component_metadata = {
                "flow_id": ObjectId(flow_id),
                "web_url": web_url,
                "type": "web",
                "processing_type": "openai_web_search",
                **source_metadata_fields(source_context),
            }
            component_id = component_collection.insert_one(component_metadata).inserted_id
            update_operation_progress(
                operation_id,
                phase="complete",
                message="Web source is ready",
                detail="The web source was saved to the workspace.",
                progress=100,
                status_value="completed",
            )
            return {
                "component_id": str(component_id),
                "type": "web",
                "message": "Web component created successfully",
            }

        update_operation_progress(
            operation_id,
            phase="ai_deriving",
            message="AI is reading the web page",
            detail="Deriving reviewable workspace nodes from the URL.",
            progress=64,
        )
        response_json = generate_web_mindmap(
            url=web_url,
            flow_id=flow_id,
            model=None,
        )
        response_json = ground_mindmap_with_source_refs(response_json, source_context)

        component_metadata = {
            "flow_id": ObjectId(flow_id),
            "web_url": web_url,
            "type": "web",
            "processing_type": "openai_web_search",
            "mindmap_json": response_json,
            **source_metadata_fields(source_context),
        }

        component_id = component_collection.insert_one(component_metadata).inserted_id
        update_operation_progress(
            operation_id,
            phase="complete",
            message="Generated web workspace is ready",
            detail="The derived mind map was saved to the workspace.",
            progress=100,
            status_value="completed",
        )

        return {
            "component_id": str(component_id),
            "type": "web",
            "mindmap_json": response_json,
            "flow_type": flow_type,
        }

    except HTTPException as exc:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Web source failed",
            detail=str(exc.detail),
            progress=100,
            status_value="failed",
        )
        raise
    except Exception as e:
        update_operation_progress(
            operation_id,
            phase="failed",
            message="Web source failed",
            detail=str(e),
            progress=100,
            status_value="failed",
        )
        print(f"Error in /component-create-crawl endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/pdf-component-qa", response_model=List[PDFNodeQueryResponse])
def PDF_QA(request: PDFNodeQueryRequest):
    return answer_component_with_responses(request, "pdf", PDFNodeQueryResponse)

@app.post("/txt-component-qa", response_model=List[TXTNodeQueryResponse])
def TXT_QA(request: TXTNodeQueryRequest):
    return answer_component_with_responses(request, "txt", TXTNodeQueryResponse)

@app.post("/md-component-qa", response_model=List[MDNodeQueryResponse])
def MD_QA(request: MDNodeQueryRequest):
    return answer_component_with_responses(request, "md", MDNodeQueryResponse)

@app.post("/html-component-qa", response_model=List[HTMLNodeQueryResponse])
def HTML_QA(request: HTMLNodeQueryRequest):
    return answer_component_with_responses(request, "html", HTMLNodeQueryResponse)

@app.post("/docx-component-qa", response_model=List[DOCXNodeQueryResponse])
def DOCX_QA(request: DOCXNodeQueryRequest):
    return answer_component_with_responses(request, "docx", DOCXNodeQueryResponse)

@app.post("/pptx-component-qa", response_model=List[PPTXNodeQueryResponse])
def PPTX_QA(request: PPTXNodeQueryRequest):
    return answer_component_with_responses(request, "pptx", PPTXNodeQueryResponse)

@app.post("/csv-component-qa", response_model=List[CSVNodeQueryResponse])
def CSV_QA(request: CSVNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "csv",
            }
        )

        if not record or "table_name" not in record:
            # Raise an HTTP 404 error if the table is not found
            raise HTTPException(
                status_code=404,
                detail="Table not found for the given flow_id and component_id",
            )

        table_name = record["table_name"]
        question_with_table = f"question - {request.query} for table name: {table_name}"
        sqlQuery = csvBot.generate_sql(question_with_table)
        if csvBot.is_sql_valid(sqlQuery):
            runSQLDF = csvBot.run_sql(sqlQuery)
            summSQL = csvBot.generate_summary(question_with_table, runSQLDF)
            code = csvBot.generate_plotly_code(
                question=question_with_table,
                sql=sqlQuery,
                df_metadata=f"Running df.dtypes gives:\n {runSQLDF.dtypes}",
            )
            fig = csvBot.get_plotly_figure(plotly_code=code, df=runSQLDF)
            plotyGraph = fig.to_json()
            df_dict = runSQLDF.to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]
            artifact_payload = build_structured_data_artifacts(
                source_type="csv",
                source_id=str(request.component_id),
                question=request.query,
                table_name=table_name,
                sql=sqlQuery,
                rows=df_dict,
                summary=summSQL,
                chart_json=plotyGraph,
            )
            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.query,
                "query": sqlQuery,
                "df": df_dict,
                "summ": summSQL,
                "graph": plotyGraph,
                "node_type": artifact_payload["node_type"],
                "artifact_type": artifact_payload["artifact_type"],
                "artifact_ids": artifact_payload["artifact_ids"],
                "review_state": artifact_payload["review_state"],
                "source_refs": artifact_payload["source_refs"],
                "generated_artifacts": artifact_payload["generated_artifacts"],
                "metadata": artifact_payload["metadata"],
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "csv",
                "is_delete": "false",
                "created_at": datetime.datetime.utcnow(),
            }
            node_id_response = node_collection.insert_one(result_document)
            question_entries = []
            response_data = {
                "id": str(request.node_id),
                "type": "CSVNode",
                "data": {
                    "question": request.query,
                    "query": sqlQuery,
                    "df": df_dict,
                    "summ": summSQL,
                    "graph": plotyGraph,
                    "node_type": artifact_payload["node_type"],
                    "artifact_type": artifact_payload["artifact_type"],
                    "artifact_ids": artifact_payload["artifact_ids"],
                    "review_state": artifact_payload["review_state"],
                    "source_refs": artifact_payload["source_refs"],
                    "generated_artifacts": artifact_payload["generated_artifacts"],
                    "metadata": artifact_payload["metadata"],
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "csv",
                },
            }
            question_entries.append(response_data)
            if request.request_type == "question":
                empty_node = {
                    "id": str(ObjectId()),
                    "type": "question",
                    "data": {
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "csv",
                    },
                }
                question_entries.append(empty_node)

            return question_entries

        else:
            print("No answer found from llm")
            df_dict = pd.DataFrame().to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]

            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.query,
                "df": df_dict,
                "summ": "",
                "graph": "",
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "csv",
                "is_delete": "false",
            }

            node_id_response = node_collection.insert_one(result_document)

        question_entries = []

        response_data = {
            "id": str(request.node_id),
            "type": "CSVNode",
            "data": {
                "question": request.query,
                "df": pd.DataFrame().to_dict(orient="records"),
                "summ": "",
                "graph": "",
                "flow_id": request.flow_id,
                "component_id": request.component_id,
                "component_type": "csv",
            },
        }

        question_entries.append(response_data)

        if request.request_type == "question":
            empty_node = {
                "id": str(ObjectId()),
                "type": "question",
                "data": {
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "csv",
                },
            }
            question_entries.append(empty_node)

            return question_entries
    except Exception as e:
        print(f"Error in /csv-component-qa: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/web-component-qa", response_model=List[WebNodeQueryResponse])
def WEB_QA(request: WebNodeQueryRequest):
    return answer_component_with_responses(request, "web", WebNodeQueryResponse)

@app.post("/img-component-qa", response_model=List[ImgNodeQueryResponse])
def IMG_QA(request: ImgNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "image",
            }
        )

        instructions = record.get("instructions", "")
        persona_name = record.get("persona_name", "TraceSpace reviewer")
        base64_image = record["base64_image"]
        mime_type = record["mime_type"]
        image_bytes = base64.b64decode(base64_image)

        image_part = {"mime_type": mime_type, "data": image_bytes}

        template = f"""
    You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:

    #### **Response Format:**
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}
    **Here is the persona:** {persona_name}

    ### **Example Output:**
    If the conversation history contains a table and a relevant graph, return:

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        response = model.generate_content(contents=[template, image_part])

        responseList = response.text
        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        print(responseList)

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "image",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            ImgNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "image",
                },
                id=request.node_id,
                type="ImageNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = ImgNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "image",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /img-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/audio-component-qa", response_model=List[AudioNodeQueryResponse])
def AUDIO_QA(request: AudioNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "audio",
            }
        )

        instructions = record.get("instructions", "")
        persona_name = record.get("persona_name", "TraceSpace reviewer")

        base64_audio = record["base64_audio"]
        mime_type = record["mime_type"]
        audio_bytes = base64.b64decode(base64_audio)

        audio_part = {"mime_type": mime_type, "data": audio_bytes}

        template = f"""You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:

    #### **Response Format:**
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}
    **Here is the persona:** {persona_name}

    ### **Example Output:**
    If the conversation history contains a table and a relevant graph, return:

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        response = model.generate_content(contents=[template, audio_part])

        responseList = response.text
        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "audio",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            AudioNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "audio",
                },
                id=request.node_id,
                type="AudioNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = AudioNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "audio",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /audio-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/youtube-component-qa", response_model=List[YoutubeNodeQueryResponse])
def YOUTUBE_QA(request: YoutubeNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "youtube",
            }
        )

        instructions = record.get("instructions", "")
        persona_name = record.get("persona_name", "TraceSpace reviewer")

        youtube_url = record["youtube_url"]
        mime_type = "video/*"

        template = f"""
      You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:

    #### **Response Format:**
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}
    **Here is the persona:** {persona_name}

    ### **Example Output:**
    If the conversation history contains a table and a relevant graph, return:

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        response = model_vertexai.generate_content(
            contents=[template, Part.from_uri(youtube_url, mime_type)]
        )

        responseList = response.text
        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "youtube",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            YoutubeNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "youtube",
                },
                id=request.node_id,
                type="YoutubeNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = YoutubeNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "youtube",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /youtube-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/video-component-qa", response_model=List[VideoNodeQueryResponse])
def VIDEO_QA(request: VideoNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "video",
            }
        )

        instructions = record.get("instructions", "")
        persona_name = record.get("persona_name", "TraceSpace reviewer")

        video_url = record.get("video_url")
        mime_type = record.get("mime_type", "video/*")
        if not video_url:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=(
                    "Video Q&A for local uploaded videos needs a transcript-backed "
                    "or externally addressable video source. Regenerate the map from "
                    "the video source or use the generated mind map review views."
                ),
            )

        template = f"""
          You are an AI assistant tasked with answering the user’s question based on the provided question and persona. Return the results in **JSON format** with the structure below:

    #### **Response Format:**
    {{
    "summ": "Your summarized response here...",
    "df": an array of JSON objects,
    "graph": "json_string_representation_of_plotly_graph"
    }}

    ### **Instructions:**
    1. Answer the question using the conversation history.
    2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
    3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
    4. If no graph is possible, return an empty string `""`.
    5. ** The graph's background will be black, so adjust the theme accordingly**.

    NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
    NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"

    **Here is the question:** {request.query}
    **Here is the persona:** {persona_name}

    ### **Example Output:**
    If the conversation history contains a table and a relevant graph, return:

    ```json
    {{
    "summ": "Based on the conversation, the key points discussed were...",
    "df": [
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }},
        {{
        "column1": "value1",
        "column2": "value2",
        "column3": "value3"
        }}
    ],
    "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
    }}
    """

        response = model_vertexai.generate_content(
            contents=[template, Part.from_uri(video_url, mime_type)]
        )

        print(response)
        print(response.text)

        responseList = response.text
        responseList = (
            responseList.replace("```json", "")
            .replace("```", "")
            .replace("\n", "")
            .strip()
        )

        response_json = json.loads(responseList)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "component_id": ObjectId(request.component_id),
            "question": request.query,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", {}),
            "type": "video",
            "is_delete": "false",
            "timestamp": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        question_entries = []

        question_entries.append(
            VideoNodeQueryResponse(
                data={
                    "question": request.query,
                    "summ": response_json.get("summ", ""),
                    "df": validate_dataframe(response_json.get("df", [])),
                    "graph": response_json.get("graph", {}),
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "video",
                },
                id=request.node_id,
                type="VideoNode",
            )
        )

        if request.request_type == "question":
            empty_question_entry = VideoNodeQueryResponse(
                id=str(ObjectId()),
                data={
                    "question": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "video",
                },
                type="question",
            )

            question_entries.append(empty_question_entry)

        return question_entries

    except Exception as e:
        print(f"Error in /video-component-qa endpoint: {e}")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.delete("/soft-delete-node/{node_id}", response_model=dict)
def soft_delete_node(node_id: str):
    try:
        result = node_collection.update_one(
            {"_id": ObjectId(node_id)}, {"$set": {"is_delete": "true"}}
        )

        if result.matched_count == 0:
            raise HTTPException(status_code=404, detail="Node not found.")

        return {"message": "Node soft deleted successfully."}

    except Exception as e:
        print(f"Error in /soft-delete-node endpoint: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.get("/get-all-flow-details/{flow_id}")
def get_flow_details(flow_id: str):
    try:
        if not ObjectId.is_valid(flow_id):
            raise HTTPException(status_code=400, detail="Invalid flow_id format.")

        flow = flow_collection.find_one({"_id": ObjectId(flow_id)})
        if not flow:
            raise HTTPException(status_code=404, detail="Flow not found.")

        components = list(component_collection.find({"flow_id": ObjectId(flow_id)}))
        if not components:
            raise HTTPException(
                status_code=404, detail="No components found for the given flow_id."
            )

        flow_details = {
            "flow": {
                "flow_id": str(flow["_id"]),
                "flow_name": flow.get("flow_name"),
                "description": flow.get("description"),
                "summary": flow.get("summary"),
            },
            "components": [],
        }

        for component in components:
            component_id = str(component["_id"])
            nodes = list(
                node_collection.find(
                    {
                        "component_id": ObjectId(component_id),
                        "flow_id": ObjectId(flow_id),
                    }
                )
            )

            component_details = {
                "component_id": component_id,
                "name": component.get("name"),
                "file_hash": component.get("file_hash"),
                "size": component.get("size"),
                "s3_path": component.get("s3_path"),
                "nodes": [],
            }

            for node in nodes:
                component_details["nodes"].append(
                    {
                        "node_id": str(node["_id"]),
                        "question": node.get("question"),
                        "answer": node.get("answer"),
                        "timestamp": node.get("timestamp"),
                    }
                )

            flow_details["components"].append(component_details)

        return flow_details

    except Exception as e:
        print(f"Error in /get-all-flow-details/{flow_id}: {e}")
        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/create_sql_component/", response_model=SQLComponentResponse)
def create_sql_component(request: SQLComponentRequest):
    try:
        update_operation_progress(
            request.operation_id,
            phase="validating",
            message="Preparing SQL source",
            detail=request.table_name,
            progress=12,
        )
        flow = flow_collection.find_one({"_id": ObjectId(request.flow_id)})

        if flow["flow_type"] != 'manual':
            raise HTTPException(status_code=400, detail="Only Manual Mindmap is supported for SQL.")

        update_operation_progress(
            request.operation_id,
            phase="schema",
            message="Reading SQL schema",
            detail=f"Looking for tables matching {request.table_name}.",
            progress=38,
        )
        df_ddl = sqlBot.run_sql("SELECT type, sql FROM sqlite_master WHERE sql IS NOT NULL AND name LIKE '%" + request.table_name + "%'")
        print(df_ddl)

        update_operation_progress(
            request.operation_id,
            phase="ai_reading",
            message="AI is reading SQL table context",
            detail="Training available table definitions for questions.",
            progress=68,
        )
        for ddl in df_ddl['sql'].to_list():
            sqlBot.train(ddl=ddl)

        training_data = sqlBot.get_training_data()
        print(training_data)
        update_operation_progress(
            request.operation_id,
            phase="saving",
            message="Adding SQL source node",
            detail=request.table_name,
            progress=88,
        )

        component_data = {
            "flow_id": ObjectId(request.flow_id),
            "type": "sql",
            "table_name": request.table_name,
            **source_metadata_fields(
                virtual_source_context(
                    label=request.table_name,
                    source_type="sql",
                    flow_id=request.flow_id,
                )
            ),
        }

        component_id = component_collection.insert_one(component_data).inserted_id
        update_operation_progress(
            request.operation_id,
            phase="complete",
            message="SQL source ready",
            detail=request.table_name,
            progress=100,
            status_value="completed",
        )

        return SQLComponentResponse(
            component_id=str(component_id),
            type="sql",
            message="Component created successfully",
        )

    except HTTPException as e:
        update_operation_progress(
            request.operation_id,
            phase="failed",
            message="SQL source could not be added",
            detail=str(e.detail),
            progress=100,
            status_value="failed",
        )
        raise e
    except Exception as e:
        print(f"Error in /create_sql_component/: {e}")
        update_operation_progress(
            request.operation_id,
            phase="failed",
            message="SQL source could not be added",
            detail=str(e),
            progress=100,
            status_value="failed",
        )
        raise HTTPException(
            status_code=500, detail=f"Error creating SQL component: {str(e)}"
        )


@app.post("/components-follow-up-questions", response_model=List[ComponentFollowUpQueryResponse])
def create_follow_up_questions(request: ComponentFollowUpQueryRequest):
    if request.component_type in RESPONSES_COMPONENT_TYPES:
        return follow_up_questions_with_responses(request)

    return follow_up_entries(
        flow_id=request.flow_id,
        component_id=request.component_id,
        component_type=request.component_type,
        questions=[],
    )

@app.post("/sql-component-qa", response_model=List[SQLNodeQueryResponse])
def SQL_QA(request: SQLNodeQueryRequest):
    try:
        record = component_collection.find_one(
            {
                "flow_id": ObjectId(request.flow_id),
                "_id": ObjectId(request.component_id),
                "type": "sql",
            }
        )

        if not record or "table_name" not in record:
            raise HTTPException(
                status_code=404,
                detail="Table not found for the given flow_id and component_id",
            )

        table_name = record["table_name"]

        question_with_table = (
            f"question - {request.question} for table name: {table_name}"
        )

        sqlQuery = sqlBot.generate_sql(question_with_table)
        if sqlBot.is_sql_valid(sqlQuery):
            runSQLDF = sqlBot.run_sql(sqlQuery)
            summSQL = sqlBot.generate_summary(question_with_table, runSQLDF)
            code = sqlBot.generate_plotly_code(
                question=question_with_table,
                sql=sqlQuery,
                df_metadata=f"Running df.dtypes gives:\n {runSQLDF.dtypes}",
            )
            fig = sqlBot.get_plotly_figure(plotly_code=code, df=runSQLDF)
            plotyGraph = fig.to_json()

            df_dict = runSQLDF.to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]
            artifact_payload = build_structured_data_artifacts(
                source_type="sql",
                source_id=str(request.component_id),
                question=request.question,
                table_name=table_name,
                sql=sqlQuery,
                rows=df_dict,
                summary=summSQL,
                chart_json=plotyGraph,
            )
            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.question,
                "query": sqlQuery,
                "df": df_dict,
                "summ": summSQL,
                "graph": plotyGraph,
                "node_type": artifact_payload["node_type"],
                "artifact_type": artifact_payload["artifact_type"],
                "artifact_ids": artifact_payload["artifact_ids"],
                "review_state": artifact_payload["review_state"],
                "source_refs": artifact_payload["source_refs"],
                "generated_artifacts": artifact_payload["generated_artifacts"],
                "metadata": artifact_payload["metadata"],
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "sql",
                "is_delete": "false",
                "created_at": datetime.datetime.utcnow(),
            }

            node_id_response = node_collection.insert_one(result_document)

            question_entries = []

            response_data = {
                "id": str(request.node_id),
                "type": "SQLNode",
                "data": {
                    "question": request.question,
                    "query": sqlQuery,
                    "df": df_dict,
                    "summ": summSQL,
                    "graph": plotyGraph,
                    "node_type": artifact_payload["node_type"],
                    "artifact_type": artifact_payload["artifact_type"],
                    "artifact_ids": artifact_payload["artifact_ids"],
                    "review_state": artifact_payload["review_state"],
                    "source_refs": artifact_payload["source_refs"],
                    "generated_artifacts": artifact_payload["generated_artifacts"],
                    "metadata": artifact_payload["metadata"],
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "sql",
                },
            }

            question_entries.append(response_data)

            if request.request_type == "question":
                empty_node = {
                    "id": str(ObjectId()),
                    "type": "question",
                    "data": {
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "sql",
                    },
                }
                question_entries.append(empty_node)

            return question_entries

        else:

            print("No answer found from llm")

            df_dict = pd.DataFrame().to_dict(orient="records")
            df_dict = [{str(k): v for k, v in row.items()} for row in df_dict]

            result_document = {
                "_id": ObjectId(request.node_id),
                "question": request.question,
                "query": "I don't know",
                "df": df_dict,
                "summ": "",
                "graph": "",
                "flow_id": ObjectId(request.flow_id),
                "component_id": ObjectId(request.component_id),
                "type": "sql",
                "is_delete": "false",
                "created_at": datetime.datetime.utcnow(),
            }

            node_id_response = node_collection.insert_one(result_document)

            question_entries = []

            response_data = {
                "id": str(request.node_id),
                "type": "SQLNode",
                "data": {
                    "question": request.question,
                    "query": "I don't know",
                    "df": pd.DataFrame().to_dict(orient="records"),
                    "summ": "",
                    "graph": "",
                    "flow_id": request.flow_id,
                    "component_id": request.component_id,
                    "component_type": "sql",
                },
            }

            question_entries.append(response_data)

            if request.request_type == "question":
                empty_node = {
                    "id": str(ObjectId()),
                    "type": "question",
                    "data": {
                        "question": "",
                        "flow_id": request.flow_id,
                        "component_id": request.component_id,
                        "component_type": "sql",
                    },
                }

                question_entries.append(empty_node)

            return question_entries

    except Exception as e:
        traceback.print_exc()
        print(f"Error in /sql-component-qa: {str(e.__traceback__)}")

        raise HTTPException(status_code=500, detail="Internal Server Error")


@app.post("/multiple-qa-summarize", response_model=MultipleQuestionAnswerQueryResponse)
def multiple_qa_summarize(request: MultipleQuestionAnswerQueryRequest):
    try:
        conversation = []

        for parent_id in request.parent_node_ids:
            question, answer = fetch_question_answer_from_node_collection(
                parent_id, request.flow_id
            )
            if question and answer:
                conversation.append({"role": "user", "content": question})
                conversation.append(
                    {"role": "assistant", "content": f"Answer: \n{answer}"}
                )

        print(conversation)

        template = """
        You are an AI assistant tasked with answering the user’s question based on the provided conversation history. Return the results in **JSON format** with the structure below:

        #### **Response Format:**
        {{
        "summ": "Your summarized response here...",
        "df": an array of JSON objects,
        "graph": "json_string_representation_of_plotly_graph"
        }}

        ### **Instructions:**
        1. Answer the question using the conversation history.
        2. Extract relevant tabular data into a JSON object compatible with Ag-Grid. If no table exists, return empty JSON object.
        3. If a dataframe is available, generate a relevant **Plotly graph**. Return it as a **valid JSON string** that can be parsed in React.js.
        4. If no graph is possible, return an empty string `""`.
        5. ** The graph's background will be black, so adjust the theme accordingly**.

        **Here is the question:** {query}
        **Here is the conversation history:** {history}

        NOTE -- "Make sure you need to return only json as response only & please don't add any comments"
        NOTE -- "Make sure you need only need the answer for which context of data is available if not available return empty json as per format"


        ### **Example Output:**
        If the conversation history contains a table and a relevant graph, return:

        ```json
        {{
        "summ": "Based on the conversation, the key points discussed were...",
        "df": [
            {{
            "column1": "value1",
            "column2": "value2",
            "column3": "value3"
            }},
            {{
            "column1": "value1",
            "column2": "value2",
            "column3": "value3"
            }}
        ],
        "graph": "{{\"data\": [{{\"x\": [\"2024-02-01\", \"2024-02-02\"], \"y\": [100, 150], \"type\": \"line\"}}], \"layout\": {{\"title\": \"Sample Graph\"}}"
        }}
        """

        json_structure = """{
            "summ": "Your summary text here...",
                "df": [
                    {
                    "column1": "value1",
                    "column2": "value2",
                    "column3": "value3"
                    },
                    {
                    "column1": "value1",
                    "column2": "value2",
                    "column3": "value3"
                    }
                ],
                "graph": {
                    "data": [ ... ],
                    "layout": { ... }
                }
            }"""

        prompt = PromptTemplate.from_template(template)

        print(prompt)

        llm_chain = prompt | llm
        answer = llm_chain.invoke({"query": request.question, "history": conversation})

        response = (
            answer.content.replace("```json", "")
            .replace("```", "")
            .strip()
            .replace("\n", "")
        )
        response_json = json.loads(response)
        print(response_json)

        node_data = {
            "_id": ObjectId(request.node_id),
            "flow_id": ObjectId(request.flow_id),
            "parent_node_ids": request.parent_node_ids,
            "question": request.question,
            "summ": response_json.get("summ", ""),
            "df": validate_dataframe(response_json.get("df", [])),
            "graph": response_json.get("graph", ""),
            "type": "MultipleQA",
            "is_delete": "false",
            "created_at": datetime.datetime.utcnow(),
        }

        node_id_response = node_collection.insert_one(node_data)

        response = MultipleQuestionAnswerQueryResponse(
            data={
                "question": request.question,
                "summ": response_json.get("summ", ""),
                "df": validate_dataframe(response_json.get("df", [])),
                "graph": response_json.get("graph", ""),
                "flow_id": request.flow_id,
                "parent_node_ids": request.parent_node_ids,
                "component_type": "MultipleQA",
            },
            id=request.node_id,
            type="MultipleQA",
            parent_node_ids=request.parent_node_ids,
        )

        return response

    except Exception as e:
        print(f"Error in /multiple-qa-summarize: {e}")
        raise HTTPException(
            status_code=500, detail=f"Error processing the request: {str(e)}"
        )


@app.post("/flow-summarizer", response_model=FlowSummarizeResponse)
def flow_summarizer(request: FlowSummarizeRequest):
    try:
        nodes = node_collection.find({"flow_id": ObjectId(request.flow_id)})
        if not nodes:
            raise HTTPException(
                status_code=404, detail="No nodes found for the given flow_id."
            )

        conversation = []
        print("This are nodes", request.flow_id)
        for node in nodes:
            node_id = node["_id"]
            print(node_id)
            question, answer = fetch_question_answer_from_node_collection(
                node_id, request.flow_id
            )
            if question and answer:
                conversation.append({"role": "user", "content": question})
                conversation.append(
                    {"role": "assistant", "content": f"Answer: \n{answer}"}
                )

        print(conversation)

        template = """You are an AI assistant tasked with generating a JSX element based on multiple conversations between the user and assistant.

        ### **Rules:**
        1. **Return only valid JSX. No explanations, comments, or extra text.**
        2. **Use only `plotly.js` and `ag-grid-community`. No other libraries are allowed.**
        3. **The output must look like a structured financial report, not just 2-3 components.**
        4. **The layout should have clear sections like:**
            - Executive Summary (must be 100-150 words)
            - Key Financial Metrics  (must highlight all the crucial points)
            - Performance Tables  (most important and data crucial)
            - Multiple Charts for Trends
            - Additional Insights
        5. **Ensure proper spacing, professional styling, and structured formatting.**
        6. **Use `ag-grid-community` for multiple tables.**
        7. **Use `plotly.js` for multiple relevant graphs.**
        8. **All sections should be visually distinct but cohesive.**
        9. **Make the design responsive, with professional styling (flexbox, grid, typography).**
        10. **Ensure the final output looks like an actual financial report from an investment firm.**
        11. For AG-Grid React always include the ref, rowClass, rowHeight, rowStyle, headerHeight and domLayout given in the **REFERENCE**.
        12. **FONT COLOR SHOULD BE BLACK FOR p tags and h1 tags**.
        13. **!important NO STYLING SHOULD BE APPLIED TO PLOT**.
        14. **BACKGROUND COLOR  OF MAIN DIV WILL BE WHITE**.
        15. **TRY TO LAYOUT AG-GRID BY GIVING HEIGHT AND WIDTH AS INLINE CSS BASED ON THE ROWDATA AND COLDEFS**.
        16. **!important GET AG-GRID TABLES IN WHITE **.
        17. **INCLUDE HR TAG after EACH SECTION**.
        18. **ALWAYS HAVE THE HEADING OF FINANCIAL REPORT AT BEGINNING AT THE CENTER**.
        19. **GIVE SOME GAPS AFTER EACH SECTION AND BETWEEN AG-GRID TABLES**.


        ### **Conversation History**
        Here is the conversation history :- {conversation}

        ### **Reference Output (Only JSX, No Comments or Extra Text, No Need to follow as given below):**

        <div>
            <p style={{ fontSize: "16px", fontWeight: "bold", marginBottom: "10px" }}>
                Summary: {{Your Answer}}
            </p>

            <div className="ag-theme-alpine">
                <AgGridReact
                    rowData=[
                        {{ column1: "value1", column2: "value2", column3: "value3" }},
                        {{ column1: "value1", column2: "value2", column3: "value3" }}
                    ]
                    columnDefs=[
                        {{ headerName: "Column 1", field: "column1" }},
                        {{ headerName: "Column 2", field: "column2" }},
                        {{ headerName: "Column 3", field: "column3" }}
                    ]
                    rowClass={{"ag-row"}}
					rowHeight={{56}}
					rowStyle={{ alignItems: "center !important" }}
					headerHeight={{56}}
                    domLayout="autoHeight"
                />
            </div>

            <div style={{ width: "100%", height: "400px" }}>
                <Plot
                    data=[
                        {{ x: [1, 2, 3], y: [10, 20, 30], type: "scatter", mode: "lines+markers", marker: {{ color: "red" }}}}
                    ]
                    layout={{ title: "Graph Title", width: 600, height: 400 }}
                />
            </div>
        </div>
        """
        prompt = PromptTemplate.from_template(template)

        llm_chain = prompt | llm
        answer = llm_chain.invoke({"conversation": conversation})

        print(answer)

        answer = answer.content.replace("```jsx", "").replace("```", "").strip()
        response = answer.replace("\n", "")

        update_result = flow_collection.update_one(
            {"_id": ObjectId(request.flow_id)}, {"$set": {"summary": response}}
        )

        response = FlowSummarizeResponse(
            flow_id=request.flow_id,
            response=response,
        )
        print(response)

        return response

    except Exception as e:
        print(f"Error in /flow-summarizer: {e}")
        traceback.print_exc()
        raise HTTPException(
            status_code=500, detail=f"Error processing the request: {str(e)}"
        )

@app.get("/sqlite-tables", response_model=List[str])
def read_sqlite_tables():
    cursor = connection.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
    tables = [row[0] for row in cursor.fetchall()]
    connection.close()
    return tables
